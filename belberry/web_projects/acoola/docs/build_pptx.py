"""Build Acoola.Team marketing kit v2 as 22-slide PowerPoint draft.

Source brief: docs/acoola_marketing_kit_brief_v2.md
Brand palette: #1F1F22 dark / #EBEBEE light / #2B7CE9 blue / #7A7A85 gray.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Constants ----------

DARK = RGBColor(0x1F, 0x1F, 0x22)
LIGHT = RGBColor(0xEB, 0xEB, 0xEE)
BLUE = RGBColor(0x2B, 0x7C, 0xE9)
GRAY = RGBColor(0x7A, 0x7A, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
PURPLE = RGBColor(0x6B, 0x4F, 0xE9)  # GEO accent
RED = RGBColor(0xC8, 0x3E, 0x3E)
GOLD = RGBColor(0xC9, 0x9A, 0x3D)
DARK_BLUE = RGBColor(0x1B, 0x3F, 0x8E)

HEADING_FONT = "Manrope"  # falls back to Helvetica/Arial if absent
BODY_FONT = "Inter"

SLIDE_W = Emu(12192000)  # 16:9 1920px @ 914400 EMU/inch (13.33in × 7.5in)
SLIDE_H = Emu(6858000)

PAGE_PAD = Emu(560000)  # ~0.61"
HEADER_TOP = Emu(280000)


def emu_in(inches):
    return Emu(int(inches * 914400))


# ---------- Low-level helpers ----------

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None, corner_radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    if corner_radius is not None:
        s.adjustments[0] = corner_radius
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=16, color=WHITE, font=BODY_FONT,
             bold=False, align="left", anchor="top"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_header(slide, dark):
    """ACOOLA TEAM header on every slide."""
    text_color = WHITE if dark else BLACK
    sub_color = GRAY
    add_text(slide, PAGE_PAD, HEADER_TOP, emu_in(3), Emu(220000),
             "WWW.ACOOLA.TEAM", size=10, color=sub_color, font=BODY_FONT)
    # ACOOLA TEAM right-top, two-line composition
    logo_w = emu_in(2.2)
    logo_x = SLIDE_W - PAGE_PAD - logo_w
    add_text(slide, logo_x, HEADER_TOP, logo_w, Emu(280000),
             "ACOOLA", size=18, color=text_color, font=HEADING_FONT,
             bold=True, align="right")
    add_text(slide, logo_x, HEADER_TOP + Emu(220000), logo_w, Emu(180000),
             "TEAM", size=10, color=sub_color, font=HEADING_FONT,
             align="right")


def add_pill(slide, x, y, label, *, w=None, h=None, fill=WHITE, text_color=BLACK):
    if w is None:
        w = emu_in(2.0)
    if h is None:
        h = emu_in(0.45)
    s = add_rect(slide, x, y, w, h, fill, corner_radius=0.5)
    add_text(slide, x, y, w, h, label, size=12, color=text_color,
             font=BODY_FONT, align="center", anchor="middle", bold=True)
    return s


def add_pill_row(slide, x, y, pills, *, fill=WHITE, text_color=BLACK,
                 gap=Emu(120000), pill_w=None):
    if pill_w is None:
        pill_w = emu_in(1.95)
    cur = x
    for label in pills:
        add_pill(slide, cur, y, label, w=pill_w, fill=fill, text_color=text_color)
        cur += pill_w + gap


def add_title(slide, text, dark, *, top=None, size=44, x=None, w=None):
    color = WHITE if dark else BLACK
    if top is None:
        top = emu_in(1.05)
    if x is None:
        x = PAGE_PAD
    if w is None:
        w = SLIDE_W - 2 * PAGE_PAD
    add_text(slide, x, top, w, emu_in(1.0), text,
             size=size, color=color, font=HEADING_FONT, bold=True)


def add_subtitle(slide, text, dark, *, top, x=None, w=None, size=18):
    color = GRAY if dark else GRAY
    if x is None:
        x = PAGE_PAD
    if w is None:
        w = SLIDE_W - 2 * PAGE_PAD
    add_text(slide, x, top, w, emu_in(0.7), text,
             size=size, color=color, font=BODY_FONT)


def add_metric(slide, x, y, w, number, caption, *, dark=True, big=72,
               number_color=None):
    if number_color is None:
        number_color = BLUE
    add_text(slide, x, y, w, emu_in(1.4), number,
             size=big, color=number_color, font=HEADING_FONT, bold=True)
    cap_color = GRAY if dark else GRAY
    add_text(slide, x, y + emu_in(1.3), w, emu_in(0.5), caption,
             size=14, color=cap_color, font=BODY_FONT)


def add_bullets(slide, x, y, w, h, items, *, dark=True, size=14, leading=1.35):
    color = WHITE if dark else BLACK
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = "•  " + item
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, fill, *, title="", subtitle="",
             title_color=None, subtitle_color=None, title_size=20,
             subtitle_size=12, border=None):
    add_rect(slide, x, y, w, h, fill, line=border, corner_radius=0.06)
    if title_color is None:
        # auto
        title_color = WHITE if (fill.rgb if hasattr(fill, "rgb") else fill) == DARK else BLACK
    if subtitle_color is None:
        subtitle_color = GRAY
    pad = emu_in(0.18)
    if title:
        add_text(slide, x + pad, y + pad, w - 2 * pad, emu_in(0.6),
                 title, size=title_size, color=title_color, font=HEADING_FONT,
                 bold=True)
    if subtitle:
        add_text(slide, x + pad, y + h - emu_in(0.7), w - 2 * pad, emu_in(0.55),
                 subtitle, size=subtitle_size, color=subtitle_color, font=BODY_FONT)


# ---------- Slide builders ----------

def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    # Center logo block
    add_text(slide, PAGE_PAD, emu_in(2.2), SLIDE_W - 2 * PAGE_PAD, emu_in(1.4),
             "ACOOLA.TEAM", size=96, color=WHITE, font=HEADING_FONT,
             bold=True, align="center")
    add_text(slide, PAGE_PAD, emu_in(3.5), SLIDE_W - 2 * PAGE_PAD, emu_in(0.6),
             "Digital-агентство полного цикла", size=22, color=GRAY,
             font=BODY_FONT, align="center")
    # Slogan anchor
    add_text(slide, PAGE_PAD, emu_in(4.4), SLIDE_W - 2 * PAGE_PAD, emu_in(0.6),
             "Сайт как актив, а не разовый проект.", size=18,
             color=GRAY, font=BODY_FONT, align="center")
    # Pill row
    pills = ["🔥 Прокачаем", "📈 Выведем в топ", "💼 Увеличим прибыль"]
    pill_w = emu_in(2.05)
    total_w = pill_w * 3 + Emu(120000) * 2
    start_x = (SLIDE_W - total_w) // 2
    add_pill_row(slide, start_x, emu_in(5.3), pills, pill_w=pill_w)

    add_text(slide, PAGE_PAD, SLIDE_H - emu_in(0.5),
             emu_in(3), emu_in(0.3), "acoola.team", size=10,
             color=GRAY, font=BODY_FONT)
    add_text(slide, SLIDE_W - emu_in(3.3), SLIDE_H - emu_in(0.5),
             emu_in(3), emu_in(0.3), "Маркетинг-кит / 2026", size=10,
             color=GRAY, font=BODY_FONT, align="right")


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)
    add_title(slide, "Немного о компании", dark=False, top=emu_in(0.85), size=42)

    # Pill row right
    pills = ["🔥 Прокачаем", "📈 Выведем в топ", "💼 Увеличим прибыль"]
    add_pill_row(slide, emu_in(7.3), emu_in(1.0), pills,
                 pill_w=emu_in(1.7), gap=Emu(80000))

    # 5 award cards in a row
    cards = [
        (DARK, "Рейтинг Рунета\n1 место", "Веб-студии для медицинской ниши", WHITE),
        (DARK_BLUE, "Серебро Tagline\nSMM 2024", "SMM для медицинской тематики", WHITE),
        (RED, "Мониторинг\n1С Битрикс", "Участник программы качества", WHITE),
        (GOLD, "Золотой партнёр\n1С Битрикс", "Автоматизируем бизнес на максимум", WHITE),
        (DARK, "ТОП-3 лучших\nагентств", "3 место Ruward 2025: разработка сайта", WHITE),
    ]
    card_w = emu_in(2.35)
    card_h = emu_in(2.0)
    gap = Emu(120000)
    total_w = card_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) // 2
    y_cards = emu_in(2.0)
    for i, (fill, title, sub, tcolor) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, y_cards, card_w, card_h, fill, corner_radius=0.08)
        add_text(slide, x + emu_in(0.18), y_cards + emu_in(0.25),
                 card_w - emu_in(0.36), emu_in(1.1), title,
                 size=18, color=tcolor, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.18), y_cards + card_h - emu_in(0.85),
                 card_w - emu_in(0.36), emu_in(0.8), sub,
                 size=11, color=RGBColor(0xCC, 0xCC, 0xD0), font=BODY_FONT)

    # Numbers card at bottom (white on light bg)
    nums_y = emu_in(4.4)
    nums_h = emu_in(2.3)
    nums_w = SLIDE_W - 2 * PAGE_PAD
    add_rect(slide, PAGE_PAD, nums_y, nums_w, nums_h, WHITE, corner_radius=0.04)
    add_text(slide, PAGE_PAD + emu_in(0.4), nums_y + emu_in(0.25),
             nums_w - emu_in(0.8), emu_in(0.4),
             "Цифры, которые говорят сами за себя",
             size=14, color=GRAY, font=BODY_FONT)
    metrics = [("50+", "Digital-экспертов"),
               ("70+", "Проектов ежемесячно"),
               (">3,5 ЛЕТ", "Средний срок сотрудничества"),
               ("95%", "Довольных клиентов")]
    col_w = (nums_w - emu_in(0.8)) // 4
    for i, (num, cap) in enumerate(metrics):
        x = PAGE_PAD + emu_in(0.4) + i * col_w
        add_text(slide, x, nums_y + emu_in(0.7), col_w, emu_in(1.15),
                 num, size=58, color=BLUE, font=HEADING_FONT, bold=True)
        add_text(slide, x, nums_y + nums_h - emu_in(0.55), col_w, emu_in(0.4),
                 cap, size=12, color=GRAY, font=BODY_FONT)


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)
    add_title(slide, "Кому мы помогаем", dark=True, top=emu_in(0.95), size=42)
    add_subtitle(slide, "Финансы, недвижимость, промышленность, e-commerce, IT — отрасли, где сайт = канал продаж",
                 dark=True, top=emu_in(1.7))

    rows = [
        ("Финансы", "БыстроБанк · Трансстройбанк · NZT Rusfond · PROFIT CONF · НПФ «Будущее»"),
        ("Недвижимость", "Главстрой · DOGMA · Запстрой · Pioneer · Prime Park"),
        ("Промышленность", "HOHMAN · Robomaster · РУСЭЛ · Yarus Networks · Эксара"),
        ("E-commerce", "J.Prettier · Титан · MARKIN · Coolstream · Latrika"),
        ("IT", "DIS Group · NGENIX · textback · SBER API · SBER DEVICES"),
    ]
    y0 = emu_in(2.6)
    row_h = emu_in(0.75)
    for i, (industry, clients) in enumerate(rows):
        y = y0 + i * row_h
        add_text(slide, PAGE_PAD, y, emu_in(2.5), row_h, industry,
                 size=18, color=WHITE, font=HEADING_FONT, bold=True,
                 anchor="middle")
        add_text(slide, PAGE_PAD + emu_in(2.7), y,
                 SLIDE_W - PAGE_PAD - emu_in(2.7) - PAGE_PAD, row_h,
                 clients, size=14, color=RGBColor(0xCC, 0xCC, 0xD0),
                 font=BODY_FONT, anchor="middle")
        # divider
        add_rect(slide, PAGE_PAD, y + row_h - Emu(8000),
                 SLIDE_W - 2 * PAGE_PAD, Emu(8000),
                 RGBColor(0x33, 0x33, 0x36))

    # Bottom: jump to slide 19
    add_text(slide, PAGE_PAD, SLIDE_H - emu_in(0.6),
             SLIDE_W - 2 * PAGE_PAD, emu_in(0.4),
             "5 отраслей × 5 клиентов = 25 Tier-1 имён  →  слайд 19",
             size=12, color=BLUE, font=BODY_FONT, bold=True)


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)
    add_title(slide, "Какие задачи решаем", dark=False, top=emu_in(0.95), size=42)
    add_subtitle(slide, "Переводим услуги на язык бизнес-результатов",
                 dark=False, top=emu_in(1.7))

    tasks = [
        ("Превратить сайт в канал продаж", "Разработка и веб-дизайн под бизнес-метрики"),
        ("Снизить CPL и поднять ROI", "Контекст и SEO в связке с аналитикой"),
        ("Стабильный органический трафик", "SEO с фокусом на обращения, а не позиции"),
        ("Запустить CRM и навести порядок", "Битрикс24 как Золотой партнёр"),
        ("Защитить и усилить репутацию", "ORM на 30+ площадках"),
        ("Появиться в ответах AI-ассистентов", "Новое направление GEO 2026"),
    ]
    cols = 3
    rows = 2
    grid_x = PAGE_PAD
    grid_y = emu_in(2.5)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    grid_h = SLIDE_H - grid_y - PAGE_PAD
    gap = Emu(160000)
    cell_w = (grid_w - gap * (cols - 1)) // cols
    cell_h = (grid_h - gap * (rows - 1)) // rows
    for i, (title, subtitle) in enumerate(tasks):
        r, c = divmod(i, cols)
        x = grid_x + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        is_geo = i == 5
        fill = WHITE if not is_geo else RGBColor(0xEC, 0xE6, 0xFE)
        border = None if not is_geo else PURPLE
        add_rect(slide, x, y, cell_w, cell_h, fill,
                 line=border, corner_radius=0.05)
        add_text(slide, x + emu_in(0.3), y + emu_in(0.3),
                 cell_w - emu_in(0.6), emu_in(0.7), title,
                 size=18, color=BLACK, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.3), y + emu_in(1.0),
                 cell_w - emu_in(0.6), cell_h - emu_in(1.3), subtitle,
                 size=13, color=GRAY, font=BODY_FONT)
        if is_geo:
            add_pill(slide, x + cell_w - emu_in(1.4), y + emu_in(0.25),
                     "New 2026", w=emu_in(1.1), h=emu_in(0.35),
                     fill=PURPLE, text_color=WHITE)


def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)
    add_title(slide, "9 направлений под одной крышей", dark=True,
              top=emu_in(0.95), size=42)
    add_subtitle(slide,
                 "Полный digital-цикл: от первой строчки кода до первой строчки в ответе ChatGPT",
                 dark=True, top=emu_in(1.7))

    services = [
        ("Разработка сайта", "Под бизнес-задачу"),
        ("Шаблон Битрикс", "Быстрый запуск"),
        ("Веб-дизайн", "Дизайн под метрики"),
        ("SEO-продвижение", "Заявки и выручка"),
        ("Контекст", "Не клики — деньги"),
        ("ORM", "Репутация как процесс"),
        ("Битрикс24", "CRM, в которой работают"),
        ("Техподдержка", "Сайт всегда на связи"),
        ("GEO", "Видимость в AI-поиске"),
    ]
    grid_x = PAGE_PAD
    grid_y = emu_in(2.4)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    grid_h = SLIDE_H - grid_y - PAGE_PAD
    cols = 3
    rows = 3
    gap = Emu(120000)
    cell_w = (grid_w - gap * (cols - 1)) // cols
    cell_h = (grid_h - gap * (rows - 1)) // rows
    for i, (name, sub) in enumerate(services):
        r, c = divmod(i, cols)
        x = grid_x + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        is_geo = i == 8
        fill = RGBColor(0x2A, 0x2A, 0x2E) if not is_geo else PURPLE
        add_rect(slide, x, y, cell_w, cell_h, fill,
                 line=RGBColor(0x44, 0x44, 0x48), corner_radius=0.06)
        add_text(slide, x + emu_in(0.3), y + emu_in(0.3),
                 cell_w - emu_in(0.6), emu_in(0.6), name,
                 size=20, color=WHITE, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.3), y + emu_in(1.0),
                 cell_w - emu_in(0.6), emu_in(0.5), sub,
                 size=12, color=RGBColor(0xCC, 0xCC, 0xD0), font=BODY_FONT)
        # arrow
        add_text(slide, x + cell_w - emu_in(0.5), y + emu_in(0.3),
                 emu_in(0.4), emu_in(0.4), "↗",
                 size=16, color=GRAY, font=BODY_FONT, align="right")
        if is_geo:
            add_pill(slide, x + emu_in(0.3), y + cell_h - emu_in(0.7),
                     "New 2026", w=emu_in(1.1), h=emu_in(0.35),
                     fill=WHITE, text_color=PURPLE)


def _service_slide(prs, *, dark, title, subtitle, pill, pain, do, includes,
                   proof, fits, price_text, cta, side_placeholder):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK if dark else LIGHT)
    add_header(slide, dark=dark)
    add_title(slide, title, dark=dark, top=emu_in(0.95), size=40)
    add_subtitle(slide, subtitle, dark=dark, top=emu_in(1.7))

    # Pill (top right under header)
    if pill:
        add_pill(slide, SLIDE_W - PAGE_PAD - emu_in(2.0), emu_in(1.0),
                 pill, w=emu_in(1.95),
                 fill=WHITE if not dark else RGBColor(0x33, 0x33, 0x36),
                 text_color=BLACK if not dark else WHITE)

    # Left 60% — bullets
    left_x = PAGE_PAD
    left_y = emu_in(2.5)
    left_w = (SLIDE_W - 2 * PAGE_PAD) * 6 // 10 - Emu(120000)
    items = [
        f"Боль клиента — {pain}",
        f"Что мы делаем — {do}",
        f"Что входит — {includes}",
        f"Доказательство — {proof}",
        f"Кому подходит — {fits}",
    ]
    add_bullets(slide, left_x, left_y, left_w, emu_in(3.5), items,
                dark=dark, size=13)

    # Right 40% — placeholder graphic
    right_x = left_x + left_w + Emu(120000)
    right_w = SLIDE_W - PAGE_PAD - right_x
    right_y = left_y
    right_h = emu_in(3.5)
    placeholder_fill = RGBColor(0x33, 0x33, 0x36) if dark else WHITE
    add_rect(slide, right_x, right_y, right_w, right_h,
             placeholder_fill, corner_radius=0.04,
             line=RGBColor(0x55, 0x55, 0x58) if dark else RGBColor(0xD0, 0xD0, 0xD3))
    add_text(slide, right_x, right_y, right_w, right_h,
             side_placeholder, size=14, color=GRAY,
             font=BODY_FONT, align="center", anchor="middle")

    # Bottom strip — price + CTA
    strip_y = SLIDE_H - emu_in(1.0)
    strip_h = emu_in(0.7)
    add_text(slide, PAGE_PAD, strip_y, emu_in(4), strip_h, price_text,
             size=22, color=BLUE, font=HEADING_FONT, bold=True, anchor="middle")
    add_pill(slide, SLIDE_W - PAGE_PAD - emu_in(3.5), strip_y, cta,
             w=emu_in(3.5), h=strip_h, fill=BLUE, text_color=WHITE)


def slide_06(prs):
    _service_slide(
        prs, dark=False,
        title="Разработка сайта",
        subtitle="Сайт сразу как актив для продвижения, а не как красивая обёртка",
        pill="🔥 Прокачаем",
        pain="устаревший сайт не приносит заявок и тяжело дорабатывается",
        do="разработка идёт в связке с SEO, дизайном и аналитикой",
        includes="UX, дизайн в системе, разработка на 1С Битрикс, интеграции, тестирование, передача под SEO",
        proof="Золотой партнёр 1С Битрикс · 95% довольных клиентов · Сбер, БыстроБанк, Главстрой",
        fits="B2B/B2C-проекты со сложной структурой каталога и амбициями на трафик",
        price_text="от __ ₽  [TODO]",
        cta="Запросить расчёт разработки",
        side_placeholder="[ МАКЕТ САЙТА / WIREFRAME ]\nдизайнер: подставить мокап",
    )


def slide_07(prs):
    _service_slide(
        prs, dark=True,
        title="Разработка на шаблоне Битрикс",
        subtitle="Быстрый запуск без технического долга, с возможностью масштабирования",
        pill="🔥 Прокачаем",
        pain="нужен сайт быстро, но кустарные шаблоны тянут технический долг",
        do="разрабатываем на шаблонах 1С Битрикс с учётом будущего масштабирования",
        includes="подбор и адаптация шаблона, базовая SEO-настройка, интеграции, тестирование, инструкция для редактора",
        proof="Золотой партнёр 1С Битрикс — Автоматизируем бизнес на максимум",
        fits="малый и средний бизнес с типовой моделью, кому сайт нужен за недели",
        price_text="от __ ₽  [TODO]",
        cta="Запросить расчёт шаблона",
        side_placeholder="[ МОДУЛЬНАЯ СЕТКА ШАБЛОНА ]\nдизайнер: подставить графику",
    )


def slide_08(prs):
    _service_slide(
        prs, dark=False,
        title="Веб-дизайн",
        subtitle="Дизайн, который окупается метриками",
        pill="📈 Выведем в топ",
        pain="сайт «красивый», но не конвертит — дизайн делали под вкус",
        do="проектируем UX по сценариям, разрабатываем визуальную систему",
        includes="UX-исследование, вайрфреймы, дизайн-концепция, ключевые экраны, дизайн-система компонентов",
        proof="Серебро Tagline SMM 2024 — медицина · 3 место Ruward 2025 — разработка",
        fits="компании, готовые менять «нравится» на «работает»; премиум-сегменты",
        price_text="от __ ₽  [TODO]",
        cta="Запросить дизайн-консультацию",
        side_placeholder="[ ОБРАЗЕЦ ДИЗАЙН-СИСТЕМЫ ]\nдизайнер: подставить экран сайта",
    )


def slide_09(prs):
    _service_slide(
        prs, dark=True,
        title="SEO-продвижение",
        subtitle="Не топ-10. Заявки и выручка",
        pill="📈 Выведем в топ",
        pain="трафик не растёт, лиды дорогие, прошлый подрядчик «накрутил позиции»",
        do="ведём SEO как процесс — аудит, семантика, контент, ссылки, отчётность по выручке",
        includes="технический аудит, семантическое ядро, контент-план, on-page, линкбилдинг, ежемесячная отчётность",
        proof="1 место Рейтинг Рунета — медицина · Гольфстрим: лид 265→50 ₽ · wwerb: +223% органики",
        fits="e-commerce с большой матрицей; B2B с длинной воронкой; премиальные ниши",
        price_text="от __ ₽/мес  [TODO]",
        cta="Запросить SEO-аудит",
        side_placeholder="[ ГРАФИК РОСТА ОБРАЩЕНИЙ ]\nдизайнер: подставить кейс",
    )


def slide_10(prs):
    _service_slide(
        prs, dark=False,
        title="Контекстная реклама",
        subtitle="Не клики. Деньги",
        pill="💼 Увеличим прибыль",
        pain="реклама съедает бюджет, заявки не превращаются в продажи — ROI отрицательный",
        do="связка контекста с аналитикой и сайтом — оптимизируем не CPL, а ROI",
        includes="медиапланирование, настройка Яндекс.Директ и VK Ads, корректировка ставок, A/B тесты, отчёт по выручке",
        proof="Kenner: ROI 140% · HOHMAN: 5,6 млн ₽ выручки за 3 месяца · Сертифицированное агентство Яндекса",
        fits="интернет-магазины с конверсионной воронкой; производители; B2B-сегменты",
        price_text="от __ ₽/мес + бюджет  [TODO]",
        cta="Запросить расчёт контекста",
        side_placeholder="[ КАРТОЧКА «ROI 140%» ]\nдизайнер: крупная цифра",
    )


def slide_11(prs):
    _service_slide(
        prs, dark=True,
        title="Управление репутацией (ORM)",
        subtitle="Репутация как процесс, а не как тушение пожара",
        pill="💼 Увеличим прибыль",
        pain="старые негативные отзывы в выдаче, конкуренты вытесняют из карточек",
        do="строим устойчивую систему сбора и управления обратной связью",
        includes="аудит, регламент работы с негативом, оптимизация карточек 2ГИС/Яндекс.Карт/Banki.ru, мониторинг, отчёт",
        proof="Серебро Tagline SMM 2024 — медицина · средний срок сотрудничества 3,5+ года",
        fits="компании с физическими точками; e-commerce с потоком заказов; премиум-бренды",
        price_text="от __ ₽/мес  [TODO]",
        cta="Запросить ORM-аудит",
        side_placeholder="[ КАРТА 30+ ПЛОЩАДОК ]\nдизайнер: схема присутствия",
    )


def slide_12(prs):
    _service_slide(
        prs, dark=False,
        title="Внедрение Битрикс24",
        subtitle="CRM, в которой работают, а не открывают раз в неделю",
        pill="💼 Увеличим прибыль",
        pain="менеджеры теряют заявки, у руководителя нет картины по воронке",
        do="не просто «ставим коробку», а собираем процесс: воронка, права, отчёты, интеграции",
        includes="аудит процессов, проектирование воронки и прав, настройка, интеграции с сайтом и телефонией, обучение",
        proof="Золотой партнёр 1С Битрикс · Участник программы мониторинга качества",
        fits="компании 5–50 человек, выросшие из таблиц; смешанные каналы заявок",
        price_text="от __ ₽  [TODO]",
        cta="Запросить расчёт Битрикс24",
        side_placeholder="[ СХЕМА ВОРОНКИ CRM ]\nдизайнер: 5–6 этапов",
    )


def slide_13(prs):
    _service_slide(
        prs, dark=True,
        title="Техническая поддержка",
        subtitle="Сайт всегда на связи",
        pill="",
        pain="на сайте «упало», менеджер у подрядчика в отпуске, форма заявки не работает",
        do="SLA на ответ в часах, регламент работ, плановые задачи и резерв под форс-мажор",
        includes="регулярные доработки и фиксы, мониторинг, резервное копирование, обновления, ежемесячный отчёт по часам",
        proof="средний срок сотрудничества 3,5+ года · 95% довольных клиентов",
        fits="все клиенты с действующим сайтом, для кого простой = упущенные заявки",
        price_text="от __ ₽/мес  [TODO]",
        cta="Запросить условия техподдержки",
        side_placeholder="[ КАРТОЧКИ «3,5+ ГОДА / 95%» ]\nдизайнер: пара метрик",
    )


def slide_14(prs):
    """GEO — special launch slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    # Pill row
    add_pill(slide, SLIDE_W - PAGE_PAD - emu_in(4.1), emu_in(1.0),
             "🔥 Прокачаем", w=emu_in(1.95), fill=RGBColor(0x33, 0x33, 0x36),
             text_color=WHITE)
    add_pill(slide, SLIDE_W - PAGE_PAD - emu_in(2.0), emu_in(1.0),
             "📈 Выведем в топ", w=emu_in(1.95), fill=RGBColor(0x33, 0x33, 0x36),
             text_color=WHITE)

    add_title(slide, "GEO — видимость бренда в AI-поиске",
              dark=True, top=emu_in(0.95), size=40)
    add_subtitle(slide,
                 "Делаем так, чтобы ChatGPT, Perplexity, Gemini, YandexGPT и Алиса упоминали ваш бренд — а не только конкурентов",
                 dark=True, top=emu_in(1.7))

    # Left 50% — AI mockup placeholder
    left_x = PAGE_PAD
    left_y = emu_in(2.6)
    left_w = (SLIDE_W - 2 * PAGE_PAD) // 2 - Emu(120000)
    left_h = emu_in(3.0)
    add_rect(slide, left_x, left_y, left_w, left_h,
             RGBColor(0x33, 0x33, 0x36), corner_radius=0.04)
    add_text(slide, left_x, left_y, left_w, emu_in(0.6),
             "ChatGPT · Perplexity · Gemini · YandexGPT · Алиса",
             size=11, color=GRAY, font=BODY_FONT, align="center", anchor="middle")
    add_text(slide, left_x, left_y, left_w, left_h,
             "[ МОКАП ОТВЕТА AI-АССИСТЕНТА ]\nдизайнер: карточка с цитатой бренда\nи логотипами LLM в шапке",
             size=14, color=GRAY, font=BODY_FONT, align="center", anchor="middle")

    # Right 50% — 5 steps
    right_x = left_x + left_w + Emu(240000)
    right_w = SLIDE_W - PAGE_PAD - right_x
    steps = ["1. Аудит видимости в LLM",
             "2. Карта AI-запросов аудитории",
             "3. Источники цитирования: контент, медиа, отзывы",
             "4. Структура данных: schema.org, llms.txt, FAQ",
             "5. Ежемесячный отчёт: доля голоса, динамика"]
    step_h = left_h // 5
    for i, step in enumerate(steps):
        y = left_y + i * step_h
        add_text(slide, right_x, y, right_w, step_h,
                 step, size=14, color=WHITE, font=BODY_FONT, anchor="middle")

    # Bottom strip — price + CTA + badge
    strip_y = SLIDE_H - emu_in(1.0)
    strip_h = emu_in(0.7)
    add_rect(slide, PAGE_PAD, strip_y, emu_in(4.5), strip_h, PURPLE,
             corner_radius=0.1)
    add_text(slide, PAGE_PAD, strip_y, emu_in(4.5), strip_h,
             "от 90 000 ₽/мес", size=22, color=WHITE,
             font=HEADING_FONT, bold=True, align="center", anchor="middle")
    add_pill(slide, PAGE_PAD + emu_in(4.7), strip_y, "Новое направление 2026",
             w=emu_in(3.0), h=strip_h, fill=WHITE, text_color=PURPLE)
    add_pill(slide, SLIDE_W - PAGE_PAD - emu_in(3.0), strip_y,
             "Запросить GEO-аудит",
             w=emu_in(3.0), h=strip_h, fill=BLUE, text_color=WHITE)


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)
    add_title(slide, "Как мы ведём проект", dark=False, top=emu_in(0.95), size=42)
    add_subtitle(slide, "6 шагов от первого письма до ежемесячного отчёта",
                 dark=False, top=emu_in(1.7))

    steps = [
        ("1", "Бриф и аудит", "Бизнес, цели, состояние сайта"),
        ("2", "Стратегия", "План работ на 3–6 месяцев"),
        ("3", "Запуск", "Разработка / SEO / реклама / ORM"),
        ("4", "Итерации", "Задачи → реализация → измерение"),
        ("5", "Отчёт", "Что сделали и измерили"),
        ("6", "Развитие", "Новые гипотезы, длинная дистанция"),
    ]
    grid_y = emu_in(2.6)
    grid_h = emu_in(3.0)
    cols = 6
    gap = Emu(80000)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    cell_w = (grid_w - gap * (cols - 1)) // cols
    for i, (num, title, sub) in enumerate(steps):
        x = PAGE_PAD + i * (cell_w + gap)
        is_hl = i >= 4  # шаги 5 и 6 — выделены
        fill = WHITE if not is_hl else RGBColor(0xDC, 0xE9, 0xFB)
        border = None if not is_hl else BLUE
        add_rect(slide, x, grid_y, cell_w, grid_h, fill,
                 line=border, corner_radius=0.05)
        add_text(slide, x + emu_in(0.25), grid_y + emu_in(0.25),
                 cell_w - emu_in(0.5), emu_in(0.6), num,
                 size=24, color=BLUE, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.25), grid_y + emu_in(1.0),
                 cell_w - emu_in(0.5), emu_in(0.6), title,
                 size=16, color=BLACK, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.25), grid_y + emu_in(1.6),
                 cell_w - emu_in(0.5), emu_in(1.2), sub,
                 size=11, color=GRAY, font=BODY_FONT)

    add_text(slide, PAGE_PAD, SLIDE_H - emu_in(0.7),
             SLIDE_W - 2 * PAGE_PAD, emu_in(0.5),
             "Средний срок сотрудничества — 3,5+ года",
             size=14, color=GRAY, font=BODY_FONT, align="center")


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)
    add_title(slide, "Почему Acoola.Team", dark=True, top=emu_in(0.95), size=42)
    add_subtitle(slide, "4 опоры, на которых держатся длинные отношения с клиентами",
                 dark=True, top=emu_in(1.7))

    cards = [
        ("Полный цикл", "9 направлений работ — снимаем «зоопарк подрядчиков»", None),
        ("Tier-1 клиенты", "Сбер · БыстроБанк · Главстрой · Pioneer · NGENIX", None),
        ("Длинные отношения", "Средний срок сотрудничества", "3,5+ года"),
        ("95% довольных", "Главный фирменный показатель", "95%"),
    ]
    grid_y = emu_in(2.6)
    grid_h = emu_in(3.4)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    cols = 2
    rows = 2
    gap = Emu(160000)
    cell_w = (grid_w - gap) // cols
    cell_h = (grid_h - gap) // rows
    for i, (title, sub, metric) in enumerate(cards):
        r, c = divmod(i, cols)
        x = PAGE_PAD + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        is_hero = i == 3
        fill = RGBColor(0x33, 0x33, 0x36) if not is_hero else BLUE
        add_rect(slide, x, y, cell_w, cell_h, fill, corner_radius=0.05)
        add_text(slide, x + emu_in(0.4), y + emu_in(0.4),
                 cell_w - emu_in(0.8), emu_in(0.7), title,
                 size=20, color=WHITE, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.4), y + emu_in(1.2),
                 cell_w - emu_in(0.8), emu_in(0.9), sub,
                 size=13, color=RGBColor(0xCC, 0xCC, 0xD0), font=BODY_FONT)
        if metric:
            add_text(slide, x + emu_in(0.4), y + cell_h - emu_in(1.4),
                     cell_w - emu_in(0.8), emu_in(1.2), metric,
                     size=54, color=WHITE, font=HEADING_FONT, bold=True)


def slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)
    add_title(slide, "Кейсы: контекст и SEO", dark=True, top=emu_in(0.95), size=40)
    add_subtitle(slide, "4 проекта, где считаем рост в деньгах, а не в позициях",
                 dark=True, top=emu_in(1.7))

    cases = [
        ("Kenner", "Мебель · контекст",
         "ROI 140%\nCPL 1 374 ₽ (корзина)\nCPL 4 067 ₽ (продажа)"),
        ("HOHMAN", "Станки · контекст",
         "5,6 млн ₽ за 3 мес\nROI 140 · 68% из РСЯ\nДо 15 регионов"),
        ("Гольфстрим", "Охрана · SEO",
         "Лид: 265 → 50 ₽\nКонверсии: 602 → 3 213\n+40% видимости · 3 года"),
        ("wwerb.com", "Часы · SEO",
         "Лид: 18 400 → 582 ₽\nОрганика: +223%\nВыручка 164+ млн ₽"),
    ]
    grid_y = emu_in(2.6)
    grid_h = emu_in(3.6)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    cols = 2
    rows = 2
    gap = Emu(160000)
    cell_w = (grid_w - gap) // cols
    cell_h = (grid_h - gap) // rows
    for i, (client, tag, metrics) in enumerate(cases):
        r, c = divmod(i, cols)
        x = PAGE_PAD + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        add_rect(slide, x, y, cell_w, cell_h, RGBColor(0x2A, 0x2A, 0x2E),
                 line=RGBColor(0x44, 0x44, 0x48), corner_radius=0.05)
        add_text(slide, x + emu_in(0.4), y + emu_in(0.3),
                 cell_w - emu_in(0.8), emu_in(0.5), client,
                 size=22, color=WHITE, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.4), y + emu_in(0.85),
                 cell_w - emu_in(0.8), emu_in(0.4), tag,
                 size=12, color=GRAY, font=BODY_FONT)
        add_text(slide, x + emu_in(0.4), y + emu_in(1.35),
                 cell_w - emu_in(0.8), cell_h - emu_in(1.55), metrics,
                 size=14, color=BLUE, font=HEADING_FONT, bold=True)


def slide_18(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)
    add_title(slide, "Кейсы: недвижимость, e-commerce", dark=True,
              top=emu_in(0.95), size=40)
    add_subtitle(slide, "Где UX, SEO и SMM работают на обращения и выручку",
                 dark=True, top=emu_in(1.7))

    cases = [
        ("ALIA", "Недвижимость · SEO",
         "Лид: 1 780 → 195 ₽\nПосещения: 5 505 → 6 754/мес\nОтказы: 13,5% → 12,7%"),
        ("SKY VIEW", "Премиум · SEO+SMM",
         "Отказы: 24,9% → 11,1%\nВремя: 2:12 → 3:52\nЛид: 3 330 → 1 480 ₽"),
        ("КВАДРОСТИЛЬ", "E-commerce · SEO",
         "×223 переходов из поиска\n×2,3 конверсий\n5 лет работы"),
    ]
    grid_y = emu_in(2.6)
    grid_h = emu_in(3.4)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    cols = 3
    gap = Emu(160000)
    cell_w = (grid_w - gap * (cols - 1)) // cols
    for i, (client, tag, metrics) in enumerate(cases):
        x = PAGE_PAD + i * (cell_w + gap)
        y = grid_y
        add_rect(slide, x, y, cell_w, grid_h, RGBColor(0x2A, 0x2A, 0x2E),
                 line=RGBColor(0x44, 0x44, 0x48), corner_radius=0.05)
        add_text(slide, x + emu_in(0.4), y + emu_in(0.3),
                 cell_w - emu_in(0.8), emu_in(0.5), client,
                 size=22, color=WHITE, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.4), y + emu_in(0.85),
                 cell_w - emu_in(0.8), emu_in(0.4), tag,
                 size=12, color=GRAY, font=BODY_FONT)
        add_text(slide, x + emu_in(0.4), y + emu_in(1.35),
                 cell_w - emu_in(0.8), grid_h - emu_in(1.55), metrics,
                 size=14, color=BLUE, font=HEADING_FONT, bold=True)

    add_text(slide, PAGE_PAD, SLIDE_H - emu_in(0.6),
             SLIDE_W - 2 * PAGE_PAD, emu_in(0.4),
             "Полный список кейсов: acoola.team/projects",
             size=12, color=GRAY, font=BODY_FONT, align="center")


def slide_19(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)
    add_title(slide, "Наши клиенты", dark=True, top=emu_in(0.95), size=42)
    add_subtitle(slide,
                 "70+ проектов в работе ежемесячно. 5 отраслей × 5 Tier-1 клиентов",
                 dark=True, top=emu_in(1.7))

    rows = [
        ("Финансы", ["БыстроБанк", "Трансстройбанк", "NZT Rusfond", "PROFIT CONF", "НПФ «Будущее»"]),
        ("Недвижимость", ["Главстрой", "DOGMA", "Запстрой", "Pioneer", "Prime Park"]),
        ("Промышленность", ["HOHMAN", "Robomaster", "РУСЭЛ", "Yarus Networks", "Эксара"]),
        ("E-commerce", ["J.Prettier", "Титан", "MARKIN", "Coolstream", "Latrika"]),
        ("IT", ["DIS Group", "NGENIX", "textback", "SBER API", "SBER DEVICES"]),
    ]
    grid_y = emu_in(2.6)
    grid_h = SLIDE_H - grid_y - PAGE_PAD - emu_in(0.3)
    grid_w = SLIDE_W - 2 * PAGE_PAD
    label_w = emu_in(2.3)
    cell_w = (grid_w - label_w) // 5
    row_h = grid_h // 5
    for r, (label, clients) in enumerate(rows):
        y = grid_y + r * row_h
        add_text(slide, PAGE_PAD, y, label_w, row_h, label,
                 size=15, color=GRAY, font=HEADING_FONT, bold=True,
                 anchor="middle")
        for c, name in enumerate(clients):
            x = PAGE_PAD + label_w + c * cell_w
            add_rect(slide, x, y + Emu(40000),
                     cell_w - Emu(60000), row_h - Emu(80000),
                     RGBColor(0x2A, 0x2A, 0x2E),
                     line=RGBColor(0x44, 0x44, 0x48), corner_radius=0.04)
            add_text(slide, x + emu_in(0.15), y + Emu(40000),
                     cell_w - Emu(60000) - emu_in(0.3), row_h - Emu(80000),
                     name, size=12, color=WHITE, font=BODY_FONT,
                     anchor="middle")
            # arrow
            add_text(slide, x + cell_w - emu_in(0.5), y + emu_in(0.1),
                     emu_in(0.3), emu_in(0.3), "↗",
                     size=11, color=GRAY, font=BODY_FONT, align="right")


def slide_20(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)
    add_title(slide, "Команда", dark=False, top=emu_in(0.95), size=42)
    add_subtitle(slide,
                 "50+ digital-специалистов в 6 направлениях — все в Москве, в одном офисе",
                 dark=False, top=emu_in(1.7))

    # Left — big metric
    left_x = PAGE_PAD
    left_y = emu_in(2.6)
    left_w = emu_in(4.0)
    left_h = emu_in(3.4)
    add_rect(slide, left_x, left_y, left_w, left_h, WHITE, corner_radius=0.05)
    add_text(slide, left_x, left_y + emu_in(0.6), left_w, emu_in(1.5),
             "50+", size=120, color=BLUE, font=HEADING_FONT, bold=True,
             align="center")
    add_text(slide, left_x, left_y + emu_in(2.1), left_w, emu_in(0.5),
             "специалистов", size=18, color=BLACK, font=BODY_FONT,
             align="center")
    add_text(slide, left_x, left_y + emu_in(2.6), left_w, emu_in(0.4),
             "в 6 направлениях", size=14, color=GRAY, font=BODY_FONT,
             align="center")

    # Right — 4 columns
    right_x = left_x + left_w + emu_in(0.3)
    right_w = SLIDE_W - PAGE_PAD - right_x
    cols = [
        ("Разработка", "Frontend\nBackend\nBitrix Framework\nАдмины"),
        ("Маркетинг и трафик", "SEO\nКонтекст\nSMM\nКонтент"),
        ("Аналитика", "Brand Analytics\nCalltouch\nЯндекс.Метрика\nGoogle Analytics"),
        ("Управление", "Менеджеры проектов\nДиректор по производству\nРуководители отделов"),
    ]
    col_gap = Emu(100000)
    col_w = (right_w - col_gap * 3) // 4
    for i, (h, body) in enumerate(cols):
        x = right_x + i * (col_w + col_gap)
        add_rect(slide, x, left_y, col_w, left_h, WHITE, corner_radius=0.05)
        add_text(slide, x + emu_in(0.2), left_y + emu_in(0.3),
                 col_w - emu_in(0.4), emu_in(0.5), h,
                 size=14, color=BLUE, font=HEADING_FONT, bold=True)
        add_text(slide, x + emu_in(0.2), left_y + emu_in(0.95),
                 col_w - emu_in(0.4), left_h - emu_in(1.1), body,
                 size=12, color=BLACK, font=BODY_FONT)

    # Certifications strip
    add_text(slide, PAGE_PAD, SLIDE_H - emu_in(0.7),
             SLIDE_W - 2 * PAGE_PAD, emu_in(0.5),
             "Сертификации: Яндекс · Google · 1С Битрикс · Brand Analytics · Calltouch",
             size=12, color=GRAY, font=BODY_FONT, align="center")


def slide_21(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)
    add_title(slide, "Внешняя валидация", dark=False, top=emu_in(0.95), size=42)
    add_subtitle(slide, "Точные формулировки из фирменной презентации",
                 dark=False, top=emu_in(1.7))

    # Left — awards grid 3×2
    left_x = PAGE_PAD
    left_y = emu_in(2.6)
    left_w = emu_in(7.5)
    left_h = emu_in(3.5)
    awards = [
        "1 место Рейтинг Рунета —\nВеб-студии для медицинской ниши",
        "Серебро Tagline SMM 2024 —\nSMM для медицинской тематики",
        "3 место Ruward 2025 —\nРазработка сайта",
        "Золотой партнёр 1С Битрикс —\nАвтоматизируем бизнес на максимум",
        "Участник программы\nмониторинга качества 1С Битрикс",
        "Сертифицированное агентство\nЯндекса · Google Partner",
    ]
    cols = 3
    rows = 2
    gap = Emu(80000)
    cell_w = (left_w - gap * (cols - 1)) // cols
    cell_h = (left_h - gap * (rows - 1)) // rows
    for i, txt in enumerate(awards):
        r, c = divmod(i, cols)
        x = left_x + c * (cell_w + gap)
        y = left_y + r * (cell_h + gap)
        add_rect(slide, x, y, cell_w, cell_h, WHITE, corner_radius=0.05)
        add_text(slide, x + emu_in(0.2), y, cell_w - emu_in(0.4), cell_h,
                 txt, size=12, color=BLACK, font=BODY_FONT, anchor="middle")

    # Right — 95% hero
    right_x = left_x + left_w + emu_in(0.3)
    right_w = SLIDE_W - PAGE_PAD - right_x
    add_rect(slide, right_x, left_y, right_w, left_h, BLUE, corner_radius=0.05)
    add_text(slide, right_x, left_y + emu_in(0.4), right_w, emu_in(2.0),
             "95%", size=140, color=WHITE, font=HEADING_FONT, bold=True,
             align="center")
    add_text(slide, right_x, left_y + emu_in(2.4), right_w, emu_in(0.5),
             "довольных клиентов", size=18, color=WHITE, font=BODY_FONT,
             align="center")
    add_text(slide, right_x, left_y + left_h - emu_in(0.7),
             right_w, emu_in(0.5),
             "NPS 100 (год) / 95 (общий период) — 1С Битрикс",
             size=10, color=RGBColor(0xDD, 0xE9, 0xFB),
             font=BODY_FONT, align="center")


def slide_22(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    pills = ["🔥 Прокачаем", "📈 Выведем в топ", "💼 Увеличим прибыль"]
    add_pill_row(slide, emu_in(7.3), emu_in(1.0), pills,
                 pill_w=emu_in(1.7), gap=Emu(80000))

    add_title(slide, "Получить план роста\nдля вашего сайта",
              dark=False, top=emu_in(1.6), size=48)
    add_subtitle(slide,
                 "Аудит за 3 рабочих дня. Считаем экономику, не давим продажей",
                 dark=False, top=emu_in(3.4))

    # Left — CTA big button
    cta_x = PAGE_PAD
    cta_y = emu_in(4.4)
    cta_w = emu_in(5.0)
    cta_h = emu_in(1.2)
    add_rect(slide, cta_x, cta_y, cta_w, cta_h, BLUE, corner_radius=0.5)
    add_text(slide, cta_x, cta_y, cta_w, cta_h,
             "ПОЛУЧИТЬ ПЛАН РОСТА  →",
             size=22, color=WHITE, font=HEADING_FONT, bold=True,
             align="center", anchor="middle")

    # Right — contacts card
    c_x = PAGE_PAD + cta_w + emu_in(0.4)
    c_y = cta_y - emu_in(0.6)
    c_w = SLIDE_W - PAGE_PAD - c_x
    c_h = emu_in(2.4)
    add_rect(slide, c_x, c_y, c_w, c_h, WHITE, corner_radius=0.05)
    contacts = [
        "Email: hello@acoola.team",
        "Телефон: +7 499 938-48-27 · 10:00–19:00 МСК",
        "Адрес: Москва, Варшавское ш. 118к1 «ВаршавскаSky»",
        "Сайт: acoola.team",
    ]
    for i, line in enumerate(contacts):
        add_text(slide, c_x + emu_in(0.3), c_y + emu_in(0.3) + i * emu_in(0.5),
                 c_w - emu_in(0.6), emu_in(0.5), line,
                 size=14, color=BLACK, font=BODY_FONT)

    # Footer
    add_text(slide, PAGE_PAD, SLIDE_H - emu_in(0.6),
             SLIDE_W - 2 * PAGE_PAD, emu_in(0.4),
             "ACOOLA TEAM · digital-агентство полного цикла",
             size=12, color=GRAY, font=BODY_FONT, align="center")


# ---------- Main ----------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01, slide_02, slide_03, slide_04, slide_05,
        slide_06, slide_07, slide_08, slide_09, slide_10,
        slide_11, slide_12, slide_13, slide_14, slide_15,
        slide_16, slide_17, slide_18, slide_19, slide_20,
        slide_21, slide_22,
    ]
    for b in builders:
        b(prs)

    out_path = "/Users/pro2kuror/Desktop/Cloude/docs/acoola_marketing_kit_v2.pptx"
    prs.save(out_path)
    print(f"OK: {len(prs.slides)} slides → {out_path}")


if __name__ == "__main__":
    main()
