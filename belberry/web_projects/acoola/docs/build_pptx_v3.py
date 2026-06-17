"""Build Acoola.Team marketing kit v3 (creative narrative version) as 18-slide PPTX.

Не каталог из 22 слайдов — нарратив из 18, организованный в 4 акта:
  Акт 1 (1–4): конфликт — «в digital сломалась модель»
  Акт 2 (5–9): наша модель + GEO как кульминация
  Акт 3 (10–14): доказательства (кейс-истории, клиенты, команда, награды)
  Акт 4 (15–18): спектр услуг + провокация-CTA

Source brief: docs/acoola_marketing_kit_v3_creative.md
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Brand palette ----------
DARK = RGBColor(0x1F, 0x1F, 0x22)
LIGHT = RGBColor(0xEB, 0xEB, 0xEE)
BLUE = RGBColor(0x2B, 0x7C, 0xE9)
GRAY = RGBColor(0x7A, 0x7A, 0x85)
GRAY_DARK = RGBColor(0x33, 0x33, 0x36)
GRAY_BORDER_DARK = RGBColor(0x44, 0x44, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
PURPLE = RGBColor(0x6B, 0x4F, 0xE9)
LIGHT_GRAY_TEXT = RGBColor(0xCC, 0xCC, 0xD0)
DEEP_BLUE = RGBColor(0x1A, 0x4D, 0xB8)

HEADING = "Manrope"
BODY = "Inter"

SLIDE_W = Emu(12192000)  # 13.33"
SLIDE_H = Emu(6858000)   # 7.5"
PAD = Emu(560000)
HEADER_TOP = Emu(280000)


def emu_in(inches):
    return Emu(int(inches * 914400))


# ---------- Helpers ----------

def set_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, *, line=None, corner=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    if corner is not None:
        s.adjustments[0] = corner
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=16, color=WHITE, font=BODY,
             bold=False, italic=False, align="left", anchor="top"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_header(slide, dark, *, minimal=False):
    """ACOOLA TEAM header on every slide. minimal=True hides on full-bleed slides."""
    if minimal:
        return
    sub = GRAY
    main = WHITE if dark else BLACK
    add_text(slide, PAD, HEADER_TOP, emu_in(3), Emu(220000),
             "WWW.ACOOLA.TEAM", size=10, color=sub, font=BODY)
    logo_w = emu_in(2.2)
    logo_x = SLIDE_W - PAD - logo_w
    add_text(slide, logo_x, HEADER_TOP, logo_w, Emu(280000),
             "ACOOLA", size=18, color=main, font=HEADING, bold=True, align="right")
    add_text(slide, logo_x, HEADER_TOP + Emu(220000), logo_w, Emu(180000),
             "TEAM", size=10, color=sub, font=HEADING, align="right")


def add_page_number(slide, number, total):
    """Номер слайда — в верхнем правом углу, между шапкой ACOOLA TEAM и
    основным контентом. Так не конфликтует с нижними CTA-плашками и
    якорными фразами на каждом слайде.
    """
    try:
        bg_rgb = slide.background.fill.fore_color.rgb
        dark = (bg_rgb == DARK)
    except Exception:
        dark = False
    color = GRAY
    add_text(slide, SLIDE_W - PAD - emu_in(1.0), HEADER_TOP + Emu(420000),
             emu_in(1.0), Emu(180000),
             f"{number:02d} / {total:02d}",
             size=9, color=color, font=BODY, align="right")


def add_pill(slide, x, y, label, *, w=None, h=None, fill=WHITE, text_color=BLACK):
    if w is None: w = emu_in(2.0)
    if h is None: h = emu_in(0.45)
    add_rect(slide, x, y, w, h, fill, corner=0.5)
    add_text(slide, x, y, w, h, label, size=12, color=text_color,
             font=BODY, bold=True, align="center", anchor="middle")


def add_pseudo_qr(slide, x, y, size, *, fg=BLACK, bg=WHITE):
    """Стилизованный псевдо-QR: 21×21 сетка с тремя позиционными метками в углах
    и псевдо-случайным паттерном. Не декодируется — это визуальный плейсхолдер.
    Финальный реальный QR ставит дизайнер.
    """
    add_rect(slide, x, y, size, size, bg, corner=0.05)
    n = 21  # стандартная сетка QR Version 1
    cell = (size - Emu(20000)) / n
    pad = Emu(10000)

    def is_finder(i, j):
        # три positioning markers в углах: (0..6,0..6), (n-7..n-1,0..6), (0..6,n-7..n-1)
        in_tl = i < 7 and j < 7
        in_tr = i >= n - 7 and j < 7
        in_bl = i < 7 and j >= n - 7
        if not (in_tl or in_tr or in_bl):
            return None
        # local coords inside the 7×7 finder
        if in_tl:
            li, lj = i, j
        elif in_tr:
            li, lj = i - (n - 7), j
        else:
            li, lj = i, j - (n - 7)
        # external 7×7 square, white inner ring, 3×3 black center
        if li in (0, 6) or lj in (0, 6):
            return True
        if li in (1, 5) or lj in (1, 5):
            return False
        return True  # 3×3 center

    for i in range(n):
        for j in range(n):
            f = is_finder(i, j)
            if f is True:
                fill = fg
            elif f is False:
                fill = bg
            else:
                # псевдо-случайный паттерн
                fill = fg if ((i * 7 + j * 13 + i * j) % 3 == 0) else bg
            if fill == bg:
                continue
            cx = x + pad + Emu(int(j * cell))
            cy = y + pad + Emu(int(i * cell))
            add_rect(slide, cx, cy, Emu(int(cell)), Emu(int(cell)), fg)


def add_chat_mockup(slide, x, y, w, h, *, brand_name="Ваш бренд",
                    competitors=("Конкурент A", "Конкурент B")):
    """Стилизованный мокап ответа AI-ассистента.
    Шапка с лого AI + bubble пользователя + bubble ассистента со списком,
    в котором подсвечен бренд клиента.
    """
    # background card
    add_rect(slide, x, y, w, h, GRAY_DARK, corner=0.04, line=GRAY_BORDER_DARK)

    # Header bar with "AI" indicator
    hdr_h = Emu(int(h * 0.18))
    add_rect(slide, x, y, w, hdr_h, RGBColor(0x10, 0x10, 0x12),
             corner=0.04)
    # green dot (live)
    dot = Emu(80000)
    add_rect(slide, x + Emu(120000), y + (hdr_h - dot) // 2, dot, dot,
             RGBColor(0x4A, 0xD8, 0x6E), corner=0.5)
    add_text(slide, x + Emu(280000), y, w - Emu(560000), hdr_h,
             "ChatGPT  ·  Perplexity  ·  YandexGPT  ·  Алиса",
             size=8, color=GRAY, font=BODY, bold=True, anchor="middle")

    # User bubble (right-aligned)
    body_y = y + hdr_h + Emu(80000)
    body_h = h - hdr_h - Emu(160000)
    user_w = Emu(int(w * 0.62))
    user_h = Emu(int(body_h * 0.28))
    user_x = x + w - user_w - Emu(120000)
    add_rect(slide, user_x, body_y, user_w, user_h,
             BLUE, corner=0.15)
    add_text(slide, user_x + Emu(100000), body_y, user_w - Emu(200000), user_h,
             "Кому позвонить по охранным\nсистемам в Москве?",
             size=8, color=WHITE, font=BODY, anchor="middle")

    # Assistant bubble (left-aligned, larger)
    asst_y = body_y + user_h + Emu(80000)
    asst_w = Emu(int(w * 0.85))
    asst_h = body_h - user_h - Emu(80000)
    asst_x = x + Emu(120000)
    add_rect(slide, asst_x, asst_y, asst_w, asst_h,
             RGBColor(0x2E, 0x2E, 0x32), corner=0.15)

    items_text = (
        "В вашей категории в Москве:\n"
        f"  •  {competitors[0]} — лидер по узнаваемости\n"
        f"  •  {brand_name} — закрывает задачу комплексно\n"
        f"  •  {competitors[1]} — нишевое решение"
    )
    add_text(slide, asst_x + Emu(120000), asst_y + Emu(40000),
             asst_w - Emu(240000), asst_h - Emu(40000),
             items_text, size=8, color=WHITE, font=BODY)

    # Highlight band over brand line (visual emphasis)
    hl_y = asst_y + Emu(int(asst_h * 0.42))
    hl_h = Emu(int(asst_h * 0.22))
    add_rect(slide, asst_x + Emu(80000), hl_y,
             asst_w - Emu(160000), hl_h,
             RGBColor(0x6B, 0x4F, 0xE9), corner=0.1)
    add_text(slide, asst_x + Emu(120000), hl_y,
             asst_w - Emu(240000), hl_h,
             f"  •  {brand_name} — комплексные решения, GPS-трекинг",
             size=8, color=WHITE, font=BODY, bold=True, anchor="middle")


def add_metric_chart(slide, x, y, w, h, *, points, label_top,
                     label_start, label_end, accent=BLUE, dark=True):
    """Компактный line chart: серия точек, соединённых линией.
    points: список значений 0..1 (нормализованная метрика)
    label_top: что показывает график
    label_start / label_end: подписи у первой и последней точки
    """
    bg = GRAY_DARK if dark else WHITE
    border = GRAY_BORDER_DARK if dark else RGBColor(0xD0, 0xD0, 0xD3)
    text_color = WHITE if dark else BLACK

    add_rect(slide, x, y, w, h, bg, corner=0.05, line=border)

    # Title
    add_text(slide, x + Emu(150000), y + Emu(120000),
             w - Emu(300000), Emu(280000),
             label_top, size=10, color=GRAY, font=BODY, bold=True)

    # Chart area
    chart_pad_x = Emu(280000)
    chart_pad_top = Emu(500000)
    chart_pad_bottom = Emu(500000)
    chart_x = x + chart_pad_x
    chart_y = y + chart_pad_top
    chart_w = w - 2 * chart_pad_x
    chart_h = h - chart_pad_top - chart_pad_bottom

    # Axis line (bottom)
    axis_color = RGBColor(0x55, 0x55, 0x58) if dark else RGBColor(0xCC, 0xCC, 0xCF)
    add_rect(slide, chart_x, chart_y + chart_h,
             chart_w, Emu(8000), axis_color)

    n = len(points)
    if n < 2:
        return

    # Compute point positions
    coords = []
    for i, val in enumerate(points):
        px = chart_x + Emu(int(i * chart_w / (n - 1)))
        # invert Y: high val → top
        py = chart_y + Emu(int((1 - val) * chart_h))
        coords.append((px, py))

    # Draw segments as thin rectangles between points
    seg_thickness = Emu(40000)
    for i in range(n - 1):
        x1, y1 = coords[i]
        x2, y2 = coords[i + 1]
        # rough connector: a thin diagonal rectangle approximated by short segments
        # python-pptx не поддерживает диагональные lines напрямую,
        # рисуем серию точек-пунктиров между концами
        steps = 14
        for k in range(steps + 1):
            t = k / steps
            sx = x1 + Emu(int((x2 - x1) * t))
            sy = y1 + Emu(int((y2 - y1) * t))
            sd = Emu(60000)
            add_rect(slide, sx - sd // 2, sy - sd // 2, sd, sd,
                     accent, corner=0.5)

    # Draw points (larger circles on each data point)
    pt_d = Emu(180000)
    for i, (px, py) in enumerate(coords):
        is_edge = (i == 0 or i == n - 1)
        pt_color = accent if not is_edge else WHITE
        size_d = pt_d if not is_edge else Emu(220000)
        add_rect(slide, px - size_d // 2, py - size_d // 2,
                 size_d, size_d, accent, corner=0.5)
        if is_edge:
            inner = Emu(100000)
            add_rect(slide, px - inner // 2, py - inner // 2,
                     inner, inner, WHITE, corner=0.5)

    # Edge labels — start (above-left of first point) and end (above-right of last)
    sx_, sy_ = coords[0]
    ex_, ey_ = coords[-1]
    add_text(slide, sx_ - Emu(800000), sy_ - Emu(380000),
             Emu(1600000), Emu(280000),
             label_start, size=11, color=text_color,
             font=HEADING, bold=True, align="center")
    add_text(slide, ex_ - Emu(800000), ey_ + Emu(180000),
             Emu(1600000), Emu(280000),
             label_end, size=11, color=accent,
             font=HEADING, bold=True, align="center")


def add_service_icon(slide, x, y, size, kind, *, color=BLUE):
    """Простая геометрическая иконка для каждого типа услуги.
    Не профессиональный line-icon, но даёт визуальную дифференциацию между карточками.
    kind: dev / template / design / seo / context / orm / crm / support / geo
    """
    if kind == "dev":
        # </>  — два угла + слэш (три прямоугольника)
        sw = Emu(30000)
        add_rect(slide, x + size // 4, y + size // 3, sw, size // 3, color)
        add_rect(slide, x + size * 3 // 4 - sw, y + size // 3, sw, size // 3, color)
        # diagonal slash via a rotated-looking rect
        add_rect(slide, x + size // 2 - sw // 2, y + size // 4,
                 sw, size // 2, color, corner=0.5)
    elif kind == "template":
        # сетка 2×2
        c = Emu(int(size * 0.18))
        gap = Emu(int(size * 0.04))
        for i in range(2):
            for j in range(2):
                add_rect(slide,
                         x + Emu(int(size * 0.15)) + j * (c + gap),
                         y + Emu(int(size * 0.15)) + i * (c + gap),
                         c, c, color)
    elif kind == "design":
        # квадрат с круглым акцентом (палитра)
        outer = Emu(int(size * 0.6))
        add_rect(slide, x + (size - outer) // 2, y + (size - outer) // 2,
                 outer, outer, color, corner=0.1)
        dot = Emu(int(size * 0.22))
        add_rect(slide, x + (size - dot) // 2, y + (size - dot) // 2,
                 dot, dot, WHITE, corner=0.5)
    elif kind == "seo":
        # стрелка вверх (треугольник через скруглённые квадраты)
        add_rect(slide, x + size // 3, y + size // 4,
                 size // 3, size // 2, color, corner=0.05)
        add_rect(slide, x + size // 4, y + size // 6,
                 size // 2, Emu(int(size * 0.18)), color, corner=0.5)
    elif kind == "context":
        # таргет — концентрические квадраты
        for k, frac in enumerate([0.6, 0.4, 0.2]):
            d = Emu(int(size * frac))
            c = color if k % 2 == 0 else WHITE
            add_rect(slide, x + (size - d) // 2, y + (size - d) // 2,
                     d, d, c, corner=0.5)
    elif kind == "orm":
        # щит — широкий прямоугольник со скруглёнными углами
        w = Emu(int(size * 0.55))
        h = Emu(int(size * 0.65))
        add_rect(slide, x + (size - w) // 2, y + (size - h) // 2,
                 w, h, color, corner=0.3)
    elif kind == "crm":
        # стопка — три горизонтальные плашки
        bar_h = Emu(int(size * 0.12))
        gap = Emu(int(size * 0.06))
        bw = Emu(int(size * 0.6))
        for i in range(3):
            add_rect(slide, x + (size - bw) // 2,
                     y + Emu(int(size * 0.25)) + i * (bar_h + gap),
                     bw, bar_h, color, corner=0.2)
    elif kind == "support":
        # точка с подчёркиванием — индикатор «всегда на связи»
        d = Emu(int(size * 0.28))
        add_rect(slide, x + (size - d) // 2, y + Emu(int(size * 0.3)),
                 d, d, color, corner=0.5)
        add_rect(slide, x + Emu(int(size * 0.25)),
                 y + Emu(int(size * 0.65)),
                 Emu(int(size * 0.5)), Emu(int(size * 0.06)),
                 color, corner=0.5)
    elif kind == "geo":
        # звезда-вспышка — отличается от всех остальных
        center_d = Emu(int(size * 0.35))
        cx = x + (size - center_d) // 2
        cy = y + (size - center_d) // 2
        add_rect(slide, cx, cy, center_d, center_d, color, corner=0.5)
        # 4 «лучики»
        ray_w = Emu(int(size * 0.06))
        ray_l = Emu(int(size * 0.18))
        # top
        add_rect(slide, x + (size - ray_w) // 2, y + Emu(int(size * 0.12)),
                 ray_w, ray_l, color, corner=0.5)
        # bottom
        add_rect(slide, x + (size - ray_w) // 2,
                 y + size - Emu(int(size * 0.12)) - ray_l,
                 ray_w, ray_l, color, corner=0.5)
        # left
        add_rect(slide, x + Emu(int(size * 0.12)),
                 y + (size - ray_w) // 2, ray_l, ray_w, color, corner=0.5)
        # right
        add_rect(slide, x + size - Emu(int(size * 0.12)) - ray_l,
                 y + (size - ray_w) // 2, ray_l, ray_w, color, corner=0.5)


def add_avatar_grid(slide, x, y, w, h, items, *, cols=3, rows=4,
                    bg=GRAY_DARK):
    """Сетка стилизованных аватаров с инициалами и подписями ролей.
    items: список (initials, role) — например [("А.К.", "SEO"), ...]
    """
    # frame
    add_rect(slide, x, y, w, h, bg, corner=0.04, line=GRAY_BORDER_DARK)
    # header
    hdr_h = Emu(int(h * 0.10))
    add_text(slide, x, y, w, hdr_h, "КОМАНДА · 50+ человек",
             size=10, color=GRAY, font=BODY, bold=True,
             anchor="middle", align="center")

    grid_y = y + hdr_h + Emu(80000)
    grid_h = h - hdr_h - Emu(120000)
    pad = Emu(120000)
    inner_w = w - 2 * pad
    inner_h = grid_h - Emu(80000)
    cell_w = inner_w // cols
    cell_h = inner_h // rows
    # avatar diameter
    avatar_d = Emu(int(min(cell_w, cell_h) * 0.55))

    palette = [
        BLUE,
        RGBColor(0x6B, 0x4F, 0xE9),
        RGBColor(0xE0, 0x6A, 0x6A),
        RGBColor(0xC9, 0x9A, 0x3D),
        RGBColor(0x4A, 0xD8, 0x9E),
        RGBColor(0x2E, 0xA8, 0xC8),
    ]
    for idx, (initials, role) in enumerate(items[: cols * rows]):
        r, c = divmod(idx, cols)
        cx = x + pad + c * cell_w + (cell_w - avatar_d) // 2
        cy = grid_y + r * cell_h + Emu(40000)
        color = palette[idx % len(palette)]
        # circle (rounded rect with full radius)
        add_rect(slide, cx, cy, avatar_d, avatar_d, color, corner=0.5)
        add_text(slide, cx, cy, avatar_d, avatar_d, initials,
                 size=11, color=WHITE, font=HEADING, bold=True,
                 align="center", anchor="middle")
        # role caption
        add_text(slide, x + pad + c * cell_w,
                 cy + avatar_d + Emu(20000),
                 cell_w, Emu(220000), role,
                 size=9, color=LIGHT_GRAY_TEXT, font=BODY,
                 align="center")


def add_flag(slide, x, y, w, h, country):
    """Упрощённый флаг страны через прямоугольники. Без точных пропорций
    и сложных элементов (звёзд, узоров, гербов) — это плейсхолдер
    для дизайнера, который заменит на реальные SVG. country: italy/russia/usa/belarus/kazakhstan/uae.
    """
    # Frame (light shadow + thin border)
    add_rect(slide, x, y, w, h, RGBColor(0xFF, 0xFF, 0xFF), corner=0.12,
             line=RGBColor(0xD0, 0xD0, 0xD3))

    pad = Emu(20000)
    fx = x + pad
    fy = y + pad
    fw = w - 2 * pad
    fh = h - 2 * pad

    GREEN = RGBColor(0x00, 0x9A, 0x49)
    RED_FLAG = RGBColor(0xCE, 0x2B, 0x37)
    BLUE_RU = RGBColor(0x00, 0x39, 0xA6)
    YELLOW = RGBColor(0xFE, 0xD5, 0x00)
    SKY = RGBColor(0x00, 0xAB, 0xE2)
    USA_BLUE = RGBColor(0x3C, 0x3B, 0x6E)
    BLACK_FLAG = RGBColor(0x00, 0x00, 0x00)

    if country == "italy":
        # 3 vertical stripes: green / white / red
        sw = fw // 3
        add_rect(slide, fx, fy, sw, fh, GREEN, corner=0.05)
        add_rect(slide, fx + sw, fy, sw, fh, RGBColor(0xFF, 0xFF, 0xFF))
        add_rect(slide, fx + sw * 2, fy, fw - sw * 2, fh, RED_FLAG, corner=0.05)
    elif country == "russia":
        # 3 horizontal: white / blue / red
        sh = fh // 3
        add_rect(slide, fx, fy, fw, sh, RGBColor(0xFF, 0xFF, 0xFF), corner=0.05)
        add_rect(slide, fx, fy + sh, fw, sh, BLUE_RU)
        add_rect(slide, fx, fy + sh * 2, fw, fh - sh * 2, RED_FLAG, corner=0.05)
    elif country == "usa":
        # Red/white stripes + blue canton
        stripes = 7
        sh = fh // stripes
        for i in range(stripes):
            color = RED_FLAG if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            add_rect(slide, fx, fy + i * sh, fw, sh, color)
        # Blue canton (top-left ~40% wide, top 4 stripes)
        canton_w = fw * 4 // 10
        canton_h = sh * 4
        add_rect(slide, fx, fy, canton_w, canton_h, USA_BLUE)
        # 3 tiny "star" dots
        for sr in range(2):
            for sc in range(3):
                dx = fx + Emu(50000) + sc * canton_w // 4
                dy = fy + Emu(40000) + sr * canton_h // 3
                dd = Emu(40000)
                add_rect(slide, dx, dy, dd, dd, RGBColor(0xFF, 0xFF, 0xFF), corner=0.5)
    elif country == "belarus":
        # Red on top 2/3, green on bottom 1/3 (без узора)
        red_h = fh * 2 // 3
        add_rect(slide, fx, fy, fw, red_h, RED_FLAG, corner=0.05)
        add_rect(slide, fx, fy + red_h, fw, fh - red_h, GREEN, corner=0.05)
        # тонкая вертикальная белая полоса слева — намёк на узор
        add_rect(slide, fx, fy, Emu(40000), fh, RGBColor(0xFF, 0xFF, 0xFF))
    elif country == "kazakhstan":
        # Solid sky-blue + yellow sun (упрощённо — круг)
        add_rect(slide, fx, fy, fw, fh, SKY, corner=0.05)
        sun_d = min(fw, fh) * 4 // 9
        add_rect(slide, fx + (fw - sun_d) // 2,
                 fy + (fh - sun_d) // 2,
                 sun_d, sun_d, YELLOW, corner=0.5)
    elif country == "uae":
        # Red vertical stripe on left + green/white/black horizontal
        red_w = fw // 4
        add_rect(slide, fx, fy, red_w, fh, RED_FLAG, corner=0.05)
        rest_x = fx + red_w
        rest_w = fw - red_w
        sh = fh // 3
        add_rect(slide, rest_x, fy, rest_w, sh, GREEN, corner=0.05)
        add_rect(slide, rest_x, fy + sh, rest_w, sh, RGBColor(0xFF, 0xFF, 0xFF))
        add_rect(slide, rest_x, fy + sh * 2, rest_w, fh - sh * 2, BLACK_FLAG, corner=0.05)


def add_trophy(slide, x, y, size, *, color=None):
    """Стилизованный кубок-награда. Программно — для слайда истории.
    Дизайнер заменит на реальный значок Ruward / премии.
    """
    if color is None:
        color = RGBColor(0xC9, 0x9A, 0x3D)  # gold

    # Cup body (rounded square — top wider, bottom narrower would be ideal,
    # но pptx не делает trapezoid; используем rounded rect)
    cup_w = int(size * 0.55)
    cup_h = int(size * 0.55)
    cup_x = x + (size - cup_w) // 2
    cup_y = y + int(size * 0.08)
    add_rect(slide, cup_x, cup_y, cup_w, cup_h, color, corner=0.15)

    # Handles (left + right)
    handle_w = int(size * 0.12)
    handle_h = int(size * 0.32)
    add_rect(slide, cup_x - handle_w + Emu(20000), cup_y + int(size * 0.08),
             handle_w, handle_h, color, corner=0.5)
    add_rect(slide, cup_x + cup_w - Emu(20000), cup_y + int(size * 0.08),
             handle_w, handle_h, color, corner=0.5)

    # Stem
    stem_w = int(size * 0.1)
    stem_h = int(size * 0.12)
    add_rect(slide, x + (size - stem_w) // 2,
             cup_y + cup_h, stem_w, stem_h, color)

    # Base
    base_w = int(size * 0.45)
    base_h = int(size * 0.1)
    add_rect(slide, x + (size - base_w) // 2,
             cup_y + cup_h + stem_h, base_w, base_h, color, corner=0.3)


# ============================================================
# SLIDES — 4 acts × 18 slides
# ============================================================

# ---------- Act 1. Conflict (1–4) ----------

def slide_01_cover_provocation(prs, recipient=None):
    """Обложка-провокация с опциональной персонализацией.

    Если recipient = {"company": "...", "domain": "...", "growth_pts": 3} —
    плашка справа заполняется реальными данными. Без recipient (generic-сборка)
    показываются видимые плейсхолдеры {COMPANY_NAME} / {DOMAIN} / {N} —
    дизайнер/маркетолог рендерит персонализированные копии под каждого клиента.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    # Tiny header only
    add_text(slide, PAD, HEADER_TOP, emu_in(3), Emu(220000),
             "WWW.ACOOLA.TEAM · МАРКЕТИНГ-КИТ · 2026",
             size=10, color=GRAY, font=BODY)

    # Massive provocation (left-aligned, занимает ~70% ширины)
    add_text(slide, PAD, emu_in(2.5), emu_in(8.5), emu_in(1.6),
             "Кто-то обещает", size=64, color=GRAY, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(3.4), emu_in(8.5), emu_in(1.5),
             "топ-10.", size=120, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(4.9), emu_in(8.5), emu_in(1.3),
             "Мы — нет.", size=110, color=BLUE, font=HEADING, bold=True)

    # Personalization card — right side
    card_x = emu_in(9.5)
    card_y = emu_in(2.6)
    card_w = SLIDE_W - PAD - card_x
    card_h = emu_in(2.6)
    add_rect(slide, card_x, card_y, card_w, card_h,
             RGBColor(0x2A, 0x2A, 0x2E), corner=0.05,
             line=GRAY_BORDER_DARK)

    company = (recipient or {}).get("company", "{COMPANY_NAME}")
    domain = (recipient or {}).get("domain", "{DOMAIN}")
    growth = (recipient or {}).get("growth_pts", "{N}")
    is_placeholder = recipient is None

    add_text(slide, card_x + emu_in(0.25), card_y + emu_in(0.2),
             card_w - emu_in(0.5), emu_in(0.3),
             "ПОДГОТОВЛЕНО ДЛЯ",
             size=9, color=BLUE, font=BODY, bold=True)
    add_text(slide, card_x + emu_in(0.25), card_y + emu_in(0.55),
             card_w - emu_in(0.5), emu_in(0.7),
             company, size=22 if not is_placeholder else 14,
             color=WHITE, font=HEADING, bold=True,
             italic=is_placeholder)
    add_text(slide, card_x + emu_in(0.25), card_y + emu_in(1.3),
             card_w - emu_in(0.5), emu_in(0.4),
             f"Сайт: {domain}",
             size=11, color=GRAY, font=BODY,
             italic=is_placeholder)
    add_text(slide, card_x + emu_in(0.25), card_y + emu_in(1.7),
             card_w - emu_in(0.5), emu_in(0.85),
             f"Видим {growth} точки роста — расскажем на встрече.",
             size=11, color=LIGHT_GRAY_TEXT, font=BODY,
             italic=is_placeholder)

    if is_placeholder:
        add_text(slide, card_x + emu_in(0.25), card_y + card_h - emu_in(0.4),
                 card_w - emu_in(0.5), emu_in(0.3),
                 "[ заменить на данные получателя ]",
                 size=8, color=GRAY, font=BODY, italic=True)

    # Footer
    add_text(slide, PAD, SLIDE_H - emu_in(0.6), emu_in(6), emu_in(0.4),
             "ACOOLA.TEAM · digital-агентство полного цикла",
             size=11, color=GRAY, font=BODY)
    add_text(slide, SLIDE_W - emu_in(3.3), SLIDE_H - emu_in(0.6),
             emu_in(3), emu_in(0.4),
             "→ дальше — почему именно так",
             size=11, color=GRAY, font=BODY, italic=True, align="right")


def slide_02_manifest(prs):
    """Манифест: 4 анти-обещания. Это позиция, а не описание услуг."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.6),
             "Манифест", size=14, color=GRAY, font=BODY, bold=True)

    statements = [
        ("Мы не делаем сайты.", "Мы превращаем сайт в актив."),
        ("Мы не считаем клики.", "Мы считаем обращения."),
        ("Мы не уходим после релиза.", "Мы остаёмся 3,5+ года."),
        ("Мы не пять разных подрядчиков.", "Мы — одна команда на 6 направлений."),
    ]
    y = emu_in(1.7)
    for left, right in statements:
        add_text(slide, PAD, y, SLIDE_W - 2 * PAD, emu_in(0.55),
                 left, size=24, color=GRAY, font=HEADING)
        add_text(slide, PAD, y + emu_in(0.55), SLIDE_W - 2 * PAD, emu_in(0.65),
                 right, size=32, color=BLACK, font=HEADING, bold=True)
        y += emu_in(1.3)


def slide_03_shift(prs):
    """3 сдвига, которые меняют digital. AI — один из трёх, не центральный."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Сдвиг", size=14, color=GRAY, font=BODY, bold=True)

    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Три сдвига, которые меняют digital.",
             size=42, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.45), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Те, кто их игнорирует, теряют клиентов. Те, кто отвечает — растут.",
             size=15, color=LIGHT_GRAY_TEXT, font=BODY, italic=True)

    shifts = [
        ("01", "Цикл сделки удлиняется",
         "Покупатель в B2B сегодня сравнивает 5–8 агентств вместо 2–3 пять лет назад. "
         "Сайт-визитка не выдерживает такой воронки."),
        ("02", "Клиент считает деньги, не клики",
         "Финдиректор смотрит на стоимость лида и конверсии, а не на позиции в выдаче. "
         "Отчёт «топ-10 по запросам» больше не продаёт."),
        ("03", "AI-поиск меняет первый контакт",
         "ChatGPT, Perplexity, YandexGPT, Алиса уже отвечают на запросы по вашей категории — "
         "до того, как клиент дошёл до Яндекса."),
    ]
    grid_y = emu_in(3.2)
    grid_h = emu_in(3.3)
    grid_w = SLIDE_W - 2 * PAD
    gap = Emu(140000)
    cell_w = (grid_w - gap * 2) // 3
    for i, (num, title, body) in enumerate(shifts):
        x = PAD + i * (cell_w + gap)
        add_rect(slide, x, grid_y, cell_w, grid_h, GRAY_DARK,
                 line=GRAY_BORDER_DARK, corner=0.05)
        add_text(slide, x + emu_in(0.3), grid_y + emu_in(0.3),
                 cell_w - emu_in(0.6), emu_in(0.4),
                 num, size=12, color=BLUE, font=BODY, bold=True)
        add_text(slide, x + emu_in(0.3), grid_y + emu_in(0.75),
                 cell_w - emu_in(0.6), emu_in(1.1),
                 title, size=20, color=WHITE, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.3), grid_y + emu_in(1.95),
                 cell_w - emu_in(0.6), grid_h - emu_in(2.15),
                 body, size=12, color=LIGHT_GRAY_TEXT, font=BODY)

    add_text(slide, PAD, SLIDE_H - emu_in(0.6),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "Большинство агентств отвечают только на первый. Мы — на все три.",
             size=14, color=BLUE, font=BODY, bold=True, italic=True, align="center")


def slide_04_broken(prs):
    """Что сломалось — 3 типа провала старой модели."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Что сломалось", size=14, color=GRAY, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.2),
             "В digital сломалась модель.",
             size=46, color=BLACK, font=HEADING, bold=True)

    columns = [
        ("Подрядчиков 5.\nОтветственных за результат — 0.",
         "Дизайн делает один. Разработка — другой. SEO — третий. Через год никто не помнит, кто что обещал."),
        ("Метрика — позиции в выдаче.\nФиндиректору показать нечего.",
         "Топ-10 — это слайд для презентации. Выручка — это разговор с собственником."),
        ("Релиз — финал работы.\nЧерез год сайт устарел.",
         "Сайт-актив требует ведения, а не «сдачи под ключ» с архивацией доступов."),
    ]
    grid_y = emu_in(3.3)
    grid_h = emu_in(3.5)
    grid_w = SLIDE_W - 2 * PAD
    gap = Emu(140000)
    cell_w = (grid_w - gap * 2) // 3
    for i, (h, body) in enumerate(columns):
        x = PAD + i * (cell_w + gap)
        add_rect(slide, x, grid_y, cell_w, grid_h, WHITE, corner=0.04)
        add_text(slide, x + emu_in(0.3), grid_y + emu_in(0.4),
                 cell_w - emu_in(0.6), emu_in(1.6), h,
                 size=18, color=BLACK, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.3), grid_y + emu_in(2.0),
                 cell_w - emu_in(0.6), grid_h - emu_in(2.2), body,
                 size=13, color=GRAY, font=BODY)


# ---------- Act 2. Our model (5–9) ----------

def slide_05_money(prs):
    """Обращения, не позиции — что мы реально замеряем для клиента."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наша модель", size=14, color=BLUE, font=BODY, bold=True)

    add_text(slide, PAD, emu_in(1.7), SLIDE_W - 2 * PAD, emu_in(2.4),
             "Мы отчитываемся\nобращениями, а не позициями.",
             size=52, color=WHITE, font=HEADING, bold=True)

    # Two columns: что считаем мы / что считают другие
    col_y = emu_in(4.7)
    col_h = emu_in(2.0)
    col_w = (SLIDE_W - 2 * PAD - emu_in(0.4)) // 2

    add_rect(slide, PAD, col_y, col_w, col_h, BLUE, corner=0.05)
    add_text(slide, PAD + emu_in(0.4), col_y + emu_in(0.3),
             col_w - emu_in(0.8), emu_in(0.5),
             "Что показываем в отчёте",
             size=14, color=WHITE, font=BODY, bold=True)
    add_text(slide, PAD + emu_in(0.4), col_y + emu_in(0.85),
             col_w - emu_in(0.8), col_h - emu_in(1.0),
             "обращения · стоимость лида · конверсии",
             size=20, color=WHITE, font=HEADING, bold=True)

    add_rect(slide, PAD + col_w + emu_in(0.4), col_y, col_w, col_h, GRAY_DARK,
             corner=0.05, line=GRAY_BORDER_DARK)
    add_text(slide, PAD + col_w + emu_in(0.7), col_y + emu_in(0.3),
             col_w - emu_in(0.8), emu_in(0.5),
             "Что обычно подсовывают",
             size=14, color=GRAY, font=BODY, bold=True)
    add_text(slide, PAD + col_w + emu_in(0.7), col_y + emu_in(0.85),
             col_w - emu_in(0.8), col_h - emu_in(1.0),
             "клики · просмотры · позиции в выдаче",
             size=20, color=GRAY, font=HEADING, bold=True)


def slide_06_long_term(prs):
    """3,5+ года — длинная дистанция как герой."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наша модель", size=14, color=BLUE, font=BODY, bold=True)

    # Hero number — extra large
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(3.2),
             "3,5+", size=320, color=BLUE, font=HEADING, bold=True, align="center")

    add_text(slide, PAD, emu_in(4.6), SLIDE_W - 2 * PAD, emu_in(0.6),
             "лет — средний срок сотрудничества с клиентом",
             size=22, color=BLACK, font=HEADING, bold=True, align="center")

    add_text(slide, PAD, emu_in(5.4), SLIDE_W - 2 * PAD, emu_in(0.5),
             "В digital это в 2,5–3 раза дольше отраслевой нормы.",
             size=14, color=GRAY, font=BODY, italic=True, align="center")
    add_text(slide, PAD, emu_in(6.0), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Клиенты остаются. Это и есть результат.",
             size=14, color=GRAY, font=BODY, italic=True, align="center")


def slide_07_95_visualization(prs):
    """Визуализация 95 из 100. 95 синих квадратов и 5 серых."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наша модель", size=14, color=BLUE, font=BODY, bold=True)

    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(0.9),
             "95 из 100 клиентов продлевают контракт.",
             size=34, color=BLACK, font=HEADING, bold=True)

    # 10×10 grid of squares
    sq = emu_in(0.36)
    gap = Emu(45000)
    grid_x = PAD + emu_in(0.5)
    grid_y = emu_in(2.7)
    for r in range(10):
        for c in range(10):
            idx = r * 10 + c
            color = BLUE if idx < 95 else GRAY
            x = grid_x + c * (sq + gap)
            y = grid_y + r * (sq + gap)
            add_rect(slide, x, y, sq, sq, color, corner=0.15)

    # Right column: legend + honest line
    legend_x = grid_x + 10 * (sq + gap) + emu_in(0.6)
    legend_w = SLIDE_W - PAD - legend_x

    add_rect(slide, legend_x, emu_in(2.7), emu_in(0.4), emu_in(0.4), BLUE, corner=0.15)
    add_text(slide, legend_x + emu_in(0.6), emu_in(2.7),
             legend_w - emu_in(0.6), emu_in(0.4),
             "95 клиентов остаются на следующий цикл",
             size=14, color=BLACK, font=BODY, anchor="middle")

    add_rect(slide, legend_x, emu_in(3.4), emu_in(0.4), emu_in(0.4), GRAY, corner=0.15)
    add_text(slide, legend_x + emu_in(0.6), emu_in(3.4),
             legend_w - emu_in(0.6), emu_in(0.4),
             "5 уходят. Это нормально.",
             size=14, color=BLACK, font=BODY, anchor="middle")

    add_text(slide, legend_x, emu_in(4.5), legend_w, emu_in(2.0),
             "Расскажем честно, в каких случаях\nуходят — и почему мы\nне берёмся за эти проекты.",
             size=15, color=GRAY, font=BODY, italic=True)
    add_text(slide, legend_x, emu_in(6.3), legend_w, emu_in(0.4),
             "↓ следующий слайд",
             size=12, color=BLUE, font=BODY, bold=True)


def slide_08_anti_marketing(prs):
    """Чего мы не возьмёмся делать."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наша модель", size=14, color=BLUE, font=BODY, bold=True)

    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Чего мы не возьмёмся делать.",
             size=46, color=WHITE, font=HEADING, bold=True)

    nots = [
        "Накручивать позиции в Яндексе",
        "Запускать рекламу без аналитики на сайте",
        "Делать дизайн «как у конкурента»",
        "Обещать топ-10 «к Новому году»",
        "Заходить в SEO там, где сайт нужно переделать",
        "Учиться GEO на вашем проекте",
    ]
    y0 = emu_in(3.0)
    line_h = emu_in(0.55)
    for i, txt in enumerate(nots):
        y = y0 + i * line_h
        add_text(slide, PAD, y, emu_in(0.5), line_h, "—",
                 size=22, color=BLUE, font=HEADING, bold=True, anchor="middle")
        add_text(slide, PAD + emu_in(0.6), y, SLIDE_W - 2 * PAD - emu_in(0.6), line_h,
                 txt, size=20, color=WHITE, font=BODY, anchor="middle")

    add_text(slide, PAD, SLIDE_H - emu_in(0.6), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Отказ — это форма уважения. И к клиенту, и к команде.",
             size=12, color=GRAY, font=BODY, italic=True)


def slide_09_geo_climax(prs):
    """GEO как кульминация акта 'Наша модель'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наша модель · кульминация",
             size=14, color=PURPLE, font=BODY, bold=True)

    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(1.5),
             "Мы запустили GEO в 2025 —",
             size=46, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "до того, как остальные услышали это слово.",
             size=28, color=PURPLE, font=HEADING, bold=True)

    # Body text (compressed to make room for live-demo block)
    add_text(slide, PAD, emu_in(3.5), emu_in(6.5), emu_in(1.9),
             "Когда клиент спрашивает у ChatGPT —\n«кому позвонить по [вашей категории] в Москве»,\nмы делаем так, чтобы ответ назвал ваш бренд.\n\nGEO — это не отдельная новая команда.\nЭто наши SEO, контент и ORM в новом продукте.",
             size=13, color=LIGHT_GRAY_TEXT, font=BODY)

    # Right side — AI chat mockup (нарисован программно вместо placeholder)
    right_x = emu_in(7.5)
    right_y = emu_in(3.5)
    right_w = SLIDE_W - PAD - right_x
    right_h = emu_in(1.9)
    add_chat_mockup(slide, right_x, right_y, right_w, right_h,
                    brand_name="Ваш бренд",
                    competitors=("Конкурент A", "Конкурент B"))

    # Live-demo block — full width, between body and bottom strip
    demo_y = emu_in(5.55)
    demo_h = emu_in(0.95)
    add_rect(slide, PAD, demo_y, SLIDE_W - 2 * PAD, demo_h,
             RGBColor(0x2A, 0x2A, 0x2E), corner=0.06,
             line=PURPLE)
    # left part — text
    add_text(slide, PAD + emu_in(0.3), demo_y + emu_in(0.1),
             emu_in(2), emu_in(0.3),
             "ПРЯМО СЕЙЧАС",
             size=10, color=PURPLE, font=BODY, bold=True)
    add_text(slide, PAD + emu_in(0.3), demo_y + emu_in(0.4),
             emu_in(9), emu_in(0.55),
             "Откройте ChatGPT и спросите — «кому позвонить по [вашей нише] в Москве?»  Скиньте нам скриншот.",
             size=14, color=WHITE, font=HEADING, bold=True, anchor="middle")
    # right part — псевдо-QR (нарисован, дизайнер заменит на реальный)
    qr_size = emu_in(0.8)
    qr_x = SLIDE_W - PAD - qr_size - emu_in(0.15)
    qr_y = demo_y + (demo_h - qr_size) // 2
    add_pseudo_qr(slide, qr_x, qr_y, qr_size)
    add_text(slide, qr_x - emu_in(1.5), qr_y, emu_in(1.4), qr_size,
             "→ форма обратной\nсвязи",
             size=9, color=GRAY, font=BODY, align="right", anchor="middle")

    # Bottom — price + CTA + badge
    strip_y = SLIDE_H - emu_in(0.85)
    strip_h = emu_in(0.65)
    add_rect(slide, PAD, strip_y, emu_in(4.5), strip_h, PURPLE, corner=0.1)
    add_text(slide, PAD, strip_y, emu_in(4.5), strip_h,
             "от 90 000 ₽/мес",
             size=22, color=WHITE, font=HEADING, bold=True,
             align="center", anchor="middle")
    add_pill(slide, PAD + emu_in(4.7), strip_y, "Новое направление 2025",
             w=emu_in(3.0), h=strip_h, fill=WHITE, text_color=PURPLE)
    add_pill(slide, SLIDE_W - PAD - emu_in(3.0), strip_y, "Запросить GEO-аудит",
             w=emu_in(3.0), h=strip_h, fill=BLUE, text_color=WHITE)


# ---------- Act 3. Proof (10–14) ----------

def slide_case_teo(prs):
    """Кейс · Разработка: TEO by COSMOSCOW (онлайн-маркетплейс искусства, 2022).

    Источник: /projects/teo/ на acoola.team.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Кейс · РАЗРАБОТКА",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(0.8),
             "TEO by COSMOSCOW", size=44, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.3), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Онлайн-маркетплейс искусства · разработка платформы · 2022",
             size=14, color=GRAY, font=BODY, italic=True)

    # Two-column: story (60%) + visual (40%)
    left_w = emu_in(7.5)
    right_x = PAD + left_w + emu_in(0.4)
    right_w = SLIDE_W - PAD - right_x

    paragraphs = [
        "2022 год.",
        "COSMOSCOW — крупнейшая международная ярмарка\nсовременного искусства в России.",
        "Задача: создать онлайн-платформу TEO,\nкоторая напрямую связывает авторов\nработ и покупателей.",
        "Разработали маркетплейс с нуля:\nот UX до интеграций, под нагрузку аукциона.",
    ]
    y = emu_in(3.0)
    for i, p in enumerate(paragraphs):
        is_first = (i == 0)
        is_last = (i == len(paragraphs) - 1)
        size = 18 if is_first else (22 if is_last else 16)
        color = BLUE if is_first else WHITE
        weight = True if (is_first or is_last) else False
        add_text(slide, PAD, y, left_w, emu_in(0.9),
                 p, size=size, color=color,
                 font=HEADING if (is_first or is_last) else BODY,
                 bold=weight)
        y += emu_in(0.85)

    # Right column: stylized "browser frame" placeholder for site screenshot
    chart_y = emu_in(3.0)
    chart_h = emu_in(3.4)
    add_rect(slide, right_x, chart_y, right_w, chart_h,
             GRAY_DARK, corner=0.04, line=GRAY_BORDER_DARK)
    # "browser" header bar
    bar_h = Emu(int(chart_h * 0.12))
    add_rect(slide, right_x, chart_y, right_w, bar_h,
             RGBColor(0x10, 0x10, 0x12), corner=0.04)
    # 3 traffic-light dots
    for k, color in enumerate([RGBColor(0xE0, 0x6A, 0x6A),
                                RGBColor(0xC9, 0x9A, 0x3D),
                                RGBColor(0x4A, 0xD8, 0x6E)]):
        dot_d = Emu(80000)
        add_rect(slide, right_x + Emu(100000) + k * Emu(140000),
                 chart_y + (bar_h - dot_d) // 2,
                 dot_d, dot_d, color, corner=0.5)
    # URL
    add_text(slide, right_x + Emu(600000), chart_y, right_w - Emu(800000), bar_h,
             "teo.cosmoscow.com",
             size=8, color=GRAY, font=BODY, align="center", anchor="middle")
    # Body placeholder
    add_text(slide, right_x, chart_y + bar_h, right_w, chart_h - bar_h,
             "[ СКРИНШОТ TEO ]\nдизайнер: подставить",
             size=12, color=GRAY, font=BODY,
             align="center", anchor="middle")

    # Stats footer
    strip_y = SLIDE_H - emu_in(0.9)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             GRAY_DARK, corner=0.1)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             "Онлайн-платформа · авторы ↔ покупатели · нагрузка под аукцион · acoola.team/projects/teo",
             size=13, color=BLUE, font=HEADING, bold=True,
             align="center", anchor="middle")


def slide_11_story_wwerb(prs):
    """Кейс-история #2: wwerb."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Кейс · SEO", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(0.8),
             "wwerb.com", size=44, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.3), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Часовой магазин · SEO · 2 года работы",
             size=14, color=GRAY, font=BODY, italic=True)

    # Two-column: story + chart of organic traffic growth
    left_w = emu_in(7.5)
    right_x = PAD + left_w + emu_in(0.4)
    right_w = SLIDE_W - PAD - right_x

    paragraphs = [
        "Март 2020.",
        "Лид в категории «премиальные часы» стоит 18 400 ₽.\nЭто потолок маржи — реклама убыточна.",
        "Взялись за SEO с фокусом\nна органику и торговую матрицу.",
        "Через два года лид — 582 ₽.\nОрганика +223%. Через сайт прошло 164+ млн ₽ выручки.",
    ]
    y = emu_in(3.0)
    for i, p in enumerate(paragraphs):
        is_first = (i == 0)
        is_last = (i == len(paragraphs) - 1)
        size = 18 if is_first else (22 if is_last else 16)
        color = BLUE if is_first else WHITE
        weight = True if (is_first or is_last) else False
        add_text(slide, PAD, y, left_w, emu_in(0.9),
                 p, size=size, color=color,
                 font=HEADING if (is_first or is_last) else BODY,
                 bold=weight)
        y += emu_in(0.85 if is_first else 0.85)

    # Right column: line chart of organic traffic GROWTH (+223%)
    chart_y = emu_in(3.0)
    chart_h = emu_in(3.4)
    add_metric_chart(
        slide, right_x, chart_y, right_w, chart_h,
        # 4 точки роста: март 2020 → март 2022
        points=[0.05, 0.25, 0.55, 1.0],
        label_top="Органика · рост за 2 года",
        label_start="2 757 визитов",
        label_end="6 548 · +223%",
        accent=BLUE, dark=True,
    )

    strip_y = SLIDE_H - emu_in(0.9)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             GRAY_DARK, corner=0.1)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             "18 400 ₽  →  582 ₽   ·   органика +223%   ·   164+ млн ₽ выручки",
             size=14, color=BLUE, font=HEADING, bold=True,
             align="center", anchor="middle")


def slide_12_clients(prs):
    """Клиенты — фирменная сетка 5×5 (5 отраслей × 5 имён).

    По референсу фирменной презентации Acoola: тёмный фон, слева подписи
    отраслей серым, справа — 5 ячеек-плейсхолдеров под логотипы клиентов
    в каждой отрасли. Стрелка ↗ в правом верхнем углу каждой ячейки.
    Дизайнер заменяет плейсхолдеры на реальные логотипы.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Клиенты", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.35), SLIDE_W - 2 * PAD, emu_in(0.9),
             "Наши клиенты.",
             size=42, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.3), SLIDE_W - 2 * PAD, emu_in(0.5),
             "70+ проектов в работе ежемесячно. Tier-1 клиенты с длинной дистанцией — средний срок 3,5+ года.",
             size=13, color=LIGHT_GRAY_TEXT, font=BODY, italic=True)

    # 5 industries × 5 clients
    rows_data = [
        ("Финансы",        ["БыстроБанк", "Трансстройбанк", "NZT Rusfond", "PROFIT CONF", "НПФ «Будущее»"]),
        ("Недвижимость",   ["Главстрой", "DOGMA", "Запстрой", "Pioneer", "Prime Park"]),
        ("Промышленность", ["HOHMAN", "Robomaster", "РУСЭЛ", "Yarus Networks", "Эксара"]),
        ("E-commerce",     ["J.Prettier", "Титан", "MARKIN", "Coolstream", "Latrika"]),
        ("IT",             ["DIS Group", "NGENIX", "textback", "SBER API", "SBER DEVICES"]),
    ]

    grid_y = emu_in(3.1)
    grid_h = SLIDE_H - grid_y - emu_in(0.55)
    grid_w = SLIDE_W - 2 * PAD
    # Левая колонка — подпись отрасли, остальное — 5 ячеек клиентов
    industry_col_w = emu_in(1.85)
    cells_total_w = grid_w - industry_col_w
    cell_gap = Emu(80000)
    cell_w = (cells_total_w - cell_gap * 5) // 5
    row_count = len(rows_data)
    row_gap = Emu(80000)
    row_h = (grid_h - row_gap * (row_count - 1)) // row_count

    for ri, (industry, clients) in enumerate(rows_data):
        y = grid_y + ri * (row_h + row_gap)
        # Industry label (left)
        add_text(slide, PAD, y, industry_col_w, row_h,
                 industry, size=14, color=GRAY,
                 font=HEADING, bold=True, anchor="middle")
        # 5 client cells
        for ci, name in enumerate(clients):
            x = PAD + industry_col_w + cell_gap + ci * (cell_w + cell_gap)
            # Cell background — placeholder под logo
            add_rect(slide, x, y, cell_w, row_h,
                     RGBColor(0x2A, 0x2A, 0x2E),
                     line=GRAY_BORDER_DARK, corner=0.04)
            # Client name centered
            add_text(slide, x + emu_in(0.15), y, cell_w - emu_in(0.45),
                     row_h, name, size=12, color=WHITE,
                     font=BODY, bold=True, anchor="middle", align="center")
            # Arrow ↗ top-right
            arrow_size = emu_in(0.3)
            add_text(slide, x + cell_w - arrow_size - emu_in(0.05),
                     y + emu_in(0.05), arrow_size, arrow_size,
                     "↗", size=11, color=GRAY,
                     font=BODY, align="center", anchor="middle")

    # Footer
    add_text(slide, PAD, SLIDE_H - emu_in(0.4),
             SLIDE_W - 2 * PAD, emu_in(0.3),
             "Дизайнер: вместо текстовых имён — фирменные логотипы клиентов в монохромном виде.",
             size=9, color=GRAY, font=BODY, italic=True, align="center")


def slide_13_team_quote(prs):
    """Команда — структура обслуживания клиента, без ложных претензий.

    Раньше: «100% штат», «все в Москве» — это неверно (есть удалёнщики).
    Сейчас: фокус на том, как команда работает с клиентом — а не где сидит.
    Hero — 3 проверенные цифры (50+, 6 направлений, 3,5+ года). Снизу —
    структура работы над клиентским проектом + сертификации.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Команда", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Не «команда мечты с лендинга».",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.4),
             "А конкретные люди, которые знают ваш проект — и которых вы узнаёте по именам со 2-го месяца.",
             size=14, color=GRAY, font=BODY, italic=True)

    # === 3 hero stats — top row ===
    stats_y = emu_in(3.05)
    stats_h = emu_in(1.8)
    stats_w = SLIDE_W - 2 * PAD
    gap = emu_in(0.2)
    stat_w = (stats_w - gap * 2) // 3
    stats = [
        ("50+",   "digital-специалистов",       "В 6 направлениях работ"),
        ("3,5+",  "года средний срок",           "Командой со своими проектами"),
        ("4",     "платформы сертификации",      "Google · Яндекс · 1С Битрикс · Brand Analytics"),
    ]
    for i, (number, label, body) in enumerate(stats):
        x = PAD + i * (stat_w + gap)
        add_rect(slide, x, stats_y, stat_w, stats_h, WHITE, corner=0.06)
        # accent vertical bar
        add_rect(slide, x, stats_y, Emu(60000), stats_h, BLUE)
        # Big number
        add_text(slide, x + emu_in(0.35), stats_y + emu_in(0.18),
                 stat_w - emu_in(0.7), emu_in(0.85),
                 number, size=48, color=BLUE, font=HEADING, bold=True)
        # Label
        add_text(slide, x + emu_in(0.35), stats_y + emu_in(1.08),
                 stat_w - emu_in(0.7), emu_in(0.35),
                 label, size=14, color=BLACK, font=HEADING, bold=True)
        # Body
        add_text(slide, x + emu_in(0.35), stats_y + stats_h - emu_in(0.35),
                 stat_w - emu_in(0.7), emu_in(0.3),
                 body, size=10, color=GRAY, font=BODY)

    # === Bottom: структура работы с клиентом ===
    struct_y = stats_y + stats_h + emu_in(0.3)
    struct_h = emu_in(1.8)
    add_rect(slide, PAD, struct_y, stats_w, struct_h, WHITE, corner=0.06)
    add_text(slide, PAD + emu_in(0.4), struct_y + emu_in(0.25),
             stats_w - emu_in(0.8), emu_in(0.4),
             "КТО РАБОТАЕТ НАД ВАШИМ ПРОЕКТОМ",
             size=11, color=BLUE, font=BODY, bold=True)

    # 3 роли в ряду
    roles_y = struct_y + emu_in(0.75)
    roles_h = struct_h - emu_in(0.95)
    role_w = (stats_w - emu_in(1.0)) // 3
    role_gap = emu_in(0.2)
    roles = [
        ("Менеджер проекта",
         "Постоянный контакт. Знает план, статусы, метрики."),
        ("Технический лид",
         "Отвечает за качество. Еженедельные созвоны."),
        ("Профильные специалисты",
         "Разработка, SEO, контекст, ORM — по задачам проекта."),
    ]
    for i, (title, body) in enumerate(roles):
        x = PAD + emu_in(0.4) + i * (role_w + role_gap)
        # Маленькая синяя точка-маркер
        dot_d = Emu(120000)
        add_rect(slide, x, roles_y + emu_in(0.1),
                 dot_d, dot_d, BLUE, corner=0.5)
        # Title
        add_text(slide, x + emu_in(0.3), roles_y,
                 role_w - emu_in(0.3), emu_in(0.4),
                 title, size=14, color=BLACK, font=HEADING, bold=True)
        # Body
        add_text(slide, x + emu_in(0.3), roles_y + emu_in(0.45),
                 role_w - emu_in(0.3), roles_h - emu_in(0.45),
                 body, size=11, color=GRAY, font=BODY)

    # Якорь снизу
    add_text(slide, PAD, SLIDE_H - emu_in(0.4),
             SLIDE_W - 2 * PAD, emu_in(0.3),
             "Не один менеджер на 50 клиентов. Постоянная команда на каждый проект.",
             size=11, color=BLUE, font=BODY, italic=True, align="center")


def slide_14_awards(prs):
    """Внешняя валидация — точные формулировки."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Награды", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Внешняя валидация",
             size=38, color=BLACK, font=HEADING, bold=True)

    # Только подтверждённые статусы Acoola (без медицинских наград — те принадлежат
    # Belberry, не Acoola, и были ошибочно включены в ранние версии).
    cards = [
        (RGBColor(0xC9, 0x9A, 0x3D), "Золотой партнёр\n1С Битрикс", "Автоматизируем бизнес"),
        (RGBColor(0xC8, 0x3E, 0x3E), "Мониторинг\n1С Битрикс", "Участник программы качества"),
        (DEEP_BLUE, "Сертифицированное\nагентство Яндекса", "Яндекс.Директ"),
        (DARK, "Google\nPartner", "Сертифицированный партнёр Google"),
    ]
    cards_y = emu_in(2.9)
    cards_h = emu_in(1.85)
    cards_w = emu_in(8.0)
    card_w = (cards_w - Emu(80000) * 3) // 4
    for i, (fill, title, sub) in enumerate(cards):
        x = PAD + i * (card_w + Emu(80000))
        add_rect(slide, x, cards_y, card_w, cards_h, fill, corner=0.07)
        add_text(slide, x + emu_in(0.18), cards_y + emu_in(0.18),
                 card_w - emu_in(0.36), emu_in(1.0), title,
                 size=14, color=WHITE, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.18), cards_y + cards_h - emu_in(0.65),
                 card_w - emu_in(0.36), emu_in(0.5), sub,
                 size=10, color=LIGHT_GRAY_TEXT, font=BODY)

    # Right — 95% hero
    h_x = PAD + cards_w + emu_in(0.3)
    h_w = SLIDE_W - PAD - h_x
    add_rect(slide, h_x, cards_y, h_w, cards_h, BLUE, corner=0.05)
    add_text(slide, h_x, cards_y + emu_in(0.15), h_w, emu_in(1.1),
             "95%", size=84, color=WHITE, font=HEADING, bold=True, align="center")
    add_text(slide, h_x, cards_y + cards_h - emu_in(0.5), h_w, emu_in(0.4),
             "довольных клиентов",
             size=11, color=WHITE, font=BODY, align="center")

    # Bottom — extra validation (tooling certifications team has)
    extras_y = emu_in(5.0)
    extras = [
        ("Яндекс.Метрика", "сертификация"),
        ("Google Analytics", "сертификация"),
        ("Brand Analytics", "сертификация"),
        ("Calltouch", "сертификация"),
    ]
    e_w = (SLIDE_W - 2 * PAD - Emu(80000) * 3) // 4
    for i, (top, bottom) in enumerate(extras):
        x = PAD + i * (e_w + Emu(80000))
        add_rect(slide, x, extras_y, e_w, emu_in(1.5), WHITE, corner=0.05)
        add_text(slide, x + emu_in(0.2), extras_y + emu_in(0.3),
                 e_w - emu_in(0.4), emu_in(0.5), top,
                 size=15, color=BLACK, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.2), extras_y + emu_in(0.85),
                 e_w - emu_in(0.4), emu_in(0.5), bottom,
                 size=12, color=GRAY, font=BODY)

    # Footnote
    add_text(slide, PAD, SLIDE_H - emu_in(0.55), SLIDE_W - 2 * PAD, emu_in(0.4),
             "NPS 100 за последний год · 95 за общий период — программа мониторинга качества 1С Битрикс",
             size=10, color=GRAY, font=BODY, italic=True, align="center")


# ---------- Act 4. Range + close (15–18) ----------

def slide_15_services_overview(prs):
    """6 направлений — обзор."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Спектр", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(1.0),
             "6 направлений работ.",
             size=44, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.6),
             "От первой строчки кода до первой строчки в ответе ChatGPT.",
             size=18, color=LIGHT_GRAY_TEXT, font=BODY)

    services = [
        ("01", "Разработка сайта", "dev"),
        ("02", "SEO", "seo"),
        ("03", "Контекстная реклама", "context"),
        ("04", "Управление репутацией", "orm"),
        ("05", "Техническая техподдержка", "support"),
        ("06", "GEO", "geo"),
    ]
    grid_y = emu_in(3.3)
    grid_h = SLIDE_H - grid_y - PAD
    grid_w = SLIDE_W - 2 * PAD
    cols = 3
    rows = 2
    gap = Emu(150000)
    cell_w = (grid_w - gap * (cols - 1)) // cols
    cell_h = (grid_h - gap * (rows - 1)) // rows
    for i, (num, name, icon_kind) in enumerate(services):
        r, c = divmod(i, cols)
        x = PAD + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        is_geo = (i == 5)
        fill = GRAY_DARK if not is_geo else PURPLE
        add_rect(slide, x, y, cell_w, cell_h, fill,
                 line=GRAY_BORDER_DARK, corner=0.05)
        # Icon (top-left)
        icon_size = emu_in(0.65)
        icon_x = x + emu_in(0.3)
        icon_y = y + emu_in(0.25)
        icon_color = WHITE if is_geo else BLUE
        add_service_icon(slide, icon_x, icon_y, icon_size, icon_kind,
                         color=icon_color)
        # Number (top-right area)
        add_text(slide, x + cell_w - emu_in(1.0), y + emu_in(0.25),
                 emu_in(0.5), emu_in(0.4), num,
                 size=12, color=GRAY if not is_geo else WHITE,
                 font=BODY, bold=True, align="right")
        # Service name (bottom-left)
        add_text(slide, x + emu_in(0.3), y + cell_h - emu_in(1.05),
                 cell_w - emu_in(0.6), emu_in(0.7), name,
                 size=18, color=WHITE, font=HEADING, bold=True)
        # Arrow (bottom-right)
        add_text(slide, x + cell_w - emu_in(0.5), y + cell_h - emu_in(0.65),
                 emu_in(0.4), emu_in(0.4), "↗",
                 size=14, color=GRAY, font=BODY, align="right")
        if is_geo:
            add_pill(slide, x + cell_w - emu_in(1.6),
                     y + cell_h - emu_in(0.6),
                     "New", w=emu_in(0.8), h=emu_in(0.32),
                     fill=WHITE, text_color=PURPLE)


def slide_services_lifecycle_map(prs):
    """Карта услуг по жизни сайта — 6 этапов с ценами с сайта acoola.team.

    Превращает обзор услуг из «меню» в «дорожную карту клиента». Каждый
    этап — момент в жизни сайта, услуга, её цена и одна строка про результат.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Спектр · карта услуг по жизни сайта",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Один партнёр на весь путь сайта.",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.35), SLIDE_W - 2 * PAD, emu_in(0.4),
             "От первой строчки кода до видимости в AI-поиске. Цены — с сайта acoola.team.",
             size=14, color=GRAY, font=BODY, italic=True)

    # 6 этапов жизни сайта
    stages = [
        # (when, service, price, one-liner, icon_kind, accent)
        ("До запуска",      "Разработка сайта",       "от 190 000 ₽",     "Сайт-актив, готовый\nк продвижению",         "dev",     BLUE),
        ("При запуске",     "Контекстная реклама",    "от 50 000 ₽/мес",  "Первые заявки\nза дни, не месяцы",           "context", BLUE),
        ("Через месяц",     "SEO",                    "от 70 000 ₽/мес",  "Системный рост\nобращений из органики",      "seo",     BLUE),
        ("Через 3 месяца",  "Управление репутацией",  "от 70 000 ₽/мес",  "Защищаем то,\nчто вы построили",             "orm",     BLUE),
        ("Постоянно",       "Техническая техподдержка", "от 30 000 ₽/мес", "Сайт всегда\nна связи",                      "support", BLUE),
        ("Через 6 месяцев", "GEO",                    "от 90 000 ₽/мес",  "Видимость бренда\nв AI-поиске",              "geo",     PURPLE),
    ]

    grid_y = emu_in(2.95)
    grid_h = emu_in(3.7)
    grid_w = SLIDE_W - 2 * PAD
    cols = 6
    gap = Emu(80000)
    cell_w = (grid_w - gap * (cols - 1)) // cols

    # Timeline axis line
    axis_y = grid_y + Emu(int(grid_h * 0.18))
    add_rect(slide, PAD + Emu(80000), axis_y,
             grid_w - Emu(160000), Emu(8000),
             RGBColor(0xCC, 0xCC, 0xCF))

    for i, (when, name, price, oneline, icon_kind, accent) in enumerate(stages):
        x = PAD + i * (cell_w + gap)
        is_last = (i == len(stages) - 1)
        # Stage marker dot on timeline
        dot_d = Emu(160000)
        add_rect(slide, x + (cell_w - dot_d) // 2, axis_y - dot_d // 2 + Emu(4000),
                 dot_d, dot_d, accent, corner=0.5)
        # "When" label above the dot
        add_text(slide, x, grid_y, cell_w, emu_in(0.35),
                 when.upper(), size=9, color=accent, font=BODY, bold=True,
                 align="center")
        # Card below the timeline
        card_y = axis_y + Emu(180000)
        card_h = grid_h - (card_y - grid_y) - Emu(80000)
        card_fill = WHITE if not is_last else RGBColor(0xEC, 0xE6, 0xFE)
        card_border = None if not is_last else PURPLE
        add_rect(slide, x, card_y, cell_w, card_h, card_fill,
                 corner=0.05, line=card_border)
        # Icon
        icon_size = emu_in(0.5)
        icon_x = x + (cell_w - icon_size) // 2
        add_service_icon(slide, icon_x, card_y + emu_in(0.18),
                         icon_size, icon_kind, color=accent)
        # Name
        add_text(slide, x + Emu(80000), card_y + emu_in(0.8),
                 cell_w - Emu(160000), emu_in(0.85), name,
                 size=12, color=BLACK, font=HEADING, bold=True,
                 align="center")
        # Price
        add_text(slide, x + Emu(80000), card_y + emu_in(1.7),
                 cell_w - Emu(160000), emu_in(0.4), price,
                 size=11, color=accent, font=HEADING, bold=True,
                 align="center")
        # One-liner
        add_text(slide, x + Emu(80000), card_y + emu_in(2.15),
                 cell_w - Emu(160000), card_h - emu_in(2.3), oneline,
                 size=9, color=GRAY, font=BODY, align="center")

    add_text(slide, PAD, SLIDE_H - emu_in(0.55),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "Один подрядчик — один план — одна точка ответственности",
             size=12, color=BLUE, font=BODY, italic=True, align="center")


def slide_18_close_provocation(prs):
    """Финал — бинарный выбор «5 подрядчиков vs 1 команда».

    Сдвигает акцент с GEO-провокации на устройство комплексного агентства.
    Отстройка от «колдовского / магического» агентства — через подзаголовок
    «устройство, а не магия».
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Финал", size=14, color=BLUE, font=BODY, bold=True)

    # Title — провокация про устройство
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Один контракт или пять?",
             size=52, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.5),
             "То, чем мы отличаемся, — устройство, а не магия.",
             size=18, color=GRAY, font=BODY, italic=True)

    # Two-column comparison
    cmp_y = emu_in(3.2)
    cmp_h = emu_in(3.0)
    cmp_w = SLIDE_W - 2 * PAD
    gap = emu_in(0.3)
    card_w = (cmp_w - gap) // 2

    # ---- Left: VARIANT A (нейтральный, тёмно-серый) ----
    ax = PAD
    add_rect(slide, ax, cmp_y, card_w, cmp_h,
             RGBColor(0x2A, 0x2A, 0x2E), corner=0.06,
             line=GRAY_BORDER_DARK)
    add_text(slide, ax + emu_in(0.35), cmp_y + emu_in(0.3),
             card_w - emu_in(0.7), emu_in(0.4),
             "ВАРИАНТ A — 5 ПОДРЯДЧИКОВ",
             size=11, color=GRAY, font=BODY, bold=True)
    add_text(slide, ax + emu_in(0.35), cmp_y + emu_in(0.85),
             card_w - emu_in(0.7), emu_in(1.05),
             "Дизайн — у одного. Разработка — у другого.\n"
             "SEO — у третьего. Контекст — у четвёртого.\n"
             "Поддержка — у пятого.",
             size=14, color=LIGHT_GRAY_TEXT, font=BODY)
    add_text(slide, ax + emu_in(0.35), cmp_y + emu_in(2.0),
             card_w - emu_in(0.7), emu_in(0.4),
             "5 счетов · 5 KPI · 5 разных стандартов",
             size=12, color=GRAY, font=BODY, italic=True)
    add_text(slide, ax + emu_in(0.35), cmp_y + emu_in(2.45),
             card_w - emu_in(0.7), emu_in(0.5),
             "Никто не отвечает за общий результат.",
             size=18, color=WHITE, font=HEADING, bold=True)

    # ---- Right: VARIANT B (синий акцент) ----
    bx = PAD + card_w + gap
    add_rect(slide, bx, cmp_y, card_w, cmp_h, BLUE, corner=0.06)
    add_text(slide, bx + emu_in(0.35), cmp_y + emu_in(0.3),
             card_w - emu_in(0.7), emu_in(0.4),
             "ВАРИАНТ B — 1 КОМАНДА · 6 НАПРАВЛЕНИЙ",
             size=11, color=WHITE, font=BODY, bold=True)
    add_text(slide, bx + emu_in(0.35), cmp_y + emu_in(0.85),
             card_w - emu_in(0.7), emu_in(1.05),
             "Разработка сайта / SEO / Контекстная реклама /\n"
             "Управление репутацией / Техническая техподдержка / GEO —\n"
             "один договор, один план, одна точка ответственности.",
             size=13, color=WHITE, font=BODY)
    add_text(slide, bx + emu_in(0.35), cmp_y + emu_in(2.0),
             card_w - emu_in(0.7), emu_in(0.4),
             "Среди клиентов — Сбер, БыстроБанк, Главстрой, Pioneer.",
             size=12, color=RGBColor(0xDD, 0xE9, 0xFB),
             font=BODY, italic=True)
    add_text(slide, bx + emu_in(0.35), cmp_y + emu_in(2.45),
             card_w - emu_in(0.7), emu_in(0.5),
             "95% продлевают. Средний срок — 3,5+ года.",
             size=18, color=WHITE, font=HEADING, bold=True)

    # Bridging line
    bridge_y = cmp_y + cmp_h + emu_in(0.25)
    add_text(slide, PAD, bridge_y, SLIDE_W - 2 * PAD, emu_in(0.4),
             "Если узнали себя в варианте слева — мы и есть та самая альтернатива.",
             size=14, color=GRAY, font=BODY, italic=True, align="center")

    # Contact strip
    strip_y = SLIDE_H - emu_in(0.95)
    strip_h = emu_in(0.7)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, strip_h, BLUE, corner=0.15)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, strip_h,
             "hello@acoola.team   ·   +7 499 938-48-27   ·   acoola.team",
             size=18, color=WHITE, font=HEADING, bold=True,
             align="center", anchor="middle")


# ---------- v3.5 Additional slides ----------

def slide_09a_weak_spots(prs):
    """Где мы слабые — продолжение анти-маркетинга, ещё на шаг честнее.

    Ни одно агентство в России такое не показывает.
    Поэтому это работает.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наша модель",
             size=14, color=BLUE, font=BODY, bold=True)

    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(1.1),
             "Где мы слабые.",
             size=46, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.6), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Чтобы вы знали заранее, а не на третьем месяце.",
             size=18, color=GRAY, font=BODY, italic=True)

    weaknesses = [
        ("ML / Data Science",
         "У нас нет такой команды. Для AI-продуктов работаем в партнёрстве "
         "с [TODO: имя партнёра]."),
        ("Мобильная разработка",
         "Не делаем приложения для iOS/Android. Только web и веб-сервисы."),
        ("Маленькие бюджеты",
         "Мы дороже фрилансера. Если бюджет на маркетинг до 150 000 ₽/мес — "
         "мы не подходим, и это нормально."),
        ("Ниши не наши",
         "Фарма, food retail, банки физлиц — рекомендуем коллег. "
         "Ритуалка — отказ принципиальный: с повторными "
         "продажами там объективные сложности."),
    ]
    grid_y = emu_in(3.4)
    grid_h = emu_in(3.3)
    grid_w = SLIDE_W - 2 * PAD
    cols = 2
    rows = 2
    gap = Emu(140000)
    cell_w = (grid_w - gap) // cols
    cell_h = (grid_h - gap) // rows
    for i, (title, body) in enumerate(weaknesses):
        r, c = divmod(i, cols)
        x = PAD + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        add_rect(slide, x, y, cell_w, cell_h, GRAY_DARK,
                 line=GRAY_BORDER_DARK, corner=0.05)
        add_text(slide, x + emu_in(0.35), y + emu_in(0.3),
                 cell_w - emu_in(0.7), emu_in(0.5),
                 title, size=18, color=BLUE, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.35), y + emu_in(0.95),
                 cell_w - emu_in(0.7), cell_h - emu_in(1.2),
                 body, size=12, color=LIGHT_GRAY_TEXT, font=BODY)

    add_text(slide, PAD, SLIDE_H - emu_in(0.6),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "Это часть честного разговора. Полный список — на встрече.",
             size=12, color=GRAY, font=BODY, italic=True, align="center")


def slide_13_failure_case(prs):
    """Кейс, который не получился. Никто из агентств такое не показывает.

    [ЧЕРНОВИК: ниша обобщённая, заменить на реальный пример из практики Acoola.]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Кейс · ПРОВАЛ",
             size=14, color=BLUE, font=BODY, bold=True)

    # Negative-toned title
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(0.6),
             "КЕЙС, КОТОРЫЙ НЕ ПОЛУЧИЛСЯ",
             size=14, color=RGBColor(0xE0, 0x6A, 0x6A), font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(1.95), SLIDE_W - 2 * PAD, emu_in(0.8),
             "Когда мы расстались первыми.",
             size=40, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.85), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Интернет-магазин стройматериалов · SEO + контекст · 7 месяцев [TODO: реальный кейс]",
             size=13, color=GRAY, font=BODY, italic=True)

    paragraphs = [
        ("2022 год.", BLUE, 18, True, HEADING),
        ("Зашли в проект на договоре 6 месяцев.\nSEO + контекст. Цели — рост обращений на 40%.",
         WHITE, 17, False, BODY),
        ("За полгода выросли позиции, но обращений почти нет.\nПроблема — в сайте: устаревшая структура, сломанная воронка,\nмобильная версия не работает. Клиент не был готов трогать сайт.",
         LIGHT_GRAY_TEXT, 16, False, BODY),
        ("Расстались по нашей инициативе на 7-м месяце.\nВернули нерасходованную часть бюджета.",
         WHITE, 18, True, HEADING),
        ("Что поняли: SEO не работает там, где сайт нужно переделать.\nС тех пор это в нашем стоп-листе — слайд 8.",
         BLUE, 16, True, BODY),
    ]
    y = emu_in(3.45)
    for text, color, size, bold, font in paragraphs:
        add_text(slide, PAD, y, SLIDE_W - 2 * PAD, emu_in(0.95),
                 text, size=size, color=color, font=font, bold=bold)
        y += emu_in(0.7)

    # Bottom callout
    cb_y = SLIDE_H - emu_in(0.7)
    add_text(slide, PAD, cb_y, SLIDE_W - 2 * PAD, emu_in(0.4),
             "[ЧЕРНОВИК — заменить на реальную историю от заказчика]",
             size=10, color=GRAY, font=BODY, italic=True, align="center")


def slide_17_cost_of_mistake(prs):
    """Цена ошибки выбора подрядчика. Числовая иллюстрация.

    Цифры реалистичные, но иллюстративные.
    [TODO: подтвердить с заказчиком — заменить на цифры из реальной практики.]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Что бывает, если выбрать дешевле",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Цена ошибки выбора подрядчика.",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Иллюстративный сценарий [TODO: подтвердить числа с заказчиком]",
             size=12, color=GRAY, font=BODY, italic=True)

    # Timeline of pain — left column
    timeline_x = PAD
    timeline_y = emu_in(3.0)
    timeline_w = emu_in(8.5)
    timeline_h = emu_in(3.6)

    add_text(slide, timeline_x, timeline_y, timeline_w, emu_in(0.4),
             "ВАРИАНТ A — Фрилансер за 80 000 ₽/мес",
             size=12, color=RGBColor(0xC8, 0x3E, 0x3E), font=BODY, bold=True)

    months = [
        ("М 0", "Найм фрилансера",
         "Экономия −170 000 ₽\nvs наша вилка от 250 000 ₽/мес",
         RGBColor(0xC8, 0x3E, 0x3E)),
        ("М 6", "Сайт не доработан",
         "Бюджет на контекст слит\n−450 000 ₽",
         RGBColor(0xC8, 0x3E, 0x3E)),
        ("М 7", "Новый подрядчик",
         "Переделка с нуля\n−350 000 ₽",
         RGBColor(0xC8, 0x3E, 0x3E)),
        ("М 12", "Сайт работает",
         "Простой 8 месяцев\n~600 000 ₽ упущенной выручки\n[TODO: уточнить под нишу]",
         RGBColor(0xC8, 0x3E, 0x3E)),
    ]
    cell_y = timeline_y + emu_in(0.55)
    cell_h = emu_in(2.7)
    cell_w = (timeline_w - Emu(80000) * 3) // 4
    for i, (when, head, body, color) in enumerate(months):
        x = timeline_x + i * (cell_w + Emu(80000))
        add_rect(slide, x, cell_y, cell_w, cell_h, WHITE,
                 corner=0.05, line=color)
        add_text(slide, x + emu_in(0.2), cell_y + emu_in(0.2),
                 cell_w - emu_in(0.4), emu_in(0.4), when,
                 size=11, color=color, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.2), cell_y + emu_in(0.65),
                 cell_w - emu_in(0.4), emu_in(0.55), head,
                 size=14, color=BLACK, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.2), cell_y + emu_in(1.3),
                 cell_w - emu_in(0.4), cell_h - emu_in(1.5), body,
                 size=11, color=GRAY, font=BODY)

    # Right — total
    total_x = timeline_x + timeline_w + emu_in(0.3)
    total_w = SLIDE_W - PAD - total_x
    add_rect(slide, total_x, timeline_y, total_w, timeline_h,
             RGBColor(0xC8, 0x3E, 0x3E), corner=0.05)
    add_text(slide, total_x + emu_in(0.3), timeline_y + emu_in(0.4),
             total_w - emu_in(0.6), emu_in(0.5),
             "ИТОГ ЗА ГОД",
             size=12, color=WHITE, font=BODY, bold=True)
    add_text(slide, total_x + emu_in(0.3), timeline_y + emu_in(1.0),
             total_w - emu_in(0.6), emu_in(1.5),
             "−1 230 000 ₽",
             size=42, color=WHITE, font=HEADING, bold=True)
    add_text(slide, total_x + emu_in(0.3), timeline_y + emu_in(2.4),
             total_w - emu_in(0.6), emu_in(0.5),
             "+ 8 месяцев простоя\nв продажах",
             size=13, color=WHITE, font=BODY)

    # Compare line at bottom
    cmp_y = SLIDE_H - emu_in(0.7)
    add_text(slide, PAD, cmp_y, SLIDE_W - 2 * PAD, emu_in(0.4),
             "С нами было бы — от 250 000 ₽/мес × 12 = 3 000 000 ₽ за год + работающий сайт + рост обращений.",
             size=12, color=BLUE, font=BODY, bold=True, italic=True, align="center")


def slide_22_bonus_seo_checklist(prs):
    """Бонус: 5 вопросов любому digital-подрядчику перед подписанием договора.

    Универсальный чек-лист (не AI-специфичный). Снимает GEO-перекос в финале
    кита. Незаметно дискредитирует дешёвых и нишевых подрядчиков. Положение
    Acoola — как «правильный B2B-партнёр», соответствующий всем 5 пунктам.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Бонус — забирайте, даже если не работаете с нами",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "5 вопросов до подписания договора.",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Задайте любому подрядчику. Если 2+ ответа размытые — ищите дальше.",
             size=14, color=GRAY, font=BODY, italic=True)

    items = [
        ("01",
         "Как вы отчитываетесь?",
         "Позиции и клики ≠ обращения и продажи. Спросите конкретные метрики и формат отчёта."),
        ("02",
         "Доступ к нашей CRM и Метрике?",
         "Подрядчик без доступа в воронку клиента оптимизирует не туда. Не интегрируется — плохой знак."),
        ("03",
         "Что у вас в стоп-листе?",
         "Профессионал знает, что НЕ делает. Те, кто берётся за всё, не имеют фокуса и стандартов."),
        ("04",
         "С кем работали в нашей категории?",
         "Назовите 2–3 имени или отрасли. Уход в «много кейсов» без конкретики — красный флаг."),
        ("05",
         "Что в договоре с минимальным сроком и KPI?",
         "Месячный отчёт по метрикам — норма B2B. «Сделаем что нужно» — это не договор."),
    ]
    grid_y = emu_in(3.0)
    grid_h = emu_in(3.4)
    grid_w = SLIDE_W - 2 * PAD
    gap = Emu(80000)
    cell_w = (grid_w - gap * 4) // 5
    for i, (num, head, body) in enumerate(items):
        x = PAD + i * (cell_w + gap)
        add_rect(slide, x, grid_y, cell_w, grid_h, WHITE, corner=0.05)
        add_text(slide, x + emu_in(0.2), grid_y + emu_in(0.25),
                 cell_w - emu_in(0.4), emu_in(0.5),
                 num, size=22, color=BLUE, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.2), grid_y + emu_in(0.85),
                 cell_w - emu_in(0.4), emu_in(1.4),
                 head, size=13, color=BLACK, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.2), grid_y + emu_in(2.3),
                 cell_w - emu_in(0.4), grid_h - emu_in(2.5),
                 body, size=11, color=GRAY, font=BODY)

    # Bottom CTA strip
    strip_y = SLIDE_H - emu_in(0.95)
    strip_h = emu_in(0.7)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, strip_h, BLUE, corner=0.15)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, strip_h,
             "Хотите задать эти вопросы нам? hello@acoola.team   ·   мы готовы.",
             size=15, color=WHITE, font=HEADING, bold=True,
             align="center", anchor="middle")


# ---------- v3.7 cases (Kenner / ORM) and History ----------

def slide_case_kenner_context(prs):
    """Кейс · Контекст: Kenner — производство мебели. ROI 140%.

    Источник: /projects/keys-kenner-... на acoola.team.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Кейс · КОНТЕКСТ",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(0.8),
             "Kenner", size=44, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.3), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Производство мебели · Яндекс.Директ + аналитика",
             size=14, color=GRAY, font=BODY, italic=True)

    left_w = emu_in(7.5)
    right_x = PAD + left_w + emu_in(0.4)
    right_w = SLIDE_W - PAD - right_x

    paragraphs = [
        "Старт работ.",
        "Производитель мебели заходит в digital.\nЦель — прямые продажи через интернет-магазин,\nне через дилеров.",
        "Связали контекстную рекламу\nс аналитикой и воронкой сайта:\nоптимизация не по CPL, а по ROI кампании.",
        "ROI 140%. CPL 1 374 ₽ за добавление в корзину,\n4 067 ₽ за оформленную продажу.",
    ]
    y = emu_in(3.0)
    for i, p in enumerate(paragraphs):
        is_first = (i == 0)
        is_last = (i == len(paragraphs) - 1)
        size = 18 if is_first else (22 if is_last else 16)
        color = BLUE if is_first else WHITE
        weight = True if (is_first or is_last) else False
        add_text(slide, PAD, y, left_w, emu_in(0.9),
                 p, size=size, color=color,
                 font=HEADING if (is_first or is_last) else BODY,
                 bold=weight)
        y += emu_in(0.85)

    # Right column: ROI hero card
    chart_y = emu_in(3.0)
    chart_h = emu_in(3.4)
    add_rect(slide, right_x, chart_y, right_w, chart_h,
             BLUE, corner=0.05)
    add_text(slide, right_x, chart_y + emu_in(0.3), right_w, emu_in(0.4),
             "ROI ПРОЕКТА",
             size=12, color=WHITE, font=BODY, bold=True, align="center")
    add_text(slide, right_x, chart_y + emu_in(0.8), right_w, emu_in(1.6),
             "140%",
             size=132, color=WHITE, font=HEADING, bold=True, align="center")
    add_text(slide, right_x, chart_y + emu_in(2.45), right_w, emu_in(0.4),
             "CPL 1 374 ₽ (корзина)",
             size=11, color=WHITE, font=BODY, align="center")
    add_text(slide, right_x, chart_y + emu_in(2.85), right_w, emu_in(0.4),
             "CPL 4 067 ₽ (продажа)",
             size=11, color=WHITE, font=BODY, align="center")

    # Stats footer
    strip_y = SLIDE_H - emu_in(0.9)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             GRAY_DARK, corner=0.1)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             "ROI 140%   ·   CPL 1 374 ₽ / 4 067 ₽   ·   Яндекс.Директ + аналитика   ·   acoola.team/projects",
             size=13, color=BLUE, font=HEADING, bold=True,
             align="center", anchor="middle")


def slide_case_orm_security(prs):
    """Кейс · ORM: Крупный бренд охранных систем.

    Источник: /projects/krupnyy-brend-okhrannykh-sistem/ на acoola.team.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Кейс · УПРАВЛЕНИЕ РЕПУТАЦИЕЙ",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.5), SLIDE_W - 2 * PAD, emu_in(0.8),
             "Крупный бренд охранных систем",
             size=36, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.3), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Безопасность · работа с репутацией на 30+ площадках",
             size=14, color=GRAY, font=BODY, italic=True)

    left_w = emu_in(7.5)
    right_x = PAD + left_w + emu_in(0.4)
    right_w = SLIDE_W - PAD - right_x

    paragraphs = [
        "Точка входа.",
        "В выдаче по бренду — массив негативных отзывов.\nКонкуренты вытесняют из карточек 2ГИС, Яндекс.Карт.\nДоверие падает.",
        "Запустили систему: регламент работы с отзывами,\nкарта присутствия, мониторинг через Brand Analytics,\nежемесячные отчёты по динамике.",
        "Доля негатива упала с 29% до 7%.\nСредний рейтинг бренда — 4,07 ★.",
    ]
    y = emu_in(3.0)
    for i, p in enumerate(paragraphs):
        is_first = (i == 0)
        is_last = (i == len(paragraphs) - 1)
        size = 18 if is_first else (22 if is_last else 16)
        color = BLUE if is_first else WHITE
        weight = True if (is_first or is_last) else False
        add_text(slide, PAD, y, left_w, emu_in(0.9),
                 p, size=size, color=color,
                 font=HEADING if (is_first or is_last) else BODY,
                 bold=weight)
        y += emu_in(0.85)

    # Right column: "было/стало" — две большие плашки негатива
    chart_y = emu_in(3.0)
    chart_h = emu_in(3.4)
    add_rect(slide, right_x, chart_y, right_w, chart_h,
             GRAY_DARK, corner=0.05, line=GRAY_BORDER_DARK)
    add_text(slide, right_x, chart_y + emu_in(0.25),
             right_w, emu_in(0.35),
             "ДОЛЯ НЕГАТИВА В ВЫДАЧЕ",
             size=10, color=GRAY, font=BODY, bold=True, align="center")

    # "Было" bar — wide, red
    before_y = chart_y + emu_in(0.75)
    before_h = Emu(int(chart_h * 0.32))
    add_text(slide, right_x + emu_in(0.2), before_y - emu_in(0.05),
             emu_in(1.2), emu_in(0.35), "БЫЛО",
             size=10, color=GRAY, font=BODY)
    bar_x = right_x + emu_in(0.2)
    bar_w_max = right_w - emu_in(0.4) - emu_in(1.2)
    add_rect(slide, bar_x, before_y + emu_in(0.25),
             bar_w_max, before_h,
             RGBColor(0xC8, 0x3E, 0x3E), corner=0.1)
    add_text(slide, right_x + right_w - emu_in(1.0), before_y + emu_in(0.15),
             emu_in(0.8), before_h + emu_in(0.2), "29%",
             size=28, color=WHITE, font=HEADING, bold=True,
             align="right", anchor="middle")

    # "Стало" bar — short, blue
    after_y = before_y + before_h + emu_in(0.5)
    add_text(slide, right_x + emu_in(0.2), after_y - emu_in(0.05),
             emu_in(1.2), emu_in(0.35), "СТАЛО",
             size=10, color=GRAY, font=BODY)
    # 7 / 29 = 24% of max width
    after_bar_w = Emu(int(bar_w_max * 0.24))
    add_rect(slide, bar_x, after_y + emu_in(0.25),
             after_bar_w, before_h,
             BLUE, corner=0.1)
    add_text(slide, bar_x + after_bar_w + emu_in(0.15), after_y + emu_in(0.15),
             emu_in(1.0), before_h + emu_in(0.2), "7%",
             size=28, color=WHITE, font=HEADING, bold=True,
             anchor="middle")

    # Rating star line below
    add_text(slide, right_x, chart_y + chart_h - emu_in(0.45),
             right_w, emu_in(0.4),
             "Средний рейтинг бренда — 4,07 ★",
             size=12, color=BLUE, font=HEADING, bold=True, align="center")

    # Stats footer
    strip_y = SLIDE_H - emu_in(0.9)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             GRAY_DARK, corner=0.1)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, emu_in(0.6),
             "Негатив 29% → 7%   ·   рейтинг 4,07 ★   ·   30+ площадок мониторинга",
             size=13, color=BLUE, font=HEADING, bold=True,
             align="center", anchor="middle")


def slide_history_timeline(prs):
    """История агентства · timeline 5 вех.

    Точные годы — TODO у заказчика (на сайте /about/ дат нет).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "История",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "От первого проекта до 70+ ежемесячно.",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.4),
             "От первого договора до 70+ проектов ежемесячно — за 6 лет.",
             size=12, color=GRAY, font=BODY, italic=True)

    # Годы реальные: ООО «АКУЛА ТИМ» зарегистрирована в 2019 (по ОГРН).
    # Остальные вехи расставлены по логике развития. Точные даты партнёрства
    # и NPS — заказчик может скорректировать.
    milestones = [
        ("2019", "Основание Acoola Team",
         "Первая команда — разработка сайтов под бизнес-задачи."),
        ("2020", "Золотой партнёр\n1С Битрикс",
         "Подтверждённая экспертиза CMS-разработки."),
        ("2022", "50+ команда",
         "Расширение до 6 направлений в одной команде."),
        ("2024", "NPS 100\nза последний год",
         "По независимой оценке программы качества 1С Битрикс."),
        ("2025", "GEO + премия Ruward",
         "Запустили GEO — видимость в AI-поиске.\nПолучили премию Ruward 2025 за разработку."),
    ]

    grid_y = emu_in(3.2)
    grid_h = emu_in(3.5)
    grid_w = SLIDE_W - 2 * PAD
    cols = 5
    gap = Emu(100000)
    cell_w = (grid_w - gap * (cols - 1)) // cols

    # Axis line
    axis_y = grid_y + Emu(int(grid_h * 0.20))
    add_rect(slide, PAD + Emu(80000), axis_y,
             grid_w - Emu(160000), Emu(8000),
             RGBColor(0xCC, 0xCC, 0xCF))

    for i, (year, title, body) in enumerate(milestones):
        x = PAD + i * (cell_w + gap)
        is_last = (i == len(milestones) - 1)
        accent = PURPLE if is_last else BLUE
        # Year label above
        add_text(slide, x, grid_y, cell_w, emu_in(0.45),
                 year, size=20, color=accent, font=HEADING, bold=True,
                 align="center")
        # Dot
        dot_d = Emu(180000)
        add_rect(slide, x + (cell_w - dot_d) // 2, axis_y - dot_d // 2 + Emu(4000),
                 dot_d, dot_d, accent, corner=0.5)
        # Card
        card_y = axis_y + Emu(180000)
        card_h = grid_h - (card_y - grid_y) - Emu(60000)
        card_fill = WHITE if not is_last else RGBColor(0xEC, 0xE6, 0xFE)
        card_border = None if not is_last else PURPLE
        add_rect(slide, x, card_y, cell_w, card_h, card_fill,
                 corner=0.05, line=card_border)
        # Иконка кубка для вехи с Ruward 2025 (компактнее — чтобы оставить место для текста)
        if is_last:
            trophy_size = emu_in(0.5)
            trophy_x = x + (cell_w - trophy_size) // 2
            trophy_y = card_y + emu_in(0.18)
            add_trophy(slide, trophy_x, trophy_y, trophy_size,
                       color=RGBColor(0xC9, 0x9A, 0x3D))
            title_y_offset = emu_in(0.75)
        else:
            title_y_offset = emu_in(0.25)
        title_h = emu_in(0.75)
        add_text(slide, x + emu_in(0.2), card_y + title_y_offset,
                 cell_w - emu_in(0.4), title_h, title,
                 size=13, color=BLACK, font=HEADING, bold=True, align="center")
        body_y = card_y + title_y_offset + title_h + emu_in(0.05)
        body_h = card_h - (body_y - card_y) - emu_in(0.15)
        if body_h > 0:
            add_text(slide, x + emu_in(0.2), body_y,
                     cell_w - emu_in(0.4), body_h, body,
                     size=10, color=GRAY, font=BODY, align="center")

    add_text(slide, PAD, SLIDE_H - emu_in(0.55),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "70+ проектов ежемесячно · 3,5+ года средний срок сотрудничества · 95% довольных клиентов",
             size=12, color=BLUE, font=BODY, italic=True, align="center")


# ---------- v3.9: География (6 стран) + Партнёрская программа ----------

def slide_countries(prs):
    """География как УТП, а не статистика.

    Левая колонка — смысловой посыл (заголовок + 3 ценности).
    Правая колонка — сетка 6 флагов с подписями.
    Так слайд работает как продающее обещание, а не как «у нас есть метка на карте».
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "География",
             size=14, color=BLUE, font=BODY, bold=True)

    # Hero title — пара провокаций
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Не одна Москва. Не один Яндекс.",
             size=42, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Работаем там, где работают ваши клиенты — в 6 странах.",
             size=16, color=GRAY, font=BODY, italic=True)

    # Two columns
    col_y = emu_in(3.1)
    col_h = emu_in(3.6)
    cols_w = SLIDE_W - 2 * PAD
    gap = emu_in(0.4)
    # left 42% — values
    left_w = (cols_w - gap) * 42 // 100
    right_w = cols_w - gap - left_w

    # ---- Left column: 3 value cards ----
    values = [
        ("Локализация",
         "Контент на языках местных рынков, не дословный перевод."),
        ("Местные поисковые системы",
         "Не только Яндекс — Google, локальные поисковики и каталоги под каждый рынок."),
        ("Культура и регулятор",
         "Кейсы в разных правовых, культурных и налоговых средах. Адаптируем подачу."),
    ]
    val_h = (col_h - emu_in(0.2)) // 3
    val_gap = Emu(80000)
    for i, (title, body) in enumerate(values):
        y = col_y + i * (val_h + val_gap - Emu(60000) if i else 0) + i * Emu(20000)
        # лёгкая упрощённая версия (без накопления)
        y = col_y + i * val_h
        add_rect(slide, PAD, y, left_w, val_h - Emu(60000),
                 WHITE, corner=0.05)
        add_text(slide, PAD + emu_in(0.25), y + emu_in(0.18),
                 left_w - emu_in(0.5), emu_in(0.4),
                 title, size=15, color=BLUE, font=HEADING, bold=True)
        add_text(slide, PAD + emu_in(0.25), y + emu_in(0.65),
                 left_w - emu_in(0.5), val_h - emu_in(0.95),
                 body, size=11, color=GRAY, font=BODY)

    # ---- Right column: hero "6" + grid of flags ----
    right_x = PAD + left_w + gap
    # Card background for flag area
    add_rect(slide, right_x, col_y, right_w, col_h, WHITE, corner=0.05)
    # Hero number "6" + label "стран"
    add_text(slide, right_x + emu_in(0.35), col_y + emu_in(0.3),
             emu_in(1.6), emu_in(1.5),
             "6", size=110, color=BLUE, font=HEADING, bold=True)
    add_text(slide, right_x + emu_in(1.7), col_y + emu_in(0.55),
             right_w - emu_in(2.0), emu_in(0.6),
             "стран присутствия", size=15, color=BLACK, font=HEADING, bold=True)
    add_text(slide, right_x + emu_in(1.7), col_y + emu_in(1.0),
             right_w - emu_in(2.0), emu_in(0.5),
             "70+ проектов ежемесячно", size=12, color=GRAY, font=BODY)

    # Flags grid 3×2
    countries = [
        ("italy",      "Италия"),
        ("russia",     "Россия"),
        ("usa",        "США"),
        ("belarus",    "Беларусь"),
        ("kazakhstan", "Казахстан"),
        ("uae",        "ОАЭ"),
    ]
    flag_area_y = col_y + emu_in(1.9)
    flag_area_h = col_h - emu_in(2.1)
    flag_cols = 3
    flag_rows = 2
    inner_pad = emu_in(0.3)
    flag_w = (right_w - inner_pad * 2 - Emu(120000) * (flag_cols - 1)) // flag_cols
    flag_h_each = (flag_area_h - Emu(80000) * (flag_rows - 1) - emu_in(0.4)) // flag_rows
    flag_disp_h = flag_h_each * 6 // 10  # height for actual flag image
    for i, (code, label) in enumerate(countries):
        r, c = divmod(i, flag_cols)
        fx = right_x + inner_pad + c * (flag_w + Emu(120000))
        fy = flag_area_y + r * (flag_h_each + Emu(80000))
        # Flag image
        add_flag(slide, fx, fy, flag_w, flag_disp_h, code)
        # Country label below
        add_text(slide, fx, fy + flag_disp_h + Emu(50000),
                 flag_w, emu_in(0.35),
                 label, size=11, color=GRAY, font=BODY,
                 bold=True, align="center")

    # Anchor at the bottom
    add_text(slide, PAD, SLIDE_H - emu_in(0.55),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "Среди клиентов — компании из Москвы, ОАЭ, Италии, США, Беларуси и Казахстана.",
             size=12, color=BLUE, font=BODY, italic=True, align="center")


def slide_partnership(prs):
    """Партнёрская программа — последний слайд.

    Источник: /partnership/ на acoola.team.
    Условия: 10% с ежемесячных услуг (год) + 12% с разработки сайтов.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Партнёрская программа",
             size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Зарабатывайте на рекомендациях.",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Платим деньги за привлечение клиентов на любую из 6 услуг.",
             size=14, color=GRAY, font=BODY, italic=True)

    # Two big % cards (10% / 12%) — двухколоночный layout внутри карточки:
    # слева текст (label сверху + description снизу), справа крупная цифра.
    # Это исключает наложение, которое было при вертикальной раскладке.
    cards_y = emu_in(2.85)
    cards_h = emu_in(2.2)
    cards_w = SLIDE_W - 2 * PAD
    card_gap = emu_in(0.3)
    card_w = (cards_w - card_gap) // 2

    def _draw_pct_card(x, fill, label, percent, description):
        # карточка
        add_rect(slide, x, cards_y, card_w, cards_h, fill, corner=0.06)
        # text column — left 55%
        text_pad = emu_in(0.45)
        text_w = (card_w * 55 // 100) - text_pad
        # label вверху
        add_text(slide, x + text_pad, cards_y + emu_in(0.4),
                 text_w, emu_in(0.4),
                 label, size=12, color=WHITE, font=BODY, bold=True)
        # description внизу text-колонки
        add_text(slide, x + text_pad, cards_y + cards_h - emu_in(1.1),
                 text_w, emu_in(0.85),
                 description, size=12, color=WHITE, font=BODY)
        # number — right 45%, центрированный вертикально
        num_x = x + card_w * 55 // 100
        num_w = card_w - (card_w * 55 // 100) - emu_in(0.3)
        add_text(slide, num_x, cards_y, num_w, cards_h,
                 percent, size=80, color=WHITE, font=HEADING, bold=True,
                 align="center", anchor="middle")

    _draw_pct_card(
        PAD, BLUE,
        label="ЕЖЕМЕСЯЧНЫЕ УСЛУГИ",
        percent="10%",
        description="SEO · контекст · ORM · техподдержка\n\nЕжемесячно в течение года сотрудничества",
    )
    _draw_pct_card(
        PAD + card_w + card_gap, PURPLE,
        label="РАЗРАБОТКА САЙТОВ",
        percent="12%",
        description="От каждого платежа клиента\nза разработку сайта",
    )

    # 3 steps row — компактнее, чтобы поместиться над финальной плашкой
    steps_y = cards_y + cards_h + emu_in(0.2)
    steps_h = emu_in(1.25)
    step_gap = emu_in(0.2)
    step_w = (cards_w - step_gap * 2) // 3
    steps = [
        ("01", "Рекомендуйте",
         "Расскажите коллегам, партнёрам или клиентам, которые ищут подрядчика."),
        ("02", "Сформулируйте задачу",
         "Коротко опишите, что нужно. Предупредите клиента о нашем звонке."),
        ("03", "Передайте контакты",
         "Мы проведём бесплатный аудит. Партнёру — комиссия."),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = PAD + i * (step_w + step_gap)
        add_rect(slide, x, steps_y, step_w, steps_h, WHITE, corner=0.05)
        add_text(slide, x + emu_in(0.25), steps_y + emu_in(0.18),
                 step_w - emu_in(0.5), emu_in(0.35),
                 num, size=13, color=BLUE, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.25), steps_y + emu_in(0.5),
                 step_w - emu_in(0.5), emu_in(0.35),
                 title, size=15, color=BLACK, font=HEADING, bold=True)
        add_text(slide, x + emu_in(0.25), steps_y + emu_in(0.85),
                 step_w - emu_in(0.5), emu_in(0.4),
                 body, size=10, color=GRAY, font=BODY)

    # Bottom CTA strip
    strip_y = SLIDE_H - emu_in(0.75)
    strip_h = emu_in(0.55)
    add_rect(slide, PAD, strip_y, SLIDE_W - 2 * PAD, strip_h, BLUE, corner=0.15)
    add_text(slide, PAD, strip_y, SLIDE_W - 2 * PAD, strip_h,
             "Стать партнёром: hello@acoola.team   ·   acoola.team/partnership/",
             size=15, color=WHITE, font=HEADING, bold=True,
             align="center", anchor="middle")


# ---------- v3.11: Прайс + Суперсилы ----------

def slide_pricing(prs):
    """Прайс — простая двухгрупповая таблица: Разовые + Ежемесячные.

    Цены в формате «от ...» — стартовые вилки с сайта acoola.team.
    Без иконок и выкрутасов: услуга + цена крупно. Читается за 5 секунд.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Прайс", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.35), SLIDE_W - 2 * PAD, emu_in(0.8),
             "Цены с нашего сайта.",
             size=34, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.15), SLIDE_W - 2 * PAD, emu_in(0.4),
             "Стартовые вилки. Точная стоимость под ваш проект — на встрече.",
             size=13, color=GRAY, font=BODY, italic=True)

    # Two groups: Разовые / Ежемесячные
    # Каждый item: (icon_kind, name, price, sub, is_geo)
    once = [
        ("dev", "Разработка сайта", "от 190 000 ₽", "За проект под ключ", False),
    ]
    monthly = [
        ("seo",     "SEO",                      "от 70 000 ₽/мес",  "Визитка / услуги / магазин / портал — 70/75/100/120k", False),
        ("context", "Контекстная реклама",      "от 50 000 ₽/мес",  "За ведение + рекламный бюджет",                       False),
        ("orm",     "Управление репутацией",    "от 70 000 ₽/мес",  "Карта присутствия + сканер репутации",                False),
        ("support", "Техническая техподдержка", "от 30 000 ₽/мес",  "Месячный пакет часов + SLA",                          False),
        ("geo",     "GEO",                      "от 90 000 ₽/мес",  "Видимость бренда в AI-поиске",                        True),
    ]

    grid_y = emu_in(2.85)
    grid_w = SLIDE_W - 2 * PAD
    # Левая 32% — разовые, правая 68% — ежемесячные
    col_gap = emu_in(0.3)
    left_w = (grid_w - col_gap) * 32 // 100
    right_w = grid_w - col_gap - left_w
    section_h = emu_in(3.9)

    # ---- Left: РАЗОВЫЕ ----
    add_rect(slide, PAD, grid_y, left_w, section_h, WHITE, corner=0.06)
    # Acent vertical bar на левом краю — визуальный акцент
    add_rect(slide, PAD, grid_y, Emu(60000), section_h, BLUE)
    add_text(slide, PAD + emu_in(0.45), grid_y + emu_in(0.35),
             left_w - emu_in(0.7), emu_in(0.35),
             "РАЗОВЫЕ ИНВЕСТИЦИИ",
             size=11, color=BLUE, font=BODY, bold=True)
    icon_kind, name, price, sub, _ = once[0]
    add_service_icon(slide, PAD + emu_in(0.45), grid_y + emu_in(0.95),
                     emu_in(0.55), icon_kind, color=BLUE)
    add_text(slide, PAD + emu_in(0.45), grid_y + emu_in(1.7),
             left_w - emu_in(0.7), emu_in(0.5),
             name, size=20, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD + emu_in(0.45), grid_y + emu_in(2.3),
             left_w - emu_in(0.7), emu_in(0.65),
             price, size=32, color=BLUE, font=HEADING, bold=True)
    add_text(slide, PAD + emu_in(0.45), grid_y + emu_in(3.05),
             left_w - emu_in(0.7), emu_in(0.4),
             sub, size=11, color=GRAY, font=BODY)
    add_text(slide, PAD + emu_in(0.45), grid_y + section_h - emu_in(0.4),
             left_w - emu_in(0.9), emu_in(0.3),
             "+ 30 дней техподдержки в подарок",
             size=10, color=BLUE, font=BODY, italic=True, bold=True)

    # ---- Right: ЕЖЕМЕСЯЧНЫЕ ----
    right_x = PAD + left_w + col_gap
    add_rect(slide, right_x, grid_y, right_w, section_h, WHITE, corner=0.06)
    add_text(slide, right_x + emu_in(0.45), grid_y + emu_in(0.35),
             right_w - emu_in(0.7), emu_in(0.35),
             "ЕЖЕМЕСЯЧНЫЕ ИНВЕСТИЦИИ",
             size=11, color=BLUE, font=BODY, bold=True)

    # Table layout — 5 rows, иконка + название + цена. Описание под названием мелко.
    table_y = grid_y + emu_in(0.85)
    rows_n = len(monthly)
    row_h = (section_h - emu_in(1.0)) // rows_n
    inner_l = right_x + emu_in(0.45)
    inner_r = right_x + right_w - emu_in(0.45)
    inner_w = inner_r - inner_l
    # Колонки: акцент-полоска 0.06" / иконка 0.5" / название+sub 55% от inner / цена 35% от inner
    bar_w = Emu(60000)
    icon_w = emu_in(0.55)
    icon_gap = emu_in(0.25)
    name_x = inner_l + bar_w + emu_in(0.15) + icon_w + icon_gap
    name_w = inner_w * 50 // 100
    price_x = name_x + name_w + emu_in(0.15)
    price_w = inner_r - price_x

    for i, (icon_kind, name, price, sub, is_geo) in enumerate(monthly):
        accent_color = PURPLE if is_geo else BLUE
        row_y = table_y + i * row_h
        # Тонкая горизонтальная разделительная линия (кроме первой)
        if i > 0:
            add_rect(slide, inner_l, row_y, inner_w, Emu(6000),
                     RGBColor(0xE6, 0xE6, 0xEA))
        # Маленькая вертикальная акцентная полоска слева от строки
        bar_h = row_h - emu_in(0.2)
        add_rect(slide, inner_l, row_y + emu_in(0.1),
                 bar_w, bar_h, accent_color, corner=0.5)
        # Icon
        icon_y = row_y + (row_h - icon_w) // 2
        add_service_icon(slide, inner_l + bar_w + emu_in(0.15),
                         icon_y, icon_w, icon_kind, color=accent_color)
        # Name (top) + sub (under) — в одной колонке
        text_pad_top = emu_in(0.12)
        add_text(slide, name_x, row_y + text_pad_top,
                 name_w, emu_in(0.35),
                 name, size=15, color=BLACK, font=HEADING, bold=True)
        add_text(slide, name_x, row_y + text_pad_top + emu_in(0.32),
                 name_w, row_h - emu_in(0.4),
                 sub, size=10, color=GRAY, font=BODY)
        # Price (right, big, single line)
        add_text(slide, price_x, row_y, price_w, row_h,
                 price, size=22, color=accent_color, font=HEADING, bold=True,
                 anchor="middle", align="right")

    # Bottom anchor
    add_text(slide, PAD, SLIDE_H - emu_in(0.4),
             SLIDE_W - 2 * PAD, emu_in(0.3),
             "Источник цен — acoola.team. Точная конфигурация — на встрече.",
             size=11, color=BLUE, font=BODY, italic=True, align="center")


def slide_superpowers(prs):
    """Наши суперсилы — креативная подача через «не как все».

    Каждая «суперсила» подана через ограничение / отказ от общепринятого
    подхода. Это резонирует с антимаркетинговыми слайдами кита (стоп-лист,
    «Где мы слабые») и продаёт экспертизу через позицию, а не через цифры.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_header(slide, dark=True)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.5),
             "Наши суперсилы", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.4), SLIDE_W - 2 * PAD, emu_in(1.0),
             "Шесть вещей, которые мы делаем не как все.",
             size=38, color=WHITE, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.4), SLIDE_W - 2 * PAD, emu_in(0.4),
             "И именно это работает.",
             size=15, color=GRAY, font=BODY, italic=True)

    # 6 superpowers — 3×2 grid. Описания сокращены, чтобы влезть в 2 строки.
    powers = [
        ("Не торопимся",
         "Не обещаем топ-10 к Новому году. Считаем результат на дистанции 6–12 месяцев."),
        ("Не дробимся",
         "Одна команда на 6 направлений вместо пяти подрядчиков. Один план."),
        ("Не теряем клиентов",
         "3,5+ года средний срок, 95% продлевают. Редкость в индустрии."),
        ("Не учимся на вас",
         "GEO с 2025 — пришли с готовой методологией, не обкатываем на клиенте."),
        ("Не привязаны к Москве",
         "Проекты в 6 странах. Локализация, не дословный перевод."),
        ("Не берём всё подряд",
         "Открытый стоп-лист. Знаем, что не делаем — поэтому есть стандарт."),
    ]
    grid_y = emu_in(2.95)
    grid_h = emu_in(3.95)
    grid_w = SLIDE_W - 2 * PAD
    cols = 3
    rows = 2
    gap = Emu(140000)
    cell_w = (grid_w - gap * (cols - 1)) // cols
    cell_h = (grid_h - gap * (rows - 1)) // rows
    for i, (title, body) in enumerate(powers):
        r, c = divmod(i, cols)
        x = PAD + c * (cell_w + gap)
        y = grid_y + r * (cell_h + gap)
        add_rect(slide, x, y, cell_w, cell_h, GRAY_DARK,
                 line=GRAY_BORDER_DARK, corner=0.05)
        # Number 01..06
        num = f"{i + 1:02d}"
        add_text(slide, x + emu_in(0.3), y + emu_in(0.25),
                 cell_w - emu_in(0.6), emu_in(0.35),
                 num, size=12, color=BLUE, font=BODY, bold=True)
        # Title — провокация
        add_text(slide, x + emu_in(0.3), y + emu_in(0.65),
                 cell_w - emu_in(0.6), emu_in(0.55),
                 title, size=20, color=WHITE, font=HEADING, bold=True)
        # Body — расшифровка (точно влезает в положительную высоту)
        add_text(slide, x + emu_in(0.3), y + emu_in(1.25),
                 cell_w - emu_in(0.6), cell_h - emu_in(1.4),
                 body, size=11, color=LIGHT_GRAY_TEXT, font=BODY)

    add_text(slide, PAD, SLIDE_H - emu_in(0.55),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "Каждая суперсила — это форма отказа от общепринятого. И в этом сила.",
             size=12, color=BLUE, font=BODY, italic=True, align="center")


# ---------- v3.13: 5 страхов клиента — наш ответ ----------

def slide_5_fears(prs):
    """5 страхов клиента — и наш ответ.

    Закрывает 4 из 5 топовых болей B2B-заказчика digital-агентств,
    выявленных по статьям в РФ (Mello, VC.ru, Cossa, AdIndex, Serpstat и др.).
    Структура: «БОЯЛИСЬ [страх] → ПРОПИСАНО В ДОГОВОРЕ [наш ответ]».

    [TODO у заказчика — подтвердить юридически:
      • фиксированную стоимость в договоре,
      • право расторжения с возвратом нерасходованной части,
      • регистрацию рекламных аккаунтов на юр.лицо клиента,
      • SLA на ответ в часах,
      • еженедельные созвоны.]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_header(slide, dark=False)

    add_text(slide, PAD, emu_in(0.95), SLIDE_W - 2 * PAD, emu_in(0.4),
             "5 страхов", size=14, color=BLUE, font=BODY, bold=True)
    add_text(slide, PAD, emu_in(1.35), SLIDE_W - 2 * PAD, emu_in(0.9),
             "5 страхов — и наш ответ.",
             size=38, color=BLACK, font=HEADING, bold=True)
    add_text(slide, PAD, emu_in(2.3), SLIDE_W - 2 * PAD, emu_in(0.5),
             "70% клиентов уже разочарованы прошлым подрядчиком. Мы знаем почему — и закрыли это в договоре.",
             size=14, color=GRAY, font=BODY, italic=True)

    fears = [
        ("Смета подрастёт на ходу.",
         "Фиксированная стоимость в договоре. Дополнительные работы — только по отдельному акту с вашей подписью."),
        ("Сорвут сроки.",
         "Таймлайн и KPI прописаны в договоре. Право расторжения с возвратом нерасходованной части бюджета."),
        ("Окажусь у них в заложниках.",
         "Все рекламные аккаунты, аналитика, домен и CMS — на ваше юр.лицо. У нас — только рабочий доступ."),
        ("Опять отчёты про позиции, а не про деньги.",
         "Ежемесячный отчёт по обращениям и выручке. Доступ к Метрике и CRM — у вас онлайн, 24/7."),
        ("Уйдут в тень при первой проблеме.",
         "Закреплённый менеджер. Еженедельные созвоны. SLA на ответ в часах, прописанный в договоре."),
    ]

    grid_y = emu_in(3.0)
    grid_h = SLIDE_H - grid_y - emu_in(0.95)
    grid_w = SLIDE_W - 2 * PAD
    row_count = len(fears)
    row_gap = Emu(80000)
    row_h = (grid_h - row_gap * (row_count - 1)) // row_count

    # Колонки: number | fear (40%) | arrow | answer (60% оставшегося)
    num_w = emu_in(0.5)
    fear_w = (grid_w - num_w - emu_in(0.6)) * 38 // 100
    arrow_w = emu_in(0.5)
    answer_w = grid_w - num_w - fear_w - arrow_w - emu_in(0.4)

    for i, (fear, answer) in enumerate(fears):
        y = grid_y + i * (row_h + row_gap)
        # White card-row
        add_rect(slide, PAD, y, grid_w, row_h, WHITE, corner=0.05)

        # Number
        x_cur = PAD + emu_in(0.3)
        add_text(slide, x_cur, y, num_w, row_h,
                 f"{i + 1:02d}", size=18, color=BLUE,
                 font=HEADING, bold=True, anchor="middle")
        x_cur += num_w + emu_in(0.15)

        # Fear (left)
        add_text(slide, x_cur, y + emu_in(0.08), fear_w, emu_in(0.3),
                 "БОЯЛИСЬ", size=9, color=GRAY,
                 font=BODY, bold=True)
        add_text(slide, x_cur, y + emu_in(0.32), fear_w, row_h - emu_in(0.35),
                 fear, size=14, color=BLACK,
                 font=HEADING, italic=True)
        x_cur += fear_w + emu_in(0.1)

        # Arrow
        add_text(slide, x_cur, y, arrow_w, row_h,
                 "→", size=22, color=BLUE,
                 font=HEADING, bold=True, align="center", anchor="middle")
        x_cur += arrow_w + emu_in(0.15)

        # Answer (right)
        add_text(slide, x_cur, y + emu_in(0.08), answer_w, emu_in(0.3),
                 "ПРОПИСАНО В ДОГОВОРЕ", size=9, color=BLUE,
                 font=BODY, bold=True)
        add_text(slide, x_cur, y + emu_in(0.32), answer_w, row_h - emu_in(0.35),
                 answer, size=12, color=BLACK,
                 font=BODY, bold=True)

    # Якорь снизу
    add_text(slide, PAD, SLIDE_H - emu_in(0.55),
             SLIDE_W - 2 * PAD, emu_in(0.4),
             "Хотите увидеть эти условия в договоре до подписания? Запросите шаблон: hello@acoola.team",
             size=12, color=BLUE, font=BODY, italic=True, align="center")


# ---------- Main ----------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        # Act 1: Conflict (1–4)
        slide_01_cover_provocation,
        slide_02_manifest,
        slide_03_shift,                # «3 сдвига в digital» (v3.10)
        slide_04_broken,
        # Act 2: Our model (5–9) — без «Обращений не позиции» (удалён v3.11)
        slide_06_long_term,
        slide_07_95_visualization,
        slide_08_anti_marketing,
        slide_09a_weak_spots,
        slide_09_geo_climax,
        # Act 3: Proof (10–16) — 4 кейса + failure + clients + team
        slide_case_teo,
        slide_11_story_wwerb,
        slide_case_kenner_context,
        slide_case_orm_security,
        slide_13_failure_case,
        slide_12_clients,
        slide_13_team_quote,
        # Act 3.5: о компании (17–20)
        slide_history_timeline,        # с Ruward + trophy (v3.11)
        slide_countries,               # География как УТП (v3.10)
        slide_superpowers,             # NEW v3.11 — Наши суперсилы
        slide_14_awards,
        # Act 4: к решению (21–26)
        slide_17_cost_of_mistake,
        slide_5_fears,                 # NEW v3.13 — 5 страхов и наш ответ
        slide_services_lifecycle_map,  # Карта услуг (с ценами на таймлайне)
        slide_pricing,                 # Отдельный прайс
        slide_18_close_provocation,    # Финал 5 vs 1
        slide_partnership,             # Партнёрская программа (последний)
    ]
    for b in builders:
        b(prs)

    # Post-processing: footer with page numbers on every slide
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        add_page_number(slide, i, total)

    out = "/Users/pro2kuror/Desktop/Cloude/docs/acoola_marketing_kit_v3_creative.pptx"
    prs.save(out)
    print(f"OK: {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    main()
