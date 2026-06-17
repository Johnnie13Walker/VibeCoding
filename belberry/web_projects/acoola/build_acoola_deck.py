"""
Acoola Team — KP deck for Совкомбанк (rich visual edition).

Visual identity: acoola.team
  bg #f2f2f4 / surface #ffffff / text #313131
  accent #3086fb / accent-strong #186fe0
  ink (deep) #0e1116

Logo: assets/acoola_logo.png
Photos: assets/photos/*.jpg (Unsplash)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Palette ----------
BG       = RGBColor(0xF2, 0xF2, 0xF4)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2 = RGBColor(0xF7, 0xF7, 0xFA)
ACCENT   = RGBColor(0x30, 0x86, 0xFB)
ACCENT_S = RGBColor(0x18, 0x6F, 0xE0)
INK      = RGBColor(0x0E, 0x11, 0x16)
TEXT     = RGBColor(0x31, 0x31, 0x31)
TEXT_S   = RGBColor(0x14, 0x14, 0x18)
MUTED    = RGBColor(0x6E, 0x6E, 0x78)
MUTED_2  = RGBColor(0xA0, 0xA0, 0xAA)
BORDER   = RGBColor(0xE3, 0xE3, 0xE8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.7)

FONT_DISPLAY = "Manrope"
FONT_BODY = "Gilroy"

LOGO = "assets/acoola_logo.png"
LOGO_WHITE = "assets/acoola_logo_white.png"
# Pre-cropped variants (no stretch)
PHOTO_BANKING_PORTRAIT = "assets/photos/banking_portrait.jpg"
PHOTO_TEAM_PORTRAIT = "assets/photos/team_portrait.jpg"
PHOTO_HANDSHAKE_169 = "assets/photos/handshake_169.jpg"
PHOTO_MEETING_169 = "assets/photos/meeting_169.jpg"
PHOTO_CODE_169 = "assets/photos/code_169.jpg"
PHOTO_DASHBOARD_169 = "assets/photos/hero_dashboard_169.jpg"
PHOTO_ANALYTICS_WIDE = "assets/photos/analytics_wide.jpg"
PHOTO_CHART_WIDE = "assets/photos/chart_wide.jpg"


# ============== Primitives ==============

def set_slide_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14, bold=False,
             color=TEXT, align="left", spacing=0, line_spacing=1.3, anchor="top"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if spacing:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return tb


def add_eyebrow(slide, x, y, text, color=MUTED):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.07),
                                  Inches(0.08), Inches(0.08))
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.shadow.inherit = False
    tb = slide.shapes.add_textbox(x + Inches(0.2), y, Inches(8), Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", "300")
    return tb


def add_card(slide, x, y, w, h, *, fill=SURFACE, border=BORDER, line_w=0.75,
             corner=0.04, shadow=False):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = corner
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(line_w)
    if not shadow:
        card.shadow.inherit = False
    return card


def add_image_card(slide, x, y, w, h, image_path, *, corner_in=0.18):
    """Place an image inside a rounded mask. Approximation: round-rect frame on top."""
    pic = slide.shapes.add_picture(image_path, x, y, width=w, height=h)
    # Add a thin border frame to imitate rounded card
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.adjustments[0] = corner_in / 2
    frame.fill.background()
    frame.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    frame.line.width = Pt(0)
    frame.shadow.inherit = False
    return pic


def add_pill(slide, x, y, text, *, fill=ACCENT, text_color=WHITE, border=ACCENT,
             font_size=9, h_in=0.34, padding=0.18):
    w = Inches(max(1.2, 0.13 * len(text) + 0.4))
    h = Inches(h_in)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.color.rgb = border
    pill.line.width = Pt(0.75)
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Inches(padding); tf.margin_right = Inches(padding)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", "300")
    return pill, w


def add_brand_top(slide, on_dark=False):
    slide.shapes.add_picture(LOGO, MARGIN_X, Inches(0.4), height=Inches(0.42))
    add_text(slide, SLIDE_W - MARGIN_X - Inches(3), Inches(0.5),
             Inches(3), Inches(0.3),
             "acoola.team", font=FONT_BODY, size=10,
             color=WHITE if on_dark else MUTED, align="right")


def add_footer(slide, n, total, on_dark=False):
    line_color = RGBColor(0x33, 0x3A, 0x44) if on_dark else BORDER
    text_color = RGBColor(0xC0, 0xC4, 0xCC) if on_dark else MUTED
    line = slide.shapes.add_connector(1, MARGIN_X, SLIDE_H - Inches(0.6),
                                       SLIDE_W - MARGIN_X, SLIDE_H - Inches(0.6))
    line.line.color.rgb = line_color
    line.line.width = Pt(0.5)
    add_text(slide, MARGIN_X, SLIDE_H - Inches(0.5),
             Inches(6), Inches(0.3),
             "Сопровождение сайта и маркетинга  /  для Совкомбанка",
             font=FONT_BODY, size=9.5, color=text_color)
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2), SLIDE_H - Inches(0.5),
             Inches(2), Inches(0.3),
             f"{n:02d} / {total:02d}", font=FONT_BODY, size=10,
             color=text_color, align="right")


def add_section_intro(slide, num, eyebrow, lines, *, accent_words=None,
                       size=40, max_w=Inches(11.0), eyebrow_y=Inches(1.05)):
    add_eyebrow(slide, MARGIN_X, eyebrow_y, f"{num} / {eyebrow}")
    accent_words = accent_words or set()
    tb = slide.shapes.add_textbox(MARGIN_X, eyebrow_y + Inches(0.45),
                                   max_w, Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = ACCENT if line in accent_words else INK


def add_decorative_ring(slide, cx, cy, diameter, *, color=ACCENT,
                         line_w=1.5, alpha_pct=None):
    """Decorative outlined circle. cx, cy is center."""
    from pptx.oxml.ns import qn
    from lxml import etree
    r = diameter / 2
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    int(cx - r), int(cy - r),
                                    int(diameter), int(diameter))
    ring.fill.background()
    ring.line.color.rgb = color
    ring.line.width = Pt(line_w)
    ring.shadow.inherit = False
    if alpha_pct is not None:
        ln = ring.line._get_or_add_ln() if hasattr(ring.line, "_get_or_add_ln") else None
    return ring


def add_dot_grid(slide, x, y, cols, rows, step=Inches(0.18),
                  dot_size=Inches(0.05), color=None):
    """Decorative dot grid pattern."""
    if color is None:
        color = MUTED_2
    for r in range(rows):
        for c in range(cols):
            dx = x + c * step
            dy = y + r * step
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL, dx, dy,
                                        dot_size, dot_size)
            d.fill.solid(); d.fill.fore_color.rgb = color
            d.line.fill.background(); d.shadow.inherit = False


def add_diagonal_accent(slide, x, y, w, h, *, color=ACCENT):
    """Slim diagonal blue strip — decorative."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    bar.rotation = 18
    return bar


def add_dark_overlay(slide, x, y, w, h, alpha_pct=55):
    """Dark scrim over an image area for legibility of text overlay."""
    from pptx.oxml.ns import qn
    from lxml import etree
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = INK
    rect.shadow.inherit = False
    sp = rect.fill.fore_color._xFill
    if sp.tag.endswith("solidFill"):
        srgb = sp.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(alpha_pct * 1000)))
    return rect


# ============== Slides ==============

TOTAL = 13


def slide_title(prs):
    """Hero photo right + dramatic typography left + decorative ring."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)

    # Decorative blue ring partially off-canvas (bottom-left)
    add_decorative_ring(s, Inches(-0.5), SLIDE_H + Inches(0.5),
                         Inches(4.0), color=ACCENT, line_w=1.5)
    # Decorative dot grid (top-left small accent)
    add_dot_grid(s, MARGIN_X + Inches(7.0), Inches(0.55),
                  cols=5, rows=2, step=Inches(0.16),
                  dot_size=Inches(0.045), color=MUTED_2)

    # Right-side hero photo (pre-cropped to exact 4.78/5.5 aspect)
    hero_x = Inches(7.85)
    hero_y = Inches(1.35)
    hero_w = SLIDE_W - hero_x - MARGIN_X
    hero_h = Inches(5.5)
    s.shapes.add_picture(PHOTO_BANKING_PORTRAIT,
                          hero_x, hero_y, width=hero_w, height=hero_h)
    # rounded inner border
    ring = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               hero_x, hero_y, hero_w, hero_h)
    ring.adjustments[0] = 0.04
    ring.fill.background()
    ring.line.color.rgb = BG
    ring.line.width = Pt(2)
    ring.shadow.inherit = False

    # Decorative ring overlapping image edge (top-right of image)
    add_decorative_ring(s, hero_x, hero_y + Inches(0.2),
                         Inches(0.9), color=ACCENT, line_w=2)

    # Floating client badge — bottom-left of photo
    overlay_w = Inches(3.0); overlay_h = Inches(0.95)
    ox = hero_x + Inches(0.4); oy = hero_y + hero_h - overlay_h - Inches(0.4)
    add_card(s, ox, oy, overlay_w, overlay_h, fill=ACCENT,
             border=ACCENT, corner=0.12)
    add_text(s, ox + Inches(0.3), oy + Inches(0.18),
             overlay_w - Inches(0.6), Inches(0.3),
             "ДЛЯ КЛИЕНТА",
             font=FONT_BODY, size=8, bold=True, color=WHITE, spacing=3)
    add_text(s, ox + Inches(0.3), oy + Inches(0.42),
             overlay_w - Inches(0.6), Inches(0.5),
             "ПАО «Совкомбанк»", font=FONT_DISPLAY,
             size=16, bold=True, color=WHITE)

    # Top brand
    s.shapes.add_picture(LOGO, MARGIN_X, Inches(0.5), height=Inches(0.42))
    add_text(s, MARGIN_X, Inches(1.0), Inches(3), Inches(0.2),
             "DIGITAL-АГЕНТСТВО",
             font=FONT_BODY, size=8, bold=True, color=MUTED, spacing=4)

    # Eyebrow
    add_eyebrow(s, MARGIN_X, Inches(1.95), "Коммерческое предложение")

    # MASSIVE heading
    text_w = Inches(6.9)
    tb = s.shapes.add_textbox(MARGIN_X, Inches(2.55),
                               text_w, Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Сопровождение", "сайта", "и маркетинга."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 0.95
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(48)
        r.font.bold = True
        r.font.color.rgb = ACCENT if i == 2 else INK

    # Subtitle below
    add_text(s, MARGIN_X, Inches(5.95), Inches(6.7), Inches(0.7),
             "Внешний маркетингово-технический контур для банка.",
             font=FONT_DISPLAY, size=14, bold=True, color=TEXT,
             line_spacing=1.3)

    # Bottom meta strip with accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               MARGIN_X, Inches(6.55),
                               Inches(0.3), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, MARGIN_X + Inches(0.45), Inches(6.5), Inches(7), Inches(0.3),
             "ACOOLA TEAM   /   2025   /   acoola.team",
             font=FONT_BODY, size=9, bold=True, color=MUTED, spacing=3)

    add_footer(s, 1, TOTAL)


def slide_audience(prs):
    """Image-text split: photo right, content left."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)

    # Decorative dot grid behind image
    add_dot_grid(s, Inches(7.7), Inches(1.0),
                  cols=3, rows=12, step=Inches(0.18),
                  dot_size=Inches(0.05), color=MUTED_2)

    # Right-side image (pre-cropped to exact aspect)
    img_x = Inches(8.0)
    img_y = Inches(1.4)
    img_w = SLIDE_W - img_x - MARGIN_X
    img_h = Inches(5.5)
    s.shapes.add_picture(PHOTO_TEAM_PORTRAIT,
                          img_x, img_y, width=img_w, height=img_h)
    ring = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               img_x, img_y, img_w, img_h)
    ring.adjustments[0] = 0.04
    ring.fill.background()
    ring.line.color.rgb = BG
    ring.line.width = Pt(2)
    ring.shadow.inherit = False
    # Decorative ring on image edge
    add_decorative_ring(s, img_x, img_y + img_h, Inches(0.7),
                         color=ACCENT, line_w=2)

    # Floating stat card on photo
    sx = img_x + Inches(0.35)
    sy = img_y + img_h - Inches(1.45)
    add_card(s, sx, sy, Inches(2.5), Inches(1.05), fill=WHITE, corner=0.1)
    add_text(s, sx + Inches(0.25), sy + Inches(0.15),
             Inches(2), Inches(0.3),
             "В КОМАНДЕ", font=FONT_BODY, size=8, bold=True,
             color=MUTED, spacing=3)
    add_text(s, sx + Inches(0.25), sy + Inches(0.4),
             Inches(2), Inches(0.5),
             "7 ролей", font=FONT_DISPLAY, size=22, bold=True,
             color=INK)

    # Left content
    add_eyebrow(s, MARGIN_X, Inches(1.4), "01 / Кому это адресовано")
    tb = s.shapes.add_textbox(MARGIN_X, Inches(1.9),
                               Inches(7.0), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Это для банка", "с большим сайтом", "и подрядчиками."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 0.95
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(38)
        r.font.bold = True
        r.font.color.rgb = ACCENT if i == 2 else INK

    body = ("Это предложение для банка с работающим сайтом, продуктовыми лендингами, "
            "активными релизами и командой подрядчиков. Маркетингу не хватает технической "
            "экспертизы под рукой.")
    add_text(s, MARGIN_X, Inches(5.55), Inches(6.8), Inches(1.3),
             body, size=13.5, color=MUTED, line_spacing=1.55)

    add_footer(s, 2, TOTAL)


def slide_problem(prs):
    """Heading + decorative image accent + pain cards."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)

    # Right-side accent photo (small, decorative — analytics dashboard)
    img_x = Inches(8.6); img_y = Inches(1.4)
    img_w = Inches(4.0); img_h = Inches(2.3)
    s.shapes.add_picture(PHOTO_DASHBOARD_169, img_x, img_y,
                          width=img_w, height=img_h)
    ring = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               img_x, img_y, img_w, img_h)
    ring.adjustments[0] = 0.05
    ring.fill.background()
    ring.line.color.rgb = BG; ring.line.width = Pt(2)
    ring.shadow.inherit = False
    # Decorative accent ring partial overlap
    add_decorative_ring(s, img_x, img_y + img_h, Inches(0.6),
                         color=ACCENT, line_w=2)

    # Heading area
    add_eyebrow(s, MARGIN_X, Inches(1.05), "02 / Проблема")
    tb = s.shapes.add_textbox(MARGIN_X, Inches(1.55),
                               Inches(11), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Сайт банка", "ломается тихо."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(44)
        r.font.bold = True
        r.font.color.rgb = ACCENT if i == 1 else INK

    intro = ("Подрядчик случайно закрыл продуктовую страницу от индексации. "
             "Маркетинг обновил тарифы, и калькулятор перестал считать. "
             "Падение трафика заметили через две недели.")
    add_text(s, MARGIN_X, Inches(4.0), Inches(11.5), Inches(1.1),
             intro, size=13, color=MUTED, line_spacing=1.5)

    # 4 pain cards, more visual with big number
    pains = [
        ("01", "Нет экспертизы", "Маркетингу не хватает технической стороны."),
        ("02", "Нет контроля", "Подрядчики выпускают задачи без проверки."),
        ("03", "Нет системности", "Никто не смотрит на сайт целиком."),
        ("04", "Нет SEO-надзора", "Доработки идут без оценки рисков для трафика."),
    ]
    cw = Inches(2.95); ch = Inches(1.7); gap = Inches(0.12)
    total_w = cw * 4 + gap * 3
    sx = (SLIDE_W - total_w) / 2
    cy = Inches(5.15)
    for i, (n, t, b) in enumerate(pains):
        x = sx + i * (cw + gap)
        add_card(s, x, cy, cw, ch)
        add_text(s, x + Inches(0.3), cy + Inches(0.2),
                 cw - Inches(0.6), Inches(0.4),
                 n, font=FONT_DISPLAY, size=11, bold=True,
                 color=ACCENT, spacing=2)
        add_text(s, x + Inches(0.3), cy + Inches(0.55),
                 cw - Inches(0.6), Inches(0.4),
                 t, font=FONT_DISPLAY, size=14, bold=True, color=INK)
        add_text(s, x + Inches(0.3), cy + Inches(0.95),
                 cw - Inches(0.6), Inches(0.7),
                 b, size=10.5, color=MUTED, line_spacing=1.4)

    add_footer(s, 3, TOTAL)


def slide_role(prs):
    """Bold heading + accent strip + 3 pillar cards."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)

    # Accent strip bar (decorative)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_X,
                              Inches(1.05), Inches(0.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, MARGIN_X + Inches(0.65), Inches(0.95),
             Inches(8), Inches(0.28),
             "03 / НАША РОЛЬ", font=FONT_BODY, size=9, bold=True,
             color=MUTED, spacing=3)

    # Heading
    tb = s.shapes.add_textbox(MARGIN_X, Inches(1.55),
                               Inches(12), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Внешний маркетингово-",
                                "технический отдел банка."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = ACCENT if i == 1 else INK

    body = ("Мы не заменяем маркетинг банка и не вытесняем подрядчиков. Помогаем "
            "им работать чище: ставим задачи, проверяем результат, замечаем "
            "риски заранее. Acoola Team становится тем самым старшим специалистом, "
            "который держит сайт в порядке.")
    add_text(s, MARGIN_X, Inches(3.95), Inches(11.5), Inches(1.5),
             body, size=14, color=MUTED, line_spacing=1.55)

    pillars = [
        ("Контроль", "Следим за подрядчиками и доработками."),
        ("Экспертиза", "SEO, разработка, контент, аналитика."),
        ("Прозрачность", "Отчёт по часам и задачам каждый месяц."),
    ]
    cw = Inches(3.85); ch = Inches(1.4); gap = Inches(0.25)
    total_w = cw * 3 + gap * 2
    sx = (SLIDE_W - total_w) / 2
    cy = Inches(5.4)
    for i, (t, b) in enumerate(pillars):
        x = sx + i * (cw + gap)
        add_card(s, x, cy, cw, ch)
        # filled accent square icon
        icon = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(0.3), cy + Inches(0.27),
                                   Inches(0.32), Inches(0.32))
        icon.adjustments[0] = 0.25
        icon.fill.solid(); icon.fill.fore_color.rgb = ACCENT
        icon.line.fill.background()
        icon.shadow.inherit = False
        # symbol
        tficon = icon.text_frame
        tficon.margin_left = 0; tficon.margin_right = 0
        tficon.margin_top = 0; tficon.margin_bottom = 0
        picon = tficon.paragraphs[0]; picon.alignment = PP_ALIGN.CENTER
        ricon = picon.add_run(); ricon.text = ["✓","◆","◉"][i]
        ricon.font.name = FONT_BODY; ricon.font.size = Pt(11)
        ricon.font.bold = True
        ricon.font.color.rgb = WHITE
        add_text(s, x + Inches(0.75), cy + Inches(0.25),
                 cw - Inches(1.05), Inches(0.4),
                 t, font=FONT_DISPLAY, size=15, bold=True, color=INK)
        add_text(s, x + Inches(0.3), cy + Inches(0.78),
                 cw - Inches(0.6), Inches(0.7),
                 b, size=11, color=MUTED, line_spacing=1.45)

    add_footer(s, 4, TOTAL)


def slide_what(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)
    add_section_intro(s, "04", "Что мы делаем",
                      ["Внутри ежемесячного", "пакета часов."], size=36)

    items = [
        "Разбираем задачи и расставляем приоритеты",
        "Ставим задачи подрядчикам и нашим специалистам",
        "Проверяем результат до и после внедрения",
        "Следим за SEO-рисками доработок",
        "Помогаем маркетингу банка с технической частью",
        "Работаем с контентом, продуктовыми страницами, аналитикой",
        "Ведём отчёт по часам и задачам каждый месяц",
        "Даём рекомендации по сайту и инфраструктуре",
    ]
    col_w = Inches(5.6); col_gap = Inches(0.3)
    left_x = MARGIN_X
    right_x = MARGIN_X + col_w + col_gap
    cy = Inches(4)
    half = (len(items) + 1) // 2
    for col_items, x in [(items[:half], left_x), (items[half:], right_x)]:
        for i, it in enumerate(col_items):
            yy = cy + Inches(i * 0.6)
            chk = s.shapes.add_shape(MSO_SHAPE.OVAL, x,
                                      yy + Inches(0.08), Inches(0.28), Inches(0.28))
            chk.fill.solid(); chk.fill.fore_color.rgb = ACCENT
            chk.line.fill.background(); chk.shadow.inherit = False
            tf = chk.text_frame
            tf.margin_left = 0; tf.margin_right = 0
            tf.margin_top = 0; tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "✓"
            r.font.name = FONT_BODY; r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = WHITE
            add_text(s, x + Inches(0.5), yy + Inches(0.07),
                     col_w - Inches(0.6), Inches(0.4),
                     it, size=12, color=INK)

    add_footer(s, 5, TOTAL)


def slide_team(prs):
    """Photo-banner top + role pills below."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)

    # Decorative dot grid top-right
    add_dot_grid(s, Inches(10.5), Inches(1.05),
                  cols=8, rows=4, step=Inches(0.18),
                  dot_size=Inches(0.05), color=MUTED_2)

    add_eyebrow(s, MARGIN_X, Inches(1.05), "05 / Команда")
    tb = s.shapes.add_textbox(MARGIN_X, Inches(1.55),
                               Inches(12), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Кто работает", "внутри пакета часов."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = INK

    add_text(s, MARGIN_X, Inches(3.95), Inches(11.5), Inches(0.7),
             "Банк общается с одним аккаунт-менеджером. Команду под задачу подбираем сами.",
             size=13, color=MUTED, line_spacing=1.5)

    roles = [
        ("Аккаунт", "PM"),
        ("SEO", "Специалист"),
        ("Разработчик", "Front / Back"),
        ("Контент-", "менеджер"),
        ("Копирайтер", "SEO-тексты"),
        ("Дизайнер", "Web / UI"),
        ("Аналитик", "Метрики, GA"),
    ]
    cw = Inches(2.85); ch = Inches(1.05); gap_x = Inches(0.18); gap_y = Inches(0.18)
    cy = Inches(4.55)
    row1, row2 = roles[:4], roles[4:]
    total_w1 = cw * 4 + gap_x * 3
    sx1 = (SLIDE_W - total_w1) / 2
    for i, r in enumerate(row1):
        _role_card(s, sx1 + i * (cw + gap_x), cy, cw, ch, r[0], r[1])
    total_w2 = cw * 3 + gap_x * 2
    sx2 = (SLIDE_W - total_w2) / 2
    for i, r in enumerate(row2):
        _role_card(s, sx2 + i * (cw + gap_x), cy + ch + gap_y, cw, ch, r[0], r[1])

    add_footer(s, 6, TOTAL)


def _role_card(s, x, y, w, h, title, subtitle):
    add_card(s, x, y, w, h)
    # accent square icon
    icon = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               x + Inches(0.3), y + Inches(0.3),
                               Inches(0.4), Inches(0.4))
    icon.adjustments[0] = 0.25
    icon.fill.solid(); icon.fill.fore_color.rgb = ACCENT
    icon.line.fill.background(); icon.shadow.inherit = False
    tf = icon.text_frame
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title[0]
    r.font.name = FONT_DISPLAY; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE
    add_text(s, x + Inches(0.83), y + Inches(0.25),
             w - Inches(1.0), Inches(0.4),
             title, font=FONT_DISPLAY, size=13, bold=True, color=INK)
    add_text(s, x + Inches(0.83), y + Inches(0.55),
             w - Inches(1.0), Inches(0.3),
             subtitle, size=10, color=MUTED)


def slide_format(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)
    add_section_intro(s, "06", "Формат работы",
                      ["Один цикл,", "повторяющийся каждый месяц."], size=36)

    steps = [
        ("01", "Запрос", "Задача поступает в закреплённый канал."),
        ("02", "Оценка", "Аккаунт фиксирует часы и приоритет."),
        ("03", "Исполнение", "Задача уходит нужному специалисту."),
        ("04", "Проверка", "Внутренний контроль качества."),
        ("05", "Отчёт", "По итогам месяца: часы и статусы задач."),
    ]
    cw = Inches(2.36); ch = Inches(2.4); gap = Inches(0.12)
    total_w = cw * 5 + gap * 4
    sx = (SLIDE_W - total_w) / 2
    cy = Inches(4.1)
    for i, (num, t, b) in enumerate(steps):
        x = sx + i * (cw + gap)
        add_card(s, x, cy, cw, ch)
        # accent number block
        numbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x + Inches(0.25), cy + Inches(0.25),
                                     Inches(0.7), Inches(0.4))
        numbox.adjustments[0] = 0.3
        numbox.fill.solid(); numbox.fill.fore_color.rgb = SURFACE2
        numbox.line.fill.background(); numbox.shadow.inherit = False
        tfn = numbox.text_frame
        tfn.margin_left = 0; tfn.margin_right = 0
        tfn.margin_top = 0; tfn.margin_bottom = 0
        pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run(); rn.text = num
        rn.font.name = FONT_DISPLAY; rn.font.size = Pt(11); rn.font.bold = True
        rn.font.color.rgb = ACCENT
        rPr = rn._r.get_or_add_rPr(); rPr.set("spc", "200")

        add_text(s, x + Inches(0.25), cy + Inches(0.85),
                 cw - Inches(0.5), Inches(0.5),
                 t, font=FONT_DISPLAY, size=15, bold=True, color=INK)
        add_text(s, x + Inches(0.25), cy + Inches(1.35),
                 cw - Inches(0.5), Inches(1.0),
                 b, size=10.5, color=MUTED, line_spacing=1.45)
        # connector arrow between cards
        if i < 4:
            arr_x = x + cw
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                      arr_x + Inches(0.0),
                                      cy + ch / 2 - Inches(0.05),
                                      Inches(0.12), Inches(0.1))
            arr.fill.solid(); arr.fill.fore_color.rgb = BORDER
            arr.line.fill.background(); arr.shadow.inherit = False

    add_text(s, MARGIN_X, Inches(6.65), Inches(11.5), Inches(0.4),
             "Раз в квартал проводим стратегическую встречу: что работает, что менять.",
             size=11, color=MUTED, align="center")
    add_footer(s, 7, TOTAL)


def slide_pricing(prs):
    """Pricing with elevated highlighted tier."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)
    add_section_intro(s, "07", "Тарифы",
                      ["Часы × ставка =", "абонентская плата."],
                      size=34, accent_words={"абонентская плата."})

    add_text(s, MARGIN_X, Inches(3.7), Inches(12), Inches(0.4),
             "Все цены без НДС.", size=11, color=MUTED, spacing=2)

    plans = [
        {"name": "Базовый", "hours": "20",
         "rate": "3 500 ₽", "monthly": "70 000",
         "over": "4 000 ₽", "highlight": False},
        {"name": "Оптимальный", "hours": "40",
         "rate": "3 200 ₽", "monthly": "128 000",
         "over": "3 500 ₽", "highlight": True},
        {"name": "Максимальный", "hours": "80",
         "rate": "2 800 ₽", "monthly": "224 000",
         "over": "3 200 ₽", "highlight": False},
    ]
    cw = Inches(3.85); ch = Inches(2.65); gap = Inches(0.2)
    total_w = cw * 3 + gap * 2
    sx = (SLIDE_W - total_w) / 2
    cy = Inches(4.0)
    for i, plan in enumerate(plans):
        x = sx + i * (cw + gap)
        if plan["highlight"]:
            # raised card — higher position with shadow effect (use accent fill)
            ec = add_card(s, x, cy - Inches(0.15), cw, ch + Inches(0.3),
                          fill=INK, border=INK, line_w=0, corner=0.06)
            text_color_name = WHITE
            text_color_price = ACCENT
            text_color_label = RGBColor(0xC0, 0xC4, 0xCC)
            divider_color = RGBColor(0x33, 0x3A, 0x44)
            # ribbon
            rib = add_card(s, x, cy - Inches(0.5),
                           cw, Inches(0.35), fill=ACCENT,
                           border=ACCENT, corner=0.4)
            tf = rib.text_frame
            tf.margin_left = 0; tf.margin_right = 0
            tf.margin_top = Inches(0.04); tf.margin_bottom = 0
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run(); rr.text = "ВЫБОР"
            rr.font.name = FONT_BODY; rr.font.size = Pt(9); rr.font.bold = True
            rr.font.color.rgb = WHITE
            rPr = rr._r.get_or_add_rPr(); rPr.set("spc", "400")
            cy_card = cy - Inches(0.15)
        else:
            add_card(s, x, cy, cw, ch)
            text_color_name = INK
            text_color_price = INK
            text_color_label = MUTED
            divider_color = BORDER
            cy_card = cy

        add_text(s, x + Inches(0.3), cy_card + Inches(0.25),
                 cw - Inches(0.6), Inches(0.4),
                 plan["name"], font=FONT_DISPLAY, size=14,
                 bold=True, color=text_color_name)
        # huge price number
        add_text(s, x + Inches(0.3), cy_card + Inches(0.7),
                 cw - Inches(0.6), Inches(1.0),
                 plan["monthly"], font=FONT_DISPLAY, size=36, bold=True,
                 color=text_color_price)
        add_text(s, x + Inches(0.3), cy_card + Inches(1.55),
                 cw - Inches(0.6), Inches(0.3),
                 "₽ в месяц, без НДС", size=10, color=text_color_label)
        ln = s.shapes.add_connector(1, x + Inches(0.3), cy_card + Inches(1.95),
                                     x + cw - Inches(0.3), cy_card + Inches(1.95))
        ln.line.color.rgb = divider_color; ln.line.width = Pt(0.5)
        rows = [
            ("Часов в месяц", plan["hours"]),
            ("Ставка часа", plan["rate"]),
            ("Сверхчас", plan["over"]),
        ]
        ry = cy_card + Inches(2.05)
        for label, val in rows:
            add_text(s, x + Inches(0.3), ry, Inches(1.6), Inches(0.25),
                     label, size=10, color=text_color_label)
            add_text(s, x + Inches(1.9), ry, cw - Inches(2.2), Inches(0.25),
                     val, size=10.5, bold=True,
                     color=text_color_name, align="right")
            ry += Inches(0.25)

    add_footer(s, 8, TOTAL)


def slide_examples(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)
    add_section_intro(s, "08", "Примеры задач",
                      ["Что можно решать", "в рамках поддержки."], size=36)

    # Wide thin chart strip behind heading (decorative accent)
    strip_x = Inches(8.0); strip_y = Inches(1.4)
    strip_w = SLIDE_W - strip_x - MARGIN_X
    strip_h = Inches(1.4)
    s.shapes.add_picture(PHOTO_CHART_WIDE, strip_x, strip_y,
                          width=strip_w, height=strip_h)
    ring = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               strip_x, strip_y, strip_w, strip_h)
    ring.adjustments[0] = 0.06
    ring.fill.background()
    ring.line.color.rgb = BG; ring.line.width = Pt(2)
    ring.shadow.inherit = False

    examples = [
        "Проверить задачу подрядчика перед запуском",
        "Оценить SEO-риски доработки продуктовой страницы",
        "Подготовить ТЗ для разработчика",
        "Проверить корректность калькуляторов",
        "Обновить тарифы и условия на сайте",
        "Доработать страницы вкладов и кредитов",
        "Проверить формы заявок и аналитику",
        "Исправить технические ошибки",
        "Проконтролировать подрядчика",
        "Дать рекомендации маркетингу",
        "Проверить посадочные перед рекламой",
        "Подготовить SEO-тексты для продуктов",
    ]
    cw = Inches(2.95); ch = Inches(0.7); gap_x = Inches(0.15); gap_y = Inches(0.15)
    cols = 4
    total_w = cw * cols + gap_x * (cols - 1)
    sx = (SLIDE_W - total_w) / 2
    cy = Inches(4)
    for i, ex in enumerate(examples):
        col = i % cols
        row = i // cols
        x = sx + col * (cw + gap_x)
        y = cy + row * (ch + gap_y)
        add_card(s, x, y, cw, ch)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.12),
                                  Inches(0.04), ch - Inches(0.24))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background(); bar.shadow.inherit = False
        add_text(s, x + Inches(0.25), y + Inches(0.18),
                 cw - Inches(0.4), Inches(0.45),
                 ex, size=10.5, color=INK, line_spacing=1.3, anchor="middle")

    add_footer(s, 9, TOTAL)


def slide_get(prs):
    """STATS slide: 3 huge numbers + image accent."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)

    add_eyebrow(s, MARGIN_X, Inches(1.05), "09 / Что банк получает")
    # Heading
    tb = s.shapes.add_textbox(MARGIN_X, Inches(1.55),
                               Inches(12), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Внешняя команда,",
                                "которая держит сайт в порядке."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = ACCENT if i == 1 else INK

    # 3 huge stat columns
    stats = [
        ("1", "Точка входа", "Один аккаунт-менеджер на стороне Acoola Team."),
        ("7", "Ролей в команде", "PM, SEO, dev, контент, копирайтер, дизайнер, аналитик."),
        ("100%", "Прозрачность", "Отчёт по каждому потраченному часу."),
    ]
    cw = Inches(3.85); ch = Inches(2.6); gap = Inches(0.25)
    total_w = cw * 3 + gap * 2
    sx = (SLIDE_W - total_w) / 2
    cy = Inches(4.05)
    for i, (n, t, b) in enumerate(stats):
        x = sx + i * (cw + gap)
        # Big stat — no card background, just top border
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cy,
                                   cw, Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background(); bar.shadow.inherit = False
        # huge number
        add_text(s, x, cy + Inches(0.25), cw, Inches(1.5),
                 n, font=FONT_DISPLAY, size=84, bold=True, color=INK,
                 line_spacing=0.95)
        add_text(s, x, cy + Inches(1.7), cw, Inches(0.4),
                 t.upper(), font=FONT_BODY, size=10, bold=True,
                 color=ACCENT, spacing=3)
        add_text(s, x, cy + Inches(2.05), cw - Inches(0.2), Inches(0.7),
                 b, size=11, color=MUTED, line_spacing=1.45)

    add_footer(s, 10, TOTAL)


def slide_excludes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)
    add_section_intro(s, "10", "Что не входит",
                      ["Границы услуги.", "Чтобы не было разночтений."], size=36)

    items = [
        "Не гарантируем рост заявок и конверсии",
        "Не выкупаем рекламные бюджеты и сторонние сервисы",
        "Не разрабатываем сайт с нуля в рамках абонплаты",
        "Не делаем брендинг и фирменный стиль",
        "Не ведём рекламу как отдельный продукт",
        "Не настраиваем CRM и автоматизацию продаж",
        "Не отвечаем за юридические и регуляторные формулировки",
        "Не отвечаем за подрядчиков, выбранных против рекомендаций",
    ]
    col_w = Inches(5.6); col_gap = Inches(0.3)
    left_x = MARGIN_X
    right_x = MARGIN_X + col_w + col_gap
    cy = Inches(4)
    half = (len(items) + 1) // 2
    for col_items, x in [(items[:half], left_x), (items[half:], right_x)]:
        for i, it in enumerate(col_items):
            yy = cy + Inches(i * 0.6)
            chk = s.shapes.add_shape(MSO_SHAPE.OVAL, x,
                                      yy + Inches(0.08), Inches(0.28), Inches(0.28))
            chk.fill.solid(); chk.fill.fore_color.rgb = SURFACE2
            chk.line.color.rgb = MUTED_2; chk.line.width = Pt(0.8)
            chk.shadow.inherit = False
            tf = chk.text_frame
            tf.margin_left = 0; tf.margin_right = 0
            tf.margin_top = 0; tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "×"
            r.font.name = FONT_BODY; r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = MUTED
            add_text(s, x + Inches(0.5), yy + Inches(0.07),
                     col_w - Inches(0.6), Inches(0.4),
                     it, size=12, color=MUTED, line_spacing=1.4)

    add_text(s, MARGIN_X, Inches(6.85), Inches(11.5), Inches(0.3),
             "Эти задачи можно подключить отдельным проектом или передать профильным партнёрам.",
             size=11, color=MUTED_2, align="center")
    add_footer(s, 11, TOTAL)


def slide_overhours(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s); add_brand_top(s)
    add_section_intro(s, "11", "Сверхчасы",
                      ["Если задач больше,", "чем покрывает пакет."], size=36)

    card_w = Inches(11.93); card_h = Inches(1.7)
    cx = (SLIDE_W - card_w) / 2
    cy = Inches(3.85)
    add_card(s, cx, cy, card_w, card_h, fill=WHITE,
             border=ACCENT, line_w=1.5)
    add_text(s, cx + Inches(0.45), cy + Inches(0.18),
             Inches(0.6), Inches(0.6),
             "❝", font=FONT_DISPLAY, size=30, bold=True, color=ACCENT)
    quote = ("Дополнительные часы выставляются отдельным счётом по итогам месяца, "
             "по ставке сверхчаса выбранного тарифа. О приближении к лимиту "
             "предупреждаем заранее и согласовываем продолжение работ.")
    add_text(s, cx + Inches(1.2), cy + Inches(0.32),
             card_w - Inches(1.6), Inches(1.3),
             quote, size=13.5, color=INK, line_spacing=1.45)

    rates = [
        ("Базовый", "4 000 ₽", "за час сверх пакета"),
        ("Оптимальный", "3 500 ₽", "за час сверх пакета"),
        ("Максимальный", "3 200 ₽", "за час сверх пакета"),
    ]
    cw = Inches(3.85); ch = Inches(0.95); gap = Inches(0.25)
    total_w = cw * 3 + gap * 2
    sx = (SLIDE_W - total_w) / 2
    ry = Inches(5.85)
    for i, (n, v, sub) in enumerate(rates):
        x = sx + i * (cw + gap)
        add_card(s, x, ry, cw, ch)
        add_text(s, x + Inches(0.3), ry + Inches(0.18),
                 Inches(2), Inches(0.3),
                 n, font=FONT_DISPLAY, size=12, bold=True, color=INK)
        add_text(s, x + Inches(0.3), ry + Inches(0.5),
                 Inches(2), Inches(0.3),
                 sub, size=9.5, color=MUTED)
        add_text(s, x + Inches(0.3), ry + Inches(0.3),
                 cw - Inches(0.6), Inches(0.4),
                 v, font=FONT_DISPLAY, size=18, bold=True,
                 color=ACCENT, align="right")

    add_footer(s, 12, TOTAL)


def slide_next(prs):
    """Full-bleed handshake photo with dark overlay + bold CTA."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # full-bleed photo (pre-cropped to 16:9)
    s.shapes.add_picture(PHOTO_HANDSHAKE_169, 0, 0,
                          width=SLIDE_W, height=SLIDE_H)
    # dark overlay
    add_dark_overlay(s, 0, 0, SLIDE_W, SLIDE_H, alpha_pct=78)
    # Decorative blue ring
    add_decorative_ring(s, SLIDE_W - Inches(1.0), Inches(1.5),
                         Inches(2.5), color=ACCENT, line_w=2)
    add_decorative_ring(s, Inches(11.0), Inches(2.0),
                         Inches(0.8), color=ACCENT, line_w=2)

    # White logo on dark
    s.shapes.add_picture(LOGO_WHITE, MARGIN_X, Inches(0.5), height=Inches(0.42))
    add_text(s, SLIDE_W - MARGIN_X - Inches(3), Inches(0.5),
             Inches(3), Inches(0.3),
             "acoola.team", font=FONT_BODY, size=10,
             color=RGBColor(0xC0, 0xC4, 0xCC), align="right")

    add_eyebrow(s, MARGIN_X, Inches(1.95), "12 / Следующий шаг",
                color=RGBColor(0xC0, 0xC4, 0xCC))

    # Massive heading
    tb = s.shapes.add_textbox(MARGIN_X, Inches(2.55),
                               Inches(12), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, line in enumerate(["Готовы начать",
                                "со следующего", "рабочего дня."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 0.95
        r = p.add_run()
        r.text = line
        r.font.name = FONT_DISPLAY
        r.font.size = Pt(58)
        r.font.bold = True
        r.font.color.rgb = ACCENT if i == 2 else WHITE

    # Bottom CTA bar — centered, white
    cta_w = Inches(11.93); cta_h = Inches(0.95)
    cx = (SLIDE_W - cta_w) / 2
    cta_y = Inches(6.05)
    cta = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              cx, cta_y, cta_w, cta_h)
    cta.adjustments[0] = 0.4
    cta.fill.solid(); cta.fill.fore_color.rgb = WHITE
    cta.line.fill.background(); cta.shadow.inherit = False
    tf = cta.text_frame
    tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "acoola.team   /   hello@acoola.team   /   +7 (___) ___ __ __"
    r.font.name = FONT_DISPLAY; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = INK

    # Footer (dark variant)
    add_footer(s, 13, TOTAL, on_dark=True)


# ============== Build ==============

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_audience(prs)
    slide_problem(prs)
    slide_role(prs)
    slide_what(prs)
    slide_team(prs)
    slide_format(prs)
    slide_pricing(prs)
    slide_examples(prs)
    slide_get(prs)
    slide_excludes(prs)
    slide_overhours(prs)
    slide_next(prs)

    out = "/Users/pro2kuror/Desktop/Cloude/Acoola_KP_Sovcombank.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
